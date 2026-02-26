"""ATMS 교육용 PPT 생성 스크립트 (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 상수 ──
NAVY = RGBColor(30, 60, 120)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 245, 255)
DARK_TEXT = RGBColor(40, 40, 40)
ACCENT = RGBColor(200, 60, 50)
GREEN = RGBColor(40, 140, 60)
ORANGE = RGBColor(220, 140, 30)
LIGHT_YELLOW = RGBColor(255, 250, 230)
GRAY = RGBColor(100, 100, 100)


def add_bg(slide, color=LIGHT_BG):
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = font_name
    return txBox


def add_title_bar(slide, text):
    """상단 제목 바."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            text, lvl, bld, clr = item
        else:
            text, lvl, bld, clr = item, 0, False, color

        p.text = text
        p.level = lvl
        p.font.size = Pt(font_size)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.space_after = Pt(4)
    return txBox


def add_box(slide, left, top, width, height, text, fill_color,
            text_color=DARK_TEXT, font_size=14, bold=False,
            alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    return shape


# ══════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, NAVY)
add_textbox(sl, 1, 1.5, 11, 1.5, "ATMS", 54, True, WHITE, PP_ALIGN.CENTER)
add_textbox(sl, 1, 3.0, 11, 0.8,
            "Assumption-based Truth Maintenance System",
            24, False, RGBColor(180, 200, 240), PP_ALIGN.CENTER)
add_textbox(sl, 1, 4.2, 11, 0.6,
            "목적  |  해결 과제  |  방식  |  의의",
            20, True, WHITE, PP_ALIGN.CENTER)
add_textbox(sl, 1, 5.8, 11, 0.8,
            "Building Problem Solvers\n"
            "Kenneth D. Forbus & Johan de Kleer, 1993",
            14, False, RGBColor(150, 170, 200), PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 슬라이드 2: 왜 ATMS가 필요한가?
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  왜 ATMS가 필요한가?")

add_textbox(sl, 0.5, 1.4, 12, 0.5,
            "기존 JTMS의 한계: 한 번에 하나의 세계만 볼 수 있다",
            18, True, ACCENT)

# 왼쪽: JTMS
add_box(sl, 0.5, 2.2, 5.8, 4.5,
        "", RGBColor(255, 235, 235))
add_textbox(sl, 0.7, 2.3, 5.4, 0.4, "JTMS (기존)", 18, True, ACCENT)
add_bullet_list(sl, 0.7, 2.9, 5.4, 3.5, [
    ("노드.label = :IN 또는 :OUT", 0, True, DARK_TEXT),
    ("한 시점에 하나의 세계만 유지", 0, False, DARK_TEXT),
    ("", 0, False, DARK_TEXT),
    ("진단 문제 예시:", 0, True, DARK_TEXT),
    ('"질병1" 가정 → 추론 → 안 맞음', 1, False, GRAY),
    ("→ 철회 → 처음부터 다시!", 1, False, ACCENT),
    ('"질병2" 가정 → 추론 → 또 안 맞음', 1, False, GRAY),
    ("→ 또 철회... 매번 재계산!", 1, False, ACCENT),
], font_size=15)

# 오른쪽: ATMS
add_box(sl, 7.0, 2.2, 5.8, 4.5,
        "", RGBColor(230, 245, 230))
add_textbox(sl, 7.2, 2.3, 5.4, 0.4, "ATMS (혁신)", 18, True, GREEN)
add_bullet_list(sl, 7.2, 2.9, 5.4, 3.5, [
    ("노드.label = [{A}, {B,C}]", 0, True, DARK_TEXT),
    ("모든 가능한 세계를 동시에 추적!", 0, False, GREEN),
    ("", 0, False, DARK_TEXT),
    ("같은 진단 문제:", 0, True, DARK_TEXT),
    ("질병1, 질병2, 질병3 가설 동시 유지", 1, False, GRAY),
    ("→ 검사 결과가 나오면 해당 가설만 제거", 1, False, GREEN),
    ("→ 나머지는 이미 계산되어 있음!", 1, False, GREEN),
    ("→ 재계산 없음, 자동 우회!", 1, False, GREEN),
], font_size=15)

# ══════════════════════════════════════════════════════════
# 슬라이드 3: ATMS 핵심 개념
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  ATMS 핵심 개념")

# 4개 박스
concepts = [
    ("가정\n(Assumption)", "독립적 선택지\n참/거짓 자유 설정\n예: A='질병1'",
     RGBColor(220, 235, 255)),
    ("환경\n(Environment)", "가정들의 집합\n= 하나의 '가능한 세계'\n예: {A, C}",
     RGBColor(220, 245, 220)),
    ("레이블\n(Label)", "최소 환경들의 집합\n= antichain\n예: [{A}, {B,C}]",
     RGBColor(255, 240, 220)),
    ("Nogood", "모순인 환경\n이 조합은 불가능!\n예: {A,B}*",
     RGBColor(255, 225, 225)),
]

for i, (title, desc, color) in enumerate(concepts):
    x = 0.5 + i * 3.2
    add_box(sl, x, 1.5, 2.9, 1.2, title, color, DARK_TEXT, 16, True)
    add_box(sl, x, 2.85, 2.9, 1.5, desc, WHITE, GRAY, 13)

# 핵심 아이디어
add_box(sl, 0.5, 4.8, 12.3, 1.0,
        '핵심 아이디어: "노드의 참/거짓을 묻지 말고,\n'
        '어떤 가정 하에서 참인지를 물어라"',
        LIGHT_YELLOW, NAVY, 18, True)

# 하단 비교
add_textbox(sl, 0.5, 6.2, 6, 0.5,
            'JTMS: "X는 참인가?" (이진 답변)', 14, False, ACCENT)
add_textbox(sl, 6.5, 6.2, 6, 0.5,
            'ATMS: "X는 어떤 가정 하에서 참인가?"', 14, True, GREEN)

# ══════════════════════════════════════════════════════════
# 슬라이드 4: JTMS vs ATMS 비교표
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  JTMS vs ATMS 구조적 비교")

rows = [
    ("측면", "JTMS", "ATMS"),
    ("노드 레이블", ":IN / :OUT (이진값)", "환경 리스트 [{A}, {B,C}]"),
    ("동시 세계", "1개", "모든 일관된 세계"),
    ("모순 처리", "멈추고 철회", "nogood 표시 → 자동 격리"),
    ("가정 변경", "철회 → 재추론", "이미 계산되어 있음"),
    ("질문 방식", '"X는 참인가?"', '"어떤 가정 하에서 참인가?"'),
    ("핵심 연산", "enable / retract", "weave (교차곱 + 필터링)"),
]

col_w = [2.5, 4.5, 5.3]
y_start = 1.6
for r, (c1, c2, c3) in enumerate(rows):
    y = y_start + r * 0.7
    bold = (r == 0)
    bg = NAVY if r == 0 else (RGBColor(248, 248, 255) if r % 2 == 0 else WHITE)
    tc = WHITE if r == 0 else DARK_TEXT

    x = 0.5
    for ci, (text, w) in enumerate(zip([c1, c2, c3], col_w)):
        fs = 14 if r == 0 else 13
        add_box(sl, x, y, w, 0.6, text, bg, tc, fs, bold)
        x += w + 0.05

# ══════════════════════════════════════════════════════════
# 슬라이드 5: justify-node 파이프라인
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  justify-node 파이프라인")

# 파이프라인 흐름
steps = [
    ("① 구조체 생성", "Justification\n(informant,\nconsequence,\nantecedents)",
     RGBColor(220, 235, 255)),
    ("② 양방향 링크", "consequence.justs\n← just\nantecedent.consequences\n← just",
     RGBColor(220, 245, 220)),
    ("③ propagate", "weave\n(교차곱)\n↓\nupdate\n(label 갱신)",
     RGBColor(255, 240, 220)),
]

for i, (title, desc, color) in enumerate(steps):
    x = 0.5 + i * 4.3
    add_box(sl, x, 1.5, 3.8, 0.7, title, NAVY, WHITE, 16, True)
    add_box(sl, x, 2.35, 3.8, 2.3, desc, color, DARK_TEXT, 14)

# 분기 다이어그램
add_textbox(sl, 0.5, 5.0, 12, 0.5, "update의 분기:", 16, True, NAVY)

add_box(sl, 0.5, 5.6, 5.8, 1.2,
        "일반 노드\n→ label 갱신\n→ 하류 justification 재전파",
        RGBColor(220, 245, 220), DARK_TEXT, 14, False, PP_ALIGN.LEFT)

add_box(sl, 7.0, 5.6, 5.8, 1.2,
        "모순 노드 (contradictory)\n→ new-nogood 5단계 실행!\n→ 금지 도장 + label 제거 + 오염",
        RGBColor(255, 225, 225), ACCENT, 14, False, PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════
# 슬라이드 6: new-nogood 5단계
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  new-nogood: 금지 처리 5단계")

stages = [
    ("① 금지 도장", "env.nogood = reason\n이 환경은 모순이다!",
     RGBColor(255, 225, 225), ACCENT),
    ("② label 제거", "remove-env-from-labels\n모순 세계에서 참은 무의미",
     RGBColor(255, 240, 220), ORANGE),
    ("③ 테이블 등록", "nogood-table에 등록\n이후 검사 시 참조",
     RGBColor(220, 235, 255), NAVY),
    ("④ 상위 정리", "{A,B} 등록 시\n{A,B,C} 제거 (최소성)",
     RGBColor(220, 245, 220), GREEN),
    ("⑤ 상향 오염", "env-table의 상위집합\n전부 자동 오염!",
     RGBColor(240, 220, 255), RGBColor(100, 50, 150)),
]

for i, (title, desc, bg_c, text_c) in enumerate(stages):
    x = 0.3 + i * 2.6
    add_box(sl, x, 1.5, 2.4, 0.7, title, NAVY, WHITE, 14, True)
    add_box(sl, x, 2.35, 2.4, 1.5, desc, bg_c, text_c, 12)

# 비유
add_box(sl, 0.5, 4.3, 12.3, 2.8,
        '비유: "금지 식재료 조합" 시스템\n\n'
        '셰프(ATMS)가 식재료(가정)를 조합해서 요리(결론)를 만든다.\n'
        '"땅콩 + 새우"가 독이라는 사실을 발견!\n\n'
        '① "땅콩+새우" 조합에 "독!" 딱지  '
        '② 이 조합을 쓴 모든 요리를 메뉴에서 제거\n'
        '③ 금지 목록에 등록  '
        '④ "땅콩+새우+우유"도 이미 금지? → 큰 건 제거\n'
        '⑤ "땅콩+새우+달걀" 같은 상위 조합도 자동 금지!',
        LIGHT_YELLOW, DARK_TEXT, 14, False, PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════
# 슬라이드 7: step-1 예제
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  실전 예제: step-1 (de Kleer 논문)")

# 왼쪽: 추론 과정
add_box(sl, 0.3, 1.4, 6.2, 5.8, "", WHITE)
add_textbox(sl, 0.5, 1.5, 5.8, 0.4, "추론 과정", 16, True, NAVY)
add_bullet_list(sl, 0.5, 2.0, 5.8, 5.0, [
    ("가정: A, B, C", 0, True, DARK_TEXT),
    ("J1: A → x=1     x=1.label = [{A}]", 0, False, GRAY),
    ("J2: B → y=x     y=x.label = [{B}]", 0, False, GRAY),
    ("J3: C → x=z     x=z.label = [{C}]", 0, False, GRAY),
    ("", 0, False, DARK_TEXT),
    ("nogood({A, B})  →  {A,B}* 금지!", 0, True, ACCENT),
    ("", 0, False, DARK_TEXT),
    ("J4: x=1 + y=x → y=1", 0, True, DARK_TEXT),
    ("  {A} x {B} = {A,B} → nogood! 차단!", 1, False, ACCENT),
    ("  y=1.label = []  (증명 불가)", 1, False, GRAY),
    ("", 0, False, DARK_TEXT),
    ("Premise: z=1  (무조건 참)", 0, True, DARK_TEXT),
    ("J5: z=1 + x=z → x=1", 0, False, GREEN),
    ("  x=1.label = [{A}, {C}]  두 경로!", 1, False, GREEN),
    ("", 0, False, DARK_TEXT),
    ("J4 재전파:", 0, True, GREEN),
    ("  {C} x {B} = {B,C} → 유효!", 1, False, GREEN),
    ("  y=1.label = [{B,C}]  우회 성공!", 1, True, GREEN),
], font_size=13)

# 오른쪽: 핵심 교훈
add_box(sl, 6.8, 1.4, 6.2, 2.5,
        "JTMS였다면?\n\n"
        "  {A,B} 모순 발견\n"
        "  → A 철회\n"
        "  → 모든 추론 재시작\n"
        "  → C 경로 다시 탐색...\n"
        "  → 비용 큼!",
        RGBColor(255, 235, 235), ACCENT, 14, False, PP_ALIGN.LEFT)

add_box(sl, 6.8, 4.1, 6.2, 3.1,
        "ATMS는?\n\n"
        "  {A,B}만 격리 (나머지 유효)\n"
        "  x=1에 {C} 경로 자동 추가\n"
        "  J4가 자동 재전파\n"
        "  {B,C} 우회 경로 자동 발견!\n\n"
        "  재계산 없음!\n"
        "  이것이 ATMS의 진가!",
        RGBColor(230, 255, 230), GREEN, 14, False, PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════
# 슬라이드 8: 응용 분야
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  ATMS의 응용 분야")

apps = [
    ("의료 진단", "여러 질병 가설을\n동시에 추적\n검사 결과로\n가설 제거",
     RGBColor(255, 230, 230)),
    ("회로 진단", "고장 부품 후보를\n병렬 추적\n(de Kleer의\n원래 목적)",
     RGBColor(230, 240, 255)),
    ("설계 탐색", "여러 설계 대안을\n동시에 유지\n제약 위반 시\n자동 격리",
     RGBColor(230, 255, 230)),
    ("계획 수립", "여러 행동 계획을\n병렬 평가\n실패 시 대안\n자동 활성화",
     RGBColor(255, 245, 220)),
]

for i, (title, desc, color) in enumerate(apps):
    x = 0.5 + i * 3.2
    add_box(sl, x, 1.5, 2.9, 0.8, title, NAVY, WHITE, 18, True)
    add_box(sl, x, 2.5, 2.9, 2.5, desc, color, DARK_TEXT, 15)

add_box(sl, 0.5, 5.5, 12.3, 1.5,
        "공통 패턴: 여러 가설/대안을 동시에 유지하다가\n"
        "모순이 발견되면 해당 가설만 제거하고 나머지는 계속 작동\n"
        "→ 백트래킹 없이 자동 우회!",
        LIGHT_YELLOW, NAVY, 16, True)

# ══════════════════════════════════════════════════════════
# 슬라이드 9: 설계 원칙
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title_bar(sl, "  de Kleer의 핵심 설계 원칙")

principles = [
    ("1. 발생 시점 처리",
     "매번 전수검사 대신\n발생 시 한 번에 전파",
     RGBColor(220, 235, 255)),
    ("2. 최소성 유지",
     "테이블에 최소 집합만\n메모리+속도 최적화",
     RGBColor(220, 245, 220)),
    ("3. 부분집합 판정",
     "{A,B} ⊆ X 관계로만\n오염 → 정확하고 안전",
     RGBColor(255, 240, 220)),
    ("4. 역참조 활용",
     "env.nodes로\nO(k) 빠른 label 정리",
     RGBColor(240, 220, 255)),
    ("5. 이유 기록",
     "'왜 금지?'를 저장\n설명 생성 가능",
     RGBColor(255, 230, 230)),
]

for i, (title, desc, color) in enumerate(principles):
    x = 0.3 + i * 2.6
    add_box(sl, x, 1.5, 2.4, 0.7, title, NAVY, WHITE, 13, True)
    add_box(sl, x, 2.4, 2.4, 1.4, desc, color, DARK_TEXT, 13)

# Label Antichain 불변량
add_box(sl, 0.5, 4.3, 12.3, 2.8,
        "Label의 Antichain 불변량 (de Kleer가 확립한 핵심 이론)\n\n"
        "어떤 노드의 label에 있는 임의의 두 환경 E1, E2에 대해,\n"
        "E1 ⊄ E2 이고 E2 ⊄ E1 이다.\n"
        "(어느 것도 다른 것의 부분집합이 아니다)\n\n"
        "효과: label이 항상 최소 크기로 유지됨\n"
        '       = "이 노드가 참이 되는 가장 약한 조건들"을 정확히 표현',
        RGBColor(255, 250, 235), NAVY, 14, False, PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════
# 슬라이드 10: 마무리
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, NAVY)
add_textbox(sl, 1, 1.0, 11, 0.8, "한 문장 정의", 20, False,
            RGBColor(180, 200, 240), PP_ALIGN.CENTER)
add_box(sl, 1.0, 2.0, 11.3, 2.0,
        'ATMS는 "여러 가능한 세계를 동시에 추적하면서,\n'
        '모순이 발견되면 해당 세계만 격리하고\n'
        '나머지는 계속 작동하게 하는" 진리 유지 시스템이다.',
        RGBColor(255, 250, 235), NAVY, 22, True)

add_textbox(sl, 1, 4.5, 11, 0.8, "패러다임 전환", 20, False,
            RGBColor(180, 200, 240), PP_ALIGN.CENTER)

add_textbox(sl, 1, 5.2, 5, 0.5,
            'JTMS: "X는 참이다"', 18, False, ACCENT, PP_ALIGN.CENTER)
add_textbox(sl, 6.5, 5.2, 6, 0.7,
            'ATMS: "X는 {A} 하에서 참이고,\n       {B,C} 하에서도 참이다"',
            18, True, RGBColor(100, 220, 120), PP_ALIGN.CENTER)

add_textbox(sl, 1, 6.5, 11, 0.5,
            "Building Problem Solvers — K.D. Forbus & J. de Kleer, 1993",
            12, False, RGBColor(150, 170, 200), PP_ALIGN.CENTER)

# ── 저장 ──
output_path = "/home/user/bps/ATMS_lecture.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
