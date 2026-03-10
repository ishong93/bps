#!/usr/bin/env python3
"""LTMS 개념 학습용 PowerPoint 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)       # 진한 남색 배경
BG_SECTION = RGBColor(0x16, 0x21, 0x3E)     # 섹션 배경
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)    # 강조 하늘색
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)   # 강조 초록
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)  # 강조 주황
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)     # 강조 빨강
ACCENT_PURPLE = RGBColor(0xA5, 0x5E, 0xFF)  # 강조 보라
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DIM_GRAY = RGBColor(0x88, 0x88, 0xAA)
CARD_BG = RGBColor(0x24, 0x2B, 0x4A)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=16,
                           default_color=LIGHT_GRAY, font_name="맑은 고딕"):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, size, color, bold, align = line_info, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = line_info[0]
            size = line_info[1] if len(line_info) > 1 else default_size
            color = line_info[2] if len(line_info) > 2 else default_color
            bold = line_info[3] if len(line_info) > 3 else False
            align = line_info[4] if len(line_info) > 4 else PP_ALIGN.LEFT

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_code_box(slide, left, top, width, height, code_text, font_size=12):
    shape = add_shape_rect(slide, left, top, width, height, RGBColor(0x0D, 0x11, 0x1E),
                           border_color=RGBColor(0x33, 0x44, 0x66))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(1)
    return shape


def add_arrow_shape(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════
#  슬라이드 생성 함수들
# ══════════════════════════════════════════════════════════

def slide_title(prs):
    """슬라이드 1: 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # 상단 장식선
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT_BLUE)

    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
                "LTMS", 54, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                "Logic-based Truth Maintenance System", 30, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
                "논리 기반 진리 유지 시스템의 구조와 원리", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(4.5), Inches(4.6), Inches(4.3), Pt(2), DIM_GRAY)

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                "Building Problem Solvers  |  Forbus & de Kleer", 16, DIM_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                "15개 파일  ·  ~3,552줄  ·  Common Lisp", 14, DIM_GRAY, False, PP_ALIGN.CENTER)


def slide_toc(prs):
    """슬라이드 2: 목차"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
                "목차", 36, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(2), Pt(3), ACCENT_BLUE)

    chapters = [
        ("01", "LTMS란 무엇인가?", "정의와 필요성"),
        ("02", "핵심 자료구조", "Node, Clause, LTMS"),
        ("03", "핵심 알고리즘", "BCP, 단위 전파, 역전파"),
        ("04", "수식 정규화", "논리식 → CNF 변환"),
        ("05", "예제: 자동 연쇄 추론", "하나의 사실에서 모든 결론 도출"),
        ("06", "예제: 모순 감지", "일관성 없는 지식 발견"),
        ("07", "예제: 가정 철회", "비단조적 추론"),
        ("08", "예제: 마르크스 형제 퍼즐", "제약 만족 문제"),
        ("09", "코드 파일 구조", "15개 파일의 역할"),
        ("10", "요약 및 현대적 의의", "SAT 솔버와의 관계"),
    ]

    for i, (num, title, desc) in enumerate(chapters):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(1.6 + row * 1.0)

        add_shape_rect(slide, x, y, Inches(5.5), Inches(0.85), CARD_BG, border_color=RGBColor(0x33, 0x44, 0x66))
        add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.4),
                    num, 14, ACCENT_BLUE, True)
        add_textbox(slide, x + Inches(0.6), y + Inches(0.05), Inches(4.5), Inches(0.4),
                    title, 16, WHITE, True)
        add_textbox(slide, x + Inches(0.6), y + Inches(0.45), Inches(4.5), Inches(0.35),
                    desc, 12, DIM_GRAY)


def slide_what_is_ltms(prs):
    """슬라이드 3: LTMS란 무엇인가?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                "01  LTMS란 무엇인가?", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 정의 박스
    defn_shape = add_shape_rect(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.3),
                                CARD_BG, ACCENT_BLUE)
    add_multiline_textbox(slide, Inches(1.1), Inches(1.6), Inches(11), Inches(1.2), [
        ("LTMS = Logic-based Truth Maintenance System", 20, WHITE, True),
        ("", 8),
        ("명제(proposition)들의 참/거짓 상태를 논리적 제약 하에서", 16, LIGHT_GRAY),
        ("일관성 있게 유지하면서 자동으로 추론하는 시스템", 16, LIGHT_GRAY),
    ])

    # 4개 핵심 능력
    capabilities = [
        ("가정 관리", "가정을 세우고\n철회할 수 있음", ACCENT_BLUE),
        ("제약 전파", "절(clause) 형태의\n제약 조건을 자동 전파", ACCENT_GREEN),
        ("모순 감지", "일관성 없는 상태를\n자동으로 발견", ACCENT_RED),
        ("원인 추적", "모순의 근본 원인을\n추적하여 해결", ACCENT_ORANGE),
    ]

    for i, (title, desc, color) in enumerate(capabilities):
        x = Inches(0.8 + i * 3.0)
        y = Inches(3.3)
        card = add_shape_rect(slide, x, y, Inches(2.7), Inches(2.0), CARD_BG, color)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.3), Inches(0.5),
                    title, 18, color, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.3), Inches(1.0),
                    desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # 하단 비유
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
                "비유: 스프레드시트의 수식처럼, 하나의 값이 바뀌면 관련된 모든 셀이 자동으로 업데이트됨.\n      단, LTMS는 '논리적 참/거짓'을 다루며, 모순이 발생하면 경고해줌.",
                14, DIM_GRAY)


def slide_data_structures(prs):
    """슬라이드 4: 핵심 자료구조"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                "02  핵심 자료구조", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # TMS-Node
    add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(3.9), Inches(3.5), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(3.5), Inches(0.5),
                "TMS-Node (명제 노드)", 18, ACCENT_BLUE, True)
    add_code_box(slide, Inches(0.7), Inches(2.1), Inches(3.5), Inches(2.6),
                 "index     ; 고유 식별자\n"
                 "datum     ; 추론 데이터\n"
                 "label     ; :UNKNOWN\n"
                 "          ; :TRUE\n"
                 "          ; :FALSE\n"
                 "support   ; 지지하는 절\n"
                 "true-clauses  ; 참인 절들\n"
                 "false-clauses ; 거짓인 절들\n"
                 "assumption?   ; 가정 여부", 11)

    # Clause
    add_shape_rect(slide, Inches(4.7), Inches(1.4), Inches(3.9), Inches(3.5), CARD_BG, ACCENT_GREEN)
    add_textbox(slide, Inches(4.9), Inches(1.5), Inches(3.5), Inches(0.5),
                "Clause (절 = 리터럴의 OR)", 18, ACCENT_GREEN, True)
    add_code_box(slide, Inches(4.9), Inches(2.1), Inches(3.5), Inches(2.6),
                 "literals  ; 리터럴 리스트\n"
                 "  ((n1 . :TRUE)\n"
                 "   (n2 . :FALSE) ...)\n"
                 "pvs    ; 미결정 리터럴 수\n"
                 "  = 0 → 모순 발생!\n"
                 "  = 1 → 단위 전파!\n"
                 "sats   ; 만족 리터럴 수\n"
                 "length ; 전체 길이\n"
                 "status ; 절 상태", 11)

    # LTMS
    add_shape_rect(slide, Inches(8.9), Inches(1.4), Inches(3.9), Inches(3.5), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(9.1), Inches(1.5), Inches(3.5), Inches(0.5),
                "LTMS (시스템 전체)", 18, ACCENT_ORANGE, True)
    add_code_box(slide, Inches(9.1), Inches(2.1), Inches(3.5), Inches(2.6),
                 "nodes    ; 노드 해시테이블\n"
                 "clauses  ; 모든 절 리스트\n"
                 "queue    ; 검사 대기 큐\n"
                 "violated-clauses\n"
                 "         ; 위반된 절들\n"
                 "contradiction-\n"
                 "  handlers ; 모순 처리기\n"
                 "complete ; 완전 LTMS?\n"
                 "enqueue-procedure", 11)

    # pvs 설명 강조 박스
    add_shape_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), CARD_BG, ACCENT_RED)
    add_multiline_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.8), Inches(1.3), [
        ("핵심 개념: pvs (Potentially Violating terms) 카운터", 18, ACCENT_RED, True),
        ("", 6),
        ("절 내에서 아직 해당 절을 '위반할 가능성이 있는' 리터럴의 수", 15, LIGHT_GRAY),
        ("pvs = 0  →  모든 리터럴이 거짓  →  모순(contradiction) 발생!", 15, ACCENT_RED, True),
        ("pvs = 1  →  딱 하나만 남음  →  그 리터럴을 참으로 만듦 (단위 전파)", 15, ACCENT_GREEN, True),
    ])


def slide_algorithm_overview(prs):
    """슬라이드 5: 핵심 알고리즘 전체 흐름"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                "03  핵심 알고리즘 — 전체 흐름", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 흐름도를 박스+화살표로 구성
    steps = [
        ("1. 가정 설정", "enable-assumption", "노드의 label을\n:TRUE 또는 :FALSE로 설정", ACCENT_BLUE),
        ("2. 진리값 전파", "set-truth", "관련 절의 pvs/sats\n카운터 업데이트", ACCENT_GREEN),
        ("3. BCP", "check-clauses", "pvs=1인 절 발견 시\n단위 전파 → 2로 복귀", ACCENT_ORANGE),
        ("4. 모순 검사", "check-for-contradictions", "pvs=0인 절 발견 시\n모순 핸들러 호출", ACCENT_RED),
    ]

    for i, (title, func, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        y = Inches(1.5)

        card = add_shape_rect(slide, x, y, Inches(2.8), Inches(2.2), CARD_BG, color)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.4),
                    title, 16, color, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.5), Inches(0.4),
                    func, 13, ACCENT_GREEN, False, PP_ALIGN.CENTER, "Consolas")
        add_textbox(slide, x + Inches(0.15), y + Inches(1.1), Inches(2.5), Inches(0.9),
                    desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

        # 화살표
        if i < 3:
            arrow_x = x + Inches(2.85)
            slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Inches(0.9),
                                   Inches(0.3), Inches(0.3)).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = DIM_GRAY
            slide.shapes[-1].line.fill.background()

    # 역전파 설명
    add_shape_rect(slide, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.8), CARD_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(0.7), Inches(4.3), Inches(5.6), Inches(0.5),
                "5. 가정 철회 시: propagate-unknownness", 17, ACCENT_PURPLE, True)
    add_code_box(slide, Inches(0.7), Inches(4.9), Inches(5.6), Inches(1.9),
                 "1. 철회된 노드 → :UNKNOWN으로 복원\n"
                 "2. 이 노드에 의존한 절들의 pvs/sats 복원\n"
                 "3. pvs가 1→2 된 절의 결론도 → :UNKNOWN\n"
                 "4. 연쇄적으로 모든 의존 추론 무효화\n"
                 "5. 남은 근거로 재전파 → 새 안정 상태", 13)

    # BCP 루프 설명
    add_shape_rect(slide, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.8), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(7.0), Inches(4.3), Inches(5.6), Inches(0.5),
                "BCP 루프 상세 (check-clauses)", 17, ACCENT_ORANGE, True)
    add_code_box(slide, Inches(7.0), Inches(4.9), Inches(5.6), Inches(1.9),
                 "WHILE 큐에 절이 남아있으면:\n"
                 "  절을 꺼내서 검사\n"
                 "  IF pvs = 0:\n"
                 "    → violated-clauses에 추가\n"
                 "  IF pvs = 1:\n"
                 "    → 미결정 리터럴 찾기\n"
                 "    → set-truth 호출 (단위 전파)\n"
                 "    → 새로 영향받는 절들 큐에 추가", 13)


def slide_formula_normalization(prs):
    """슬라이드 6: 수식 정규화"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                "04  수식 정규화 — 논리식을 CNF로 변환", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 지원 연산자
    add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(2.8), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.0), Inches(0.4),
                "지원하는 논리 연산자", 18, ACCENT_BLUE, True)
    ops = [
        (":AND      논리곱        A ∧ B"),
        (":OR       논리합        A ∨ B"),
        (":NOT      부정          ¬A"),
        (":IMPLIES  함의          A→B  ≡  ¬A∨B"),
        (":IFF      동치          A↔B  ≡  (A→B)∧(B→A)"),
        (":TAXONOMY 상호배타 선택  정확히 하나만 참"),
    ]
    add_code_box(slide, Inches(0.7), Inches(2.0), Inches(5.1), Inches(2.0),
                 "\n".join(ops), 12)

    # 변환 예시
    add_shape_rect(slide, Inches(6.3), Inches(1.4), Inches(6.5), Inches(2.8), CARD_BG, ACCENT_GREEN)
    add_textbox(slide, Inches(6.5), Inches(1.5), Inches(6.0), Inches(0.4),
                "변환 예시: (:IMPLIES (:AND p q) r)", 17, ACCENT_GREEN, True)
    add_code_box(slide, Inches(6.5), Inches(2.0), Inches(6.1), Inches(2.0),
                 "입력: (p ∧ q) → r\n"
                 "\n"
                 "1단계 - 함의 제거:\n"
                 "   ¬(p ∧ q) ∨ r\n"
                 "\n"
                 "2단계 - 드모르간:\n"
                 "   (¬p ∨ ¬q) ∨ r  =  ¬p ∨ ¬q ∨ r\n"
                 "\n"
                 "결과: 절 하나  {¬p, ¬q, r}", 12)

    # 왜 CNF인가?
    add_shape_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.2), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
                "왜 CNF(Conjunctive Normal Form)로 변환하는가?", 18, ACCENT_ORANGE, True)
    add_multiline_textbox(slide, Inches(0.7), Inches(5.2), Inches(11.8), Inches(1.5), [
        ("CNF = 절들의 AND  /  각 절 = 리터럴들의 OR", 16, WHITE, True),
        ("", 6),
        ("• BCP(Boolean Constraint Propagation)는 CNF 형태에서만 효율적으로 동작", 14, LIGHT_GRAY),
        ("• 절 하나가 위반되면(pvs=0) 즉시 모순을 감지할 수 있음", 14, LIGHT_GRAY),
        ("• 절에 미결정 리터럴이 하나만 남으면(pvs=1) 단위 전파가 가능", 14, LIGHT_GRAY),
        ("• 현대 SAT 솔버(MiniSat, Z3 등)도 동일한 원리를 사용", 14, ACCENT_GREEN),
    ])


def slide_example_chain(prs):
    """슬라이드 7: 예제 - 자동 연쇄 추론"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "05  예제: 자동 연쇄 추론", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 문제 설정
    add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(5.0), Inches(2.0), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(4.5), Inches(0.4),
                "문제 설정", 18, ACCENT_BLUE, True)
    add_code_box(slide, Inches(0.7), Inches(2.0), Inches(4.6), Inches(1.2),
                 "절 1:  x ∨ y      (x 또는 y는 참)\n"
                 "절 2:  ¬y ∨ z     (y이면 z)\n"
                 "절 3:  ¬z ∨ r     (z이면 r)\n"
                 "가정:  x = FALSE", 13)

    # LTMS 추론 과정 - 단계별
    add_shape_rect(slide, Inches(5.8), Inches(1.4), Inches(7.0), Inches(5.5), CARD_BG, ACCENT_GREEN)
    add_textbox(slide, Inches(6.0), Inches(1.5), Inches(6.5), Inches(0.4),
                "LTMS 자동 추론 과정", 18, ACCENT_GREEN, True)

    steps_code = (
        "1단계: x = FALSE (가정 설정)\n"
        "  → 절1 (x∨y)에서 x가 FALSE\n"
        "  → pvs: 2→1  →  단위 전파!\n"
        "  → y = TRUE\n"
        "\n"
        "2단계: y = TRUE (전파 결과)\n"
        "  → 절2 (¬y∨z)에서 ¬y가 FALSE\n"
        "  → pvs: 2→1  →  단위 전파!\n"
        "  → z = TRUE\n"
        "\n"
        "3단계: z = TRUE (전파 결과)\n"
        "  → 절3 (¬z∨r)에서 ¬z가 FALSE\n"
        "  → pvs: 2→1  →  단위 전파!\n"
        "  → r = TRUE"
    )
    add_code_box(slide, Inches(6.0), Inches(2.0), Inches(6.6), Inches(4.1), steps_code, 13)

    # 결과 요약
    add_shape_rect(slide, Inches(0.5), Inches(3.8), Inches(5.0), Inches(1.5), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(4.5), Inches(0.4),
                "추론 연쇄 결과", 18, ACCENT_ORANGE, True)
    add_code_box(slide, Inches(0.7), Inches(4.4), Inches(4.6), Inches(0.7),
                 "x=FALSE → y=TRUE → z=TRUE → r=TRUE\n"
                 "하나의 가정으로 3개의 결론 자동 도출!", 13)

    # 핵심 교훈
    add_shape_rect(slide, Inches(0.5), Inches(5.7), Inches(5.0), Inches(1.2), CARD_BG, DIM_GRAY)
    add_multiline_textbox(slide, Inches(0.7), Inches(5.8), Inches(4.5), Inches(1.0), [
        ("핵심: 사람이 하나씩 따질 필요 없이", 14, WHITE, True),
        ("LTMS가 BCP로 모든 결론을 자동 도출", 14, ACCENT_GREEN, True),
    ])


def slide_example_contradiction(prs):
    """슬라이드 8: 예제 - 모순 감지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "06  예제: 모순 감지", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 문제 설정
    add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(2.2), CARD_BG, ACCENT_RED)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.0), Inches(0.4),
                "문제: 서로 충돌하는 가정", 18, ACCENT_RED, True)
    add_code_box(slide, Inches(0.7), Inches(2.0), Inches(5.1), Inches(1.4),
                 "가정 1:  n1 = FALSE\n"
                 "가정 2:  n2 = FALSE\n"
                 "제약:    n1 ∨ n2  (둘 중 하나는 참)\n"
                 "\n"
                 "→ 모순! 둘 다 거짓일 수 없음", 13)

    # 감지 과정
    add_shape_rect(slide, Inches(6.3), Inches(1.4), Inches(6.5), Inches(2.2), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(6.5), Inches(1.5), Inches(6.0), Inches(0.4),
                "LTMS의 모순 감지 과정", 18, ACCENT_ORANGE, True)
    add_code_box(slide, Inches(6.5), Inches(2.0), Inches(6.1), Inches(1.4),
                 "절 (n1 ∨ n2) 검사:\n"
                 "  n1=FALSE → pvs: 2→1\n"
                 "  n2=FALSE → pvs: 1→0\n"
                 "  pvs = 0 → 위반! 모순 발생!\n"
                 "→ contradiction-handler 호출", 13)

    # 모순 처리 방법들
    add_textbox(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.5),
                "모순 처리 전략 3가지", 20, WHITE, True)

    strategies = [
        ("ask-user-handler", "사용자에게 질문", "모순 원인 가정을 보여주고\n어느 것을 철회할지 선택", ACCENT_BLUE),
        ("avoid-all", "자동 회피", "첫 번째 원인 가정을\n자동으로 철회", ACCENT_GREEN),
        ("nogood 학습", "재발 방지", "모순 조합을 기억하여\n같은 실패 반복 방지", ACCENT_PURPLE),
    ]

    for i, (name, title, desc, color) in enumerate(strategies):
        x = Inches(0.5 + i * 4.2)
        y = Inches(4.6)
        add_shape_rect(slide, x, y, Inches(3.8), Inches(2.2), CARD_BG, color)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.4),
                    title, 17, color, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(0.35),
                    name, 12, ACCENT_GREEN, False, PP_ALIGN.CENTER, "Consolas")
        add_textbox(slide, x + Inches(0.15), y + Inches(1.1), Inches(3.5), Inches(0.9),
                    desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)


def slide_example_retraction(prs):
    """슬라이드 9: 예제 - 가정 철회와 비단조 추론"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "07  예제: 가정 철회 — 비단조적 추론", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 문제 설정
    add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(4.0), Inches(1.6), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(3.5), Inches(0.4),
                "문제 설정", 18, ACCENT_BLUE, True)
    add_code_box(slide, Inches(0.7), Inches(2.0), Inches(3.6), Inches(0.8),
                 "절 1:  x ∨ z\n"
                 "절 2:  y ∨ z", 13)

    # 단계 1
    add_shape_rect(slide, Inches(4.8), Inches(1.4), Inches(3.8), Inches(1.6), CARD_BG, ACCENT_GREEN)
    add_textbox(slide, Inches(5.0), Inches(1.5), Inches(3.4), Inches(0.4),
                "단계 1: 가정 설정", 16, ACCENT_GREEN, True)
    add_code_box(slide, Inches(5.0), Inches(2.0), Inches(3.4), Inches(0.8),
                 "x=FALSE → z=TRUE (절1)\n"
                 "y=FALSE\n"
                 "결과: z=TRUE (근거: 절1)", 12)

    # 단계 2
    add_shape_rect(slide, Inches(8.9), Inches(1.4), Inches(3.8), Inches(1.6), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(9.1), Inches(1.5), Inches(3.4), Inches(0.4),
                "단계 2: x 가정 철회!", 16, ACCENT_ORANGE, True)
    add_code_box(slide, Inches(9.1), Inches(2.0), Inches(3.4), Inches(0.8),
                 "retract-assumption(x)\n"
                 "x → :UNKNOWN\n"
                 "z의 근거(절1) 무효화!", 12)

    # 역전파 과정
    add_shape_rect(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.5), CARD_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(0.7), Inches(3.5), Inches(11.8), Inches(0.4),
                "propagate-unknownness 상세 과정", 18, ACCENT_PURPLE, True)

    add_code_box(slide, Inches(0.7), Inches(4.0), Inches(5.5), Inches(2.7),
                 "1. x 철회 → x = :UNKNOWN\n"
                 "2. z의 support가 절1이었음\n"
                 "   절1의 pvs 업데이트: pvs 증가\n"
                 "3. z = :UNKNOWN (근거 상실)\n"
                 "\n"
                 "4. 재전파 시작:\n"
                 "   절2 (y ∨ z) 검사\n"
                 "   y=FALSE → pvs=1\n"
                 "   → 단위 전파: z = TRUE!\n"
                 "   (이번엔 절2가 근거)", 12)

    add_shape_rect(slide, Inches(6.5), Inches(4.0), Inches(6.1), Inches(1.5), RGBColor(0x1A, 0x3A, 0x1A), ACCENT_GREEN)
    add_multiline_textbox(slide, Inches(6.7), Inches(4.1), Inches(5.7), Inches(1.3), [
        ("최종 결과", 17, ACCENT_GREEN, True),
        ("", 6),
        ("x = :UNKNOWN  (철회됨)", 15, LIGHT_GRAY),
        ("y = FALSE", 15, LIGHT_GRAY),
        ("z = TRUE  (근거가 절1 → 절2로 전환!)", 15, WHITE, True),
    ])

    add_shape_rect(slide, Inches(6.5), Inches(5.8), Inches(6.1), Inches(1.0), CARD_BG, DIM_GRAY)
    add_multiline_textbox(slide, Inches(6.7), Inches(5.9), Inches(5.7), Inches(0.8), [
        ("핵심: 가정을 바꿔도 다른 근거가 있으면", 14, WHITE),
        ("결론을 유지함 → 'What-if' 분석이 가능", 14, ACCENT_GREEN, True),
    ])


def slide_example_marx(prs):
    """슬라이드 10: 마르크스 형제 퍼즐"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "08  예제: 마르크스 형제 퍼즐 — 제약 만족 문제", 30, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 단서
    add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(3.0), CARD_BG, ACCENT_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.4), Inches(0.4),
                "7가지 단서", 18, ACCENT_BLUE, True)
    add_code_box(slide, Inches(0.7), Inches(2.0), Inches(5.4), Inches(2.2),
                 "1. 피아니스트, 하프연주자, 말잘하는사람\n"
                 "   = 서로 다른 형제\n"
                 "2. 돈/도박/동물 좋아하는 사람 = 서로 다름\n"
                 "3. 말잘하는 사람은 도박 싫어함\n"
                 "4. 동물 좋아하는 사람 = 하프 연주자\n"
                 "5. Groucho는 동물 싫어함\n"
                 "6. Harpo는 항상 말이 없음\n"
                 "7. Chico는 피아노 연주", 12)

    # 논리식 변환
    add_shape_rect(slide, Inches(6.6), Inches(1.4), Inches(6.2), Inches(3.0), CARD_BG, ACCENT_GREEN)
    add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.4),
                "LTMS 논리식으로 변환 (marxdata.lisp)", 16, ACCENT_GREEN, True)
    add_code_box(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2),
                 ";; 단서 1: 피아노≠하프\n"
                 "(pairwise-nogood plays-piano plays-harp)\n"
                 ";; 단서 4: 동물=하프\n"
                 "(same-entity likes-animals plays-harp)\n"
                 ";; 단서 5: Groucho≠동물\n"
                 "(:not (likes-animals groucho))\n"
                 ";; 단서 6: Harpo≠말잘함\n"
                 "(:not (smooth-talker harpo))\n"
                 ";; 단서 7: Chico=피아노\n"
                 "(plays-piano chico)", 11)

    # DDS 탐색 과정
    add_shape_rect(slide, Inches(0.5), Inches(4.7), Inches(7.5), Inches(2.2), CARD_BG, ACCENT_ORANGE)
    add_textbox(slide, Inches(0.7), Inches(4.8), Inches(7.0), Inches(0.4),
                "DDS(의존성 지향 탐색) 풀이 과정", 17, ACCENT_ORANGE, True)
    add_code_box(slide, Inches(0.7), Inches(5.3), Inches(7.1), Inches(1.4),
                 "가정: (plays-piano groucho)\n"
                 "  → Chico가 이미 피아노 + pairwise-nogood → 모순!\n"
                 "  → nogood 학습: 다시는 이 조합 시도 안 함\n"
                 "가정: (plays-harp harpo)\n"
                 "  → same-entity 규칙 → (likes-animals harpo) = TRUE\n"
                 "  → 소거법 → (smooth-talker groucho) = TRUE → 해 발견!", 11)

    # 결과
    add_shape_rect(slide, Inches(8.3), Inches(4.7), Inches(4.5), Inches(2.2),
                   RGBColor(0x1A, 0x3A, 0x1A), ACCENT_GREEN)
    add_multiline_textbox(slide, Inches(8.5), Inches(4.8), Inches(4.1), Inches(2.0), [
        ("정답", 18, ACCENT_GREEN, True),
        ("", 8),
        ("Chico  = 피아노, 도박", 16, WHITE, True),
        ("Harpo  = 하프, 동물", 16, WHITE, True),
        ("Groucho = 말잘함, 돈", 16, WHITE, True),
        ("", 8),
        ("모순에서 배우며(nogood) 효율적 탐색!", 13, ACCENT_GREEN),
    ])


def slide_file_structure(prs):
    """슬라이드 11: 코드 파일 구조"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                "09  코드 파일 구조 — 15개 파일의 역할", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 3개 계층
    layers = [
        ("핵심 엔진", ACCENT_BLUE, [
            ("ltms.lisp", "746줄", "노드, 절, BCP, 단위전파, 역전파"),
            ("cltms.lisp", "426줄", "완전 LTMS, 해소, 트라이 인덱싱"),
        ]),
        ("추론 통합 (LTRE)", ACCENT_GREEN, [
            ("linter.lisp", "101줄", "LTRE 정의 (LTMS+DB 통합)"),
            ("ldata.lisp", "305줄", "사실/가정 관리 (assert!, assume!)"),
            ("lrules.lisp", "286줄", "패턴 매칭 규칙 엔진"),
            ("funify.lisp", "131줄", "패턴 통합(unification)"),
        ]),
        ("확장 모듈", ACCENT_ORANGE, [
            ("dds.lisp", "83줄", "의존성 지향 탐색"),
            ("cwa.lisp", "225줄", "닫힌 세계 가정(CWA)"),
            ("indirect.lisp", "47줄", "간접 증명"),
            ("setrule.lisp", "41줄", "집합 제약 규칙"),
        ]),
    ]

    y_pos = Inches(1.4)
    for layer_name, color, files in layers:
        add_shape_rect(slide, Inches(0.5), y_pos, Inches(12.3),
                       Inches(0.35 + len(files) * 0.4), CARD_BG, color)
        add_textbox(slide, Inches(0.7), y_pos + Inches(0.05), Inches(3), Inches(0.35),
                    layer_name, 16, color, True)

        for j, (fname, lines, desc) in enumerate(files):
            fy = y_pos + Inches(0.4 + j * 0.4)
            add_textbox(slide, Inches(1.0), fy, Inches(2.2), Inches(0.35),
                        fname, 13, ACCENT_GREEN, False, PP_ALIGN.LEFT, "Consolas")
            add_textbox(slide, Inches(3.2), fy, Inches(1.0), Inches(0.35),
                        lines, 12, DIM_GRAY)
            add_textbox(slide, Inches(4.2), fy, Inches(8.0), Inches(0.35),
                        desc, 13, LIGHT_GRAY)

        y_pos += Inches(0.55 + len(files) * 0.4)

    # 예제/테스트
    add_shape_rect(slide, Inches(0.5), y_pos, Inches(12.3), Inches(1.55), CARD_BG, ACCENT_PURPLE)
    add_textbox(slide, Inches(0.7), y_pos + Inches(0.05), Inches(3), Inches(0.35),
                "예제 / 테스트", 16, ACCENT_PURPLE, True)
    test_files = [
        ("ltms-ex.lisp", "896줄", "종합 테스트 및 예제"),
        ("marx.lisp", "54줄", "마르크스 형제 퍼즐"),
        ("marxdata.lisp", "81줄", "퍼즐 데이터 및 규칙"),
    ]
    for j, (fname, lines, desc) in enumerate(test_files):
        fy = y_pos + Inches(0.4 + j * 0.35)
        add_textbox(slide, Inches(1.0), fy, Inches(2.2), Inches(0.35),
                    fname, 13, ACCENT_GREEN, False, PP_ALIGN.LEFT, "Consolas")
        add_textbox(slide, Inches(3.2), fy, Inches(1.0), Inches(0.35),
                    lines, 12, DIM_GRAY)
        add_textbox(slide, Inches(4.2), fy, Inches(8.0), Inches(0.35),
                    desc, 13, LIGHT_GRAY)


def slide_summary(prs):
    """슬라이드 12: 요약 및 현대적 의의"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                "10  요약 — LTMS가 해결하는 6대 과제", 32, ACCENT_BLUE, True)
    add_shape_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)

    # 6대 과제 테이블
    tasks = [
        ("자동 추론", "A가 참이면 무엇이 따라오는가?", "BCP + 단위 전파", ACCENT_BLUE),
        ("모순 감지", "현재 지식에 모순이 있는가?", "pvs=0 감지", ACCENT_RED),
        ("원인 추적", "모순의 원인은 어떤 가정인가?", "assumptions-of-clause", ACCENT_ORANGE),
        ("가정 관리", "가정을 바꾸면 결론이 어떻게?", "propagate-unknownness", ACCENT_PURPLE),
        ("제약 만족", "모든 제약을 만족하는 해는?", "DDS + nogood 학습", ACCENT_GREEN),
        ("완전 추론", "도출 가능한 모든 사실은?", "Resolution (cltms)", ACCENT_BLUE),
    ]

    for i, (task, question, mechanism, color) in enumerate(tasks):
        row = i // 2
        col = i % 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.4 + row * 1.25)
        add_shape_rect(slide, x, y, Inches(5.9), Inches(1.1), CARD_BG, color)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.0), Inches(0.4),
                    task, 17, color, True)
        add_textbox(slide, x + Inches(2.2), y + Inches(0.05), Inches(3.5), Inches(0.4),
                    question, 13, LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(5.5), Inches(0.4),
                    mechanism, 12, ACCENT_GREEN, False, PP_ALIGN.LEFT, "Consolas")

    # 현대적 의의
    add_shape_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.7),
                   RGBColor(0x1A, 0x2A, 0x3A), ACCENT_BLUE)
    add_multiline_textbox(slide, Inches(0.7), Inches(5.3), Inches(11.8), Inches(1.5), [
        ("현대적 의의", 20, ACCENT_BLUE, True),
        ("", 6),
        ("• SAT 솔버 (MiniSat, CaDiCaL 등): LTMS의 BCP + nogood 학습이 핵심 엔진", 14, LIGHT_GRAY),
        ("• SMT 솔버 (Z3, CVC5): LTMS 위에 이론 추론기를 결합", 14, LIGHT_GRAY),
        ("• 지식 그래프 추론: 진리 유지 개념으로 일관성 관리", 14, LIGHT_GRAY),
        ("• AI 계획 시스템: 가정 기반 탐색 + 의존성 추적의 근간", 14, LIGHT_GRAY),
    ])


def slide_closing(prs):
    """슬라이드 13: 마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT_BLUE)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
                "Thank You", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(4.5), Inches(3.2), Inches(4.3), Pt(2), DIM_GRAY)

    add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
                "LTMS는 단순한 데이터 저장소가 아니라,", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                "지식의 일관성을 능동적으로 관리하는 추론 엔진입니다.", 20, WHITE, True, PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(2), Inches(5.4), Inches(9.33), Inches(1.0), CARD_BG, DIM_GRAY)
    add_multiline_textbox(slide, Inches(2.2), Inches(5.5), Inches(9), Inches(0.8), [
        ("효율성: BCP를 통한 빠른 제약 전파", 14, ACCENT_BLUE, False, PP_ALIGN.CENTER),
        ("정확성: 의존성 추적으로 정확한 추론 유지/철회", 14, ACCENT_GREEN, False, PP_ALIGN.CENTER),
        ("유연성: 가정 기반 what-if 분석", 14, ACCENT_ORANGE, False, PP_ALIGN.CENTER),
    ])


# ══════════════════════════════════════════════════════════
#  메인 실행
# ══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    # 와이드스크린 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_toc(prs)
    slide_what_is_ltms(prs)
    slide_data_structures(prs)
    slide_algorithm_overview(prs)
    slide_formula_normalization(prs)
    slide_example_chain(prs)
    slide_example_contradiction(prs)
    slide_example_retraction(prs)
    slide_example_marx(prs)
    slide_file_structure(prs)
    slide_summary(prs)
    slide_closing(prs)

    output_path = "/home/user/bps/ltms/LTMS_개념학습.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")
    print(f"총 {len(prs.slides)}장의 슬라이드")


if __name__ == "__main__":
    main()
