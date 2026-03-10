#!/usr/bin/env python3
"""JTMS vs LTMS vs ATMS 비교 학습용 PowerPoint 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_PURPLE = RGBColor(0xA5, 0x5E, 0xFF)
ACCENT_YELLOW = RGBColor(0xFF, 0xE0, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
DIM_GRAY = RGBColor(0x88, 0x88, 0xAA)
CARD_BG = RGBColor(0x24, 0x2B, 0x4A)

# 시스템별 대표 색상
JTMS_COLOR = ACCENT_BLUE
LTMS_COLOR = ACCENT_GREEN
ATMS_COLOR = ACCENT_ORANGE


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font="맑은 고딕"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_multiline(slide, left, top, width, height, lines, default_size=16,
                   default_color=LIGHT_GRAY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, size, color, bold, align = line_info, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = line_info[0]
            size = line_info[1] if len(line_info) > 1 else default_size
            color = line_info[2] if len(line_info) > 2 else default_color
            bold = line_info[3] if len(line_info) > 3 else False
            align = line_info[4] if len(line_info) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "맑은 고딕"
        p.alignment = align
        p.space_after = Pt(3)
    return box


def add_code(slide, left, top, width, height, code_text, font_size=11):
    shape = add_rect(slide, left, top, width, height,
                     RGBColor(0x0D, 0x11, 0x1E), RGBColor(0x33, 0x44, 0x66))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_right = Pt(10)
    tf.margin_bottom = Pt(6)
    for i, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(1)
    return shape


def slide_header(slide, number, title):
    """공통 슬라이드 헤더"""
    add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             f"{number}  {title}", 32, ACCENT_BLUE, True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(3), Pt(3), ACCENT_BLUE)


def three_columns(slide, y, height, contents):
    """3개 시스템 컬럼 (JTMS/LTMS/ATMS)"""
    colors = [JTMS_COLOR, LTMS_COLOR, ATMS_COLOR]
    names = ["JTMS", "LTMS", "ATMS"]
    for i in range(3):
        x = Inches(0.5 + i * 4.2)
        w = Inches(3.9)
        card = add_rect(slide, x, y, w, height, CARD_BG, colors[i])
        add_text(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.4),
                 names[i], 18, colors[i], True, PP_ALIGN.CENTER)
        # content
        if isinstance(contents[i], str):
            add_code(slide, x + Inches(0.1), y + Inches(0.5),
                     w - Inches(0.2), height - Inches(0.6), contents[i], 10)
        else:
            add_multiline(slide, x + Inches(0.15), y + Inches(0.5),
                          w - Inches(0.3), height - Inches(0.6), contents[i], 13)


# ══════════════════════════════════════════════════════════
#  개별 슬라이드 생성
# ══════════════════════════════════════════════════════════

def slide_01_title(prs):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT_BLUE)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "JTMS  vs  LTMS  vs  ATMS", 48, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.7),
             "세 가지 진리 유지 시스템의 구조와 원리 비교", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.5), Inches(3.6), Inches(4.3), Pt(2), DIM_GRAY)

    # 3개 시스템 뱃지
    badges = [
        ("JTMS", "Justification-based", JTMS_COLOR),
        ("LTMS", "Logic-based", LTMS_COLOR),
        ("ATMS", "Assumption-based", ATMS_COLOR),
    ]
    for i, (name, full, color) in enumerate(badges):
        x = Inches(1.5 + i * 3.8)
        add_rect(slide, x, Inches(4.2), Inches(3.0), Inches(1.2), CARD_BG, color)
        add_text(slide, x + Inches(0.1), Inches(4.3), Inches(2.8), Inches(0.5),
                 name, 24, color, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), Inches(4.8), Inches(2.8), Inches(0.4),
                 full, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Building Problem Solvers  |  Forbus & de Kleer", 16, DIM_GRAY, False, PP_ALIGN.CENTER)


def slide_02_toc(prs):
    """목차"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "", "목차")

    chapters = [
        ("01", "TMS란 무엇인가?", "세 시스템의 공통 기반"),
        ("02", "핵심 차이: 진리값 표현", "IN/OUT vs TRUE/FALSE/UNKNOWN vs 환경 집합"),
        ("03", "자료구조 비교", "Node, Justification, Clause, Environment"),
        ("04", "알고리즘 비교", "전파, 철회, 모순 처리 방식"),
        ("05", "예제: 같은 문제, 세 가지 풀이", "A→D, B→D를 각 시스템으로"),
        ("06", "가정 관리 비교", "설정, 철회, 비용 차이"),
        ("07", "추론 능력 스펙트럼", "단방향 → 양방향 → 다중 시나리오"),
        ("08", "장단점 및 적용 분야", "언제 어떤 TMS를 쓸 것인가"),
        ("09", "코드 구조 대응표", "파일 레벨 매핑"),
        ("10", "요약 및 핵심 정리", "한눈에 보는 비교표"),
    ]

    for i, (num, title, desc) in enumerate(chapters):
        row, col = i // 2, i % 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(1.5 + row * 1.0)
        add_rect(slide, x, y, Inches(5.5), Inches(0.85), CARD_BG, RGBColor(0x33, 0x44, 0x66))
        add_text(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.4),
                 num, 14, ACCENT_BLUE, True)
        add_text(slide, x + Inches(0.6), y + Inches(0.05), Inches(4.5), Inches(0.4),
                 title, 16, WHITE, True)
        add_text(slide, x + Inches(0.6), y + Inches(0.45), Inches(4.5), Inches(0.35),
                 desc, 12, DIM_GRAY)


def slide_03_what_is_tms(prs):
    """TMS란 무엇인가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "01", "TMS란 무엇인가?")

    # 정의 박스
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.2), CARD_BG, ACCENT_BLUE)
    add_multiline(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(1.0), [
        ("TMS = Truth Maintenance System (진리 유지 시스템)", 20, WHITE, True),
        ("", 6),
        ("추론 시스템에서 '무엇이 참인지'를 일관성 있게 관리하는 하위 시스템.", 16, LIGHT_GRAY),
        ("가정이 바뀌면 결론도 자동으로 업데이트하고, 모순이 발생하면 원인을 추적한다.", 16, LIGHT_GRAY),
    ])

    # 공통 구성요소
    add_text(slide, Inches(0.5), Inches(2.9), Inches(12), Inches(0.5),
             "세 시스템의 공통 구성 요소", 20, WHITE, True)

    components = [
        ("노드 (Node)", "명제/사실을 표현\n참/거짓 상태를 가짐", ACCENT_BLUE),
        ("정당화 (Justification)", "\"전제들이 참이면\n결론도 참\" 관계", ACCENT_GREEN),
        ("가정 (Assumption)", "참이라고 가정한 명제\n나중에 철회 가능", ACCENT_ORANGE),
        ("모순 처리", "일관성 없는 상태를\n감지하고 해결", ACCENT_RED),
    ]
    for i, (title, desc, color) in enumerate(components):
        x = Inches(0.5 + i * 3.15)
        add_rect(slide, x, Inches(3.5), Inches(2.9), Inches(1.7), CARD_BG, color)
        add_text(slide, x + Inches(0.1), Inches(3.6), Inches(2.7), Inches(0.4),
                 title, 16, color, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), Inches(4.15), Inches(2.7), Inches(0.9),
                 desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # 차이점 한줄
    add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3), CARD_BG, ACCENT_PURPLE)
    add_multiline(slide, Inches(0.8), Inches(5.7), Inches(11.8), Inches(1.1), [
        ("핵심 차이: '진리값을 어떻게 표현하고 관리하느냐'", 18, ACCENT_PURPLE, True),
        ("", 4),
        ("JTMS: 노드 상태 IN/OUT 하나  |  LTMS: 3값 논리 + CNF 절  |  ATMS: 환경 집합으로 모든 시나리오 동시 관리", 14, LIGHT_GRAY),
    ])


def slide_04_truth_representation(prs):
    """핵심 차이: 진리값 표현 방식"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "02", "핵심 차이 — 진리값 표현 방식")

    three_columns(slide, Inches(1.4), Inches(3.0), [
        # JTMS
        "label = :IN 또는 :OUT\n"
        "support = 하나의 just\n"
        "    또는 :ENABLED-ASSUMPTION\n"
        "\n"
        "▶ 2값 논리\n"
        "▶ 한 번에 하나의 시나리오\n"
        "▶ 다른 가정 시험하려면\n"
        "  철회 → 재설정 필요\n"
        "\n"
        "예: node-A.label = :IN\n"
        "    node-B.label = :OUT",

        # LTMS
        "label = :TRUE | :FALSE\n"
        "      | :UNKNOWN\n"
        "support = 하나의 clause\n"
        "\n"
        "▶ 3값 논리 (:UNKNOWN 추가)\n"
        "▶ 한 번에 하나의 시나리오\n"
        "▶ CNF 절 기반 양방향 추론\n"
        "▶ pvs 카운터로 효율적 전파\n"
        "\n"
        "예: node-x.label = :FALSE\n"
        "    node-y.label = :UNKNOWN",

        # ATMS
        "label = [env1, env2, ...]\n"
        "  각 env = {가정 노드 집합}\n"
        "\n"
        "▶ 환경(environment) 집합\n"
        "▶ 모든 시나리오를 동시 관리\n"
        "▶ 컨텍스트 전환 불필요\n"
        "▶ 최소성(minimality) 유지\n"
        "\n"
        "예: D.label = [{A}, {B}]\n"
        "  \"A만으로도, B만으로도\n"
        "   D가 참이 된다\"",
    ])

    # 비유
    add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.2), CARD_BG, DIM_GRAY)
    add_text(slide, Inches(0.7), Inches(4.8), Inches(11.8), Inches(0.4),
             "비유로 이해하기", 18, WHITE, True)

    analogies = [
        ("JTMS", "손전등 하나로\n방 하나씩 비추기", JTMS_COLOR),
        ("LTMS", "더 밝은 손전등으로\n앞뒤 모두 비추기", LTMS_COLOR),
        ("ATMS", "모든 방에 조명을\n동시에 켜놓기", ATMS_COLOR),
    ]
    for i, (name, desc, color) in enumerate(analogies):
        x = Inches(0.7 + i * 4.2)
        add_text(slide, x, Inches(5.3), Inches(3.8), Inches(0.4),
                 name, 16, color, True, PP_ALIGN.CENTER)
        add_text(slide, x, Inches(5.8), Inches(3.8), Inches(0.8),
                 desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)


def slide_05_data_structures(prs):
    """자료구조 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "03", "자료구조 비교")

    three_columns(slide, Inches(1.3), Inches(3.2), [
        # JTMS
        "tms-node:\n"
        "  index, datum\n"
        "  label: :IN | :OUT\n"
        "  support: just\n"
        "  justs: [가능한 정당화들]\n"
        "  consequences: [의존 just]\n"
        "  assumption?, contradictory?\n"
        "\n"
        "just (정당화):\n"
        "  antecedents → consequence\n"
        "  \"전제가 모두 :IN이면\n"
        "   결론도 :IN\"",

        # LTMS
        "tms-node:\n"
        "  index, datum\n"
        "  label: :TRUE|:FALSE|:UNKNOWN\n"
        "  support: clause\n"
        "  true-clauses / false-clauses\n"
        "  assumption?\n"
        "\n"
        "clause (CNF 절):\n"
        "  literals: [(n . :TRUE)...]\n"
        "  pvs: 미결정 리터럴 수\n"
        "    0 → 모순!  1 → 단위전파!\n"
        "  sats: 만족 리터럴 수",

        # ATMS
        "tms-node:\n"
        "  index, datum\n"
        "  label: [env1, env2, ...]\n"
        "    ← 최소 환경들의 집합\n"
        "  justs / consequences\n"
        "  assumption?\n"
        "\n"
        "env (환경):\n"
        "  assumptions: [A1, A3, ...]\n"
        "  nogood?: nil | just\n"
        "  nodes: [이 env에서 참인 것들]\n"
        "  ← 구조적 공유로 메모리 절약",
    ])

    # 핵심 차이 강조
    add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0), CARD_BG, ACCENT_PURPLE)
    add_multiline(slide, Inches(0.8), Inches(4.9), Inches(11.8), Inches(1.8), [
        ("핵심 자료구조 차이", 18, ACCENT_PURPLE, True),
        ("", 4),
        ("JTMS:  justification(정당화) = 전제 → 결론 관계. 순방향 그래프.", 14, JTMS_COLOR, True),
        ("LTMS:  clause(절) = CNF 리터럴의 OR. pvs 카운터가 핵심 최적화.", 14, LTMS_COLOR, True),
        ("ATMS:  environment(환경) = 가정들의 조합. label이 환경 집합.", 14, ATMS_COLOR, True),
        ("", 4),
        ("JTMS/ATMS는 justification 구조가 동일. LTMS만 clause(절) 기반으로 다름.", 14, DIM_GRAY),
    ])


def slide_06_algorithms(prs):
    """알고리즘 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "04", "알고리즘 비교 — 전파 · 철회 · 모순 처리")

    three_columns(slide, Inches(1.3), Inches(2.6), [
        # JTMS
        [
            ("전파: propagate-inness", 12, JTMS_COLOR, True),
            ("순방향만 (IN→IN)", 11, LIGHT_GRAY),
            ("전제가 모두 :IN → 결론 :IN", 11, LIGHT_GRAY),
            ("", 4),
            ("철회: propagate-outness", 12, ACCENT_RED, True),
            ("→ 의존 노드 :OUT 전파", 11, LIGHT_GRAY),
            ("→ find-alternative-support", 11, LIGHT_GRAY),
            ("  (다른 정당화 탐색)", 11, DIM_GRAY),
            ("", 4),
            ("모순: contradiction 노드 :IN 시", 12, ACCENT_ORANGE, True),
            ("→ 핸들러 호출", 11, LIGHT_GRAY),
        ],
        # LTMS
        [
            ("전파: BCP (check-clauses)", 12, LTMS_COLOR, True),
            ("양방향 (TRUE↔FALSE)", 11, LIGHT_GRAY),
            ("pvs=1 → 단위 전파", 11, LIGHT_GRAY),
            ("", 4),
            ("철회: propagate-unknownness", 12, ACCENT_RED, True),
            ("→ 의존 노드 :UNKNOWN 전파", 11, LIGHT_GRAY),
            ("→ 재전파로 새 안정 상태", 11, LIGHT_GRAY),
            ("  (다른 근거 있으면 유지)", 11, DIM_GRAY),
            ("", 4),
            ("모순: pvs = 0 감지", 12, ACCENT_ORANGE, True),
            ("→ 핸들러 / nogood 학습", 11, LIGHT_GRAY),
        ],
        # ATMS
        [
            ("전파: weave + update-label", 12, ATMS_COLOR, True),
            ("환경 교차곱 계산", 11, LIGHT_GRAY),
            ("최소성 필터링 유지", 11, LIGHT_GRAY),
            ("", 4),
            ("철회: 필요 없음!", 12, ACCENT_GREEN, True),
            ("→ 환경을 nogood 표시만", 11, LIGHT_GRAY),
            ("→ 라벨에서 자동 제거", 11, LIGHT_GRAY),
            ("  (재계산 비용 O(1))", 11, DIM_GRAY),
            ("", 4),
            ("모순: env.nogood 설정", 12, ACCENT_ORANGE, True),
            ("→ 상위집합도 자동 nogood", 11, LIGHT_GRAY),
        ],
    ])

    # 전파 흐름 비교
    add_text(slide, Inches(0.5), Inches(4.1), Inches(12), Inches(0.4),
             "전파 흐름 비교", 18, WHITE, True)

    add_code(slide, Inches(0.5), Inches(4.6), Inches(4.0), Inches(2.5),
             "JTMS:\n"
             "enable(A)\n"
             "  → A = :IN\n"
             "  → just(A→D) 검사\n"
             "    A가 :IN? → D = :IN\n"
             "  → just(D→G) 검사\n"
             "    D가 :IN? → G = :IN\n"
             "  → 순방향 연쇄만 가능", 11)

    add_code(slide, Inches(4.7), Inches(4.6), Inches(4.0), Inches(2.5),
             "LTMS:\n"
             "enable(x, :FALSE)\n"
             "  → 절(x∨y)의 pvs: 2→1\n"
             "  → 단위전파: y = :TRUE\n"
             "  → 절(¬y∨z)의 pvs: 2→1\n"
             "  → 단위전파: z = :TRUE\n"
             "  → 양방향(TRUE/FALSE) 전파\n"
             "  → Resolution으로 완전 추론", 11)

    add_code(slide, Inches(8.9), Inches(4.6), Inches(4.0), Inches(2.5),
             "ATMS:\n"
             "justify(J1, D, [A])\n"
             "  → weave: [{A}]\n"
             "  → D.label = [{A}]\n"
             "justify(J2, D, [B])\n"
             "  → weave: [{B}]\n"
             "  → D.label = [{A},{B}]\n"
             "  → 모든 경로 동시 기록!", 11)


def slide_07_same_problem(prs):
    """같은 문제, 세 가지 풀이"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "05", "예제: 같은 문제를 세 시스템으로 풀기")

    # 문제 설정
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.9), CARD_BG, ACCENT_PURPLE)
    add_multiline(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(0.8), [
        ("문제:  \"A이면 D이다. B이면 D이다. A와 B 중 어느 것이 참인지 확인하라.\"", 17, WHITE, True),
        ("규칙:  A → D,  B → D    |    가정 후보: A, B", 15, LIGHT_GRAY),
    ])

    three_columns(slide, Inches(2.5), Inches(4.5), [
        # JTMS
        "시도 1: enable(A)\n"
        "  → A = :IN\n"
        "  → just(A→D) 만족\n"
        "  → D = :IN ✓\n"
        "  → 탐색 계속...\n"
        "\n"
        "시도 2: retract(A)\n"
        "  → A = :OUT\n"
        "  → D = :OUT (지지 상실)\n"
        "  → propagate-outness\n"
        "  → enable(B)\n"
        "  → B = :IN → D = :IN ✓\n"
        "\n"
        "※ 시나리오마다 철회+재전파\n"
        "  비용이 발생함",

        # LTMS
        "절 추가:\n"
        "  ¬A ∨ D  (A이면 D)\n"
        "  ¬B ∨ D  (B이면 D)\n"
        "\n"
        "시도 1: enable(A, :TRUE)\n"
        "  → ¬A = FALSE\n"
        "  → 절(¬A∨D) pvs=1\n"
        "  → 단위전파: D = :TRUE ✓\n"
        "\n"
        "시도 2: retract(A)\n"
        "  → D = :UNKNOWN\n"
        "  → enable(B, :TRUE)\n"
        "  → 재전파: D = :TRUE ✓\n"
        "\n"
        "※ 양방향이라 추론력 더 강함",

        # ATMS
        "assume(A)  A.label=[{A}]\n"
        "assume(B)  B.label=[{B}]\n"
        "\n"
        "justify(J1, D, [A])\n"
        "  → D.label = [{A}]\n"
        "\n"
        "justify(J2, D, [B])\n"
        "  → D.label = [{A}, {B}]\n"
        "\n"
        "끝! 한 번에 완료.\n"
        "\n"
        "\"D는 A만 가정해도 참,\n"
        " B만 가정해도 참\"\n"
        "\n"
        "※ 철회/재계산 필요 없음!",
    ])


def slide_08_assumption_mgmt(prs):
    """가정 관리 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "06", "가정 관리 비교")

    # 테이블 형태
    headers = ["", "JTMS", "LTMS", "ATMS"]
    rows = [
        ("가정 설정", "enable-assumption\n→ 노드 :IN으로", "enable-assumption\n(node, :TRUE/:FALSE)", "노드 생성 시 자동\n(assumptionp t)"),
        ("가정 철회", "retract-assumption\n→ OUT 전파\n→ 대안 지지 탐색", "retract-assumption\n→ UNKNOWN 전파\n→ 재전파", "철회 개념 없음\nnogood만 추가\n→ 라벨에서 제거"),
        ("철회 비용", "비쌈\n(재계산 필요)", "비쌈\n(재계산 필요)", "O(1)\n(환경 표시만)"),
        ("동시 시나리오", "1개", "1개", "모두 동시"),
    ]

    colors_h = [DIM_GRAY, JTMS_COLOR, LTMS_COLOR, ATMS_COLOR]

    # 헤더
    for j, h in enumerate(headers):
        x = Inches(0.5 + j * 3.15)
        add_rect(slide, x, Inches(1.4), Inches(2.9), Inches(0.45),
                 CARD_BG if j == 0 else CARD_BG, colors_h[j])
        add_text(slide, x + Inches(0.1), Inches(1.43), Inches(2.7), Inches(0.4),
                 h, 16, colors_h[j], True, PP_ALIGN.CENTER)

    # 행
    for i, (label, *cells) in enumerate(rows):
        y = Inches(2.0 + i * 1.3)
        # 행 라벨
        add_rect(slide, Inches(0.5), y, Inches(2.9), Inches(1.15), CARD_BG, DIM_GRAY)
        add_text(slide, Inches(0.6), y + Inches(0.1), Inches(2.7), Inches(0.9),
                 label, 14, WHITE, True, PP_ALIGN.CENTER)
        # 셀
        for j, cell in enumerate(cells):
            x = Inches(3.65 + j * 3.15)
            color = colors_h[j + 1]
            add_rect(slide, x, y, Inches(2.9), Inches(1.15), CARD_BG,
                     RGBColor(0x33, 0x44, 0x66))
            add_text(slide, x + Inches(0.1), y + Inches(0.08), Inches(2.7), Inches(1.0),
                     cell, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)


def slide_09_inference_spectrum(prs):
    """추론 능력 스펙트럼"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "07", "추론 능력 스펙트럼")

    # 스펙트럼 바
    add_rect(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5),
             RGBColor(0x15, 0x1D, 0x35), DIM_GRAY)
    # 그라디언트 효과 - 3 구간
    add_rect(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.5), JTMS_COLOR)
    add_rect(slide, Inches(4.3), Inches(1.7), Inches(3.8), Inches(0.5), LTMS_COLOR)
    add_rect(slide, Inches(8.1), Inches(1.7), Inches(4.2), Inches(0.5), ATMS_COLOR)

    add_text(slide, Inches(0.8), Inches(1.73), Inches(3.5), Inches(0.45),
             "JTMS", 16, BG_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(4.3), Inches(1.73), Inches(3.8), Inches(0.45),
             "LTMS", 16, BG_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(8.1), Inches(1.73), Inches(4.2), Inches(0.45),
             "ATMS", 16, BG_DARK, True, PP_ALIGN.CENTER)

    add_text(slide, Inches(0.8), Inches(2.3), Inches(3), Inches(0.3),
             "추론 능력 약함", 12, DIM_GRAY)
    add_text(slide, Inches(9.3), Inches(2.3), Inches(3), Inches(0.3),
             "추론 능력 강함", 12, DIM_GRAY, False, PP_ALIGN.RIGHT)

    # 각 시스템의 능력 상세
    specs = [
        ("JTMS", JTMS_COLOR, [
            ("전파 방향", "순방향만 (A→B→C)"),
            ("추론 모드", "단일 시나리오, 순차 탐색"),
            ("강점", "단순, 빠름, 메모리 적음"),
            ("한계", "¬B→¬A 역추론 불가"),
            ("시나리오 수", "한 번에 1개"),
        ]),
        ("LTMS", LTMS_COLOR, [
            ("전파 방향", "양방향 (A→B, ¬B→¬A)"),
            ("추론 모드", "단일 시나리오, BCP + 해소"),
            ("강점", "단위전파, 완전 추론 가능"),
            ("한계", "여전히 단일 시나리오"),
            ("시나리오 수", "한 번에 1개"),
        ]),
        ("ATMS", ATMS_COLOR, [
            ("전파 방향", "환경 교차곱 (weave)"),
            ("추론 모드", "다중 시나리오 동시"),
            ("강점", "철회 비용 0, 모든 경로"),
            ("한계", "환경 수 폭발 가능 (2^n)"),
            ("시나리오 수", "모두 동시"),
        ]),
    ]

    for i, (name, color, items) in enumerate(specs):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(2.9), Inches(3.9), Inches(3.8), CARD_BG, color)
        add_text(slide, x + Inches(0.1), Inches(3.0), Inches(3.7), Inches(0.4),
                 name, 18, color, True, PP_ALIGN.CENTER)
        for j, (key, val) in enumerate(items):
            y = Inches(3.5 + j * 0.6)
            add_text(slide, x + Inches(0.15), y, Inches(1.5), Inches(0.3),
                     key, 12, DIM_GRAY, True)
            add_text(slide, x + Inches(1.6), y, Inches(2.1), Inches(0.5),
                     val, 12, LIGHT_GRAY)


def slide_10_pros_cons(prs):
    """장단점 및 적용 분야"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "08", "장단점 및 적용 분야 — 언제 어떤 TMS를 쓸 것인가")

    three_columns(slide, Inches(1.3), Inches(2.5), [
        # JTMS
        [
            ("장점", 14, ACCENT_GREEN, True),
            ("• 구조가 가장 단순", 12, LIGHT_GRAY),
            ("• 메모리 효율적", 12, LIGHT_GRAY),
            ("• 이해하기 쉬움", 12, LIGHT_GRAY),
            ("", 6),
            ("단점", 14, ACCENT_RED, True),
            ("• 순방향 추론만 가능", 12, LIGHT_GRAY),
            ("• 시나리오 전환 비용 큼", 12, LIGHT_GRAY),
            ("• 역추론 불가", 12, LIGHT_GRAY),
        ],
        # LTMS
        [
            ("장점", 14, ACCENT_GREEN, True),
            ("• 양방향 추론 (BCP)", 12, LIGHT_GRAY),
            ("• 완전 추론 가능 (Resolution)", 12, LIGHT_GRAY),
            ("• SAT 솔버의 원리", 12, LIGHT_GRAY),
            ("", 6),
            ("단점", 14, ACCENT_RED, True),
            ("• 여전히 단일 시나리오", 12, LIGHT_GRAY),
            ("• CNF 변환 오버헤드", 12, LIGHT_GRAY),
            ("• 철회 시 재계산 필요", 12, LIGHT_GRAY),
        ],
        # ATMS
        [
            ("장점", 14, ACCENT_GREEN, True),
            ("• 모든 시나리오 동시 관리", 12, LIGHT_GRAY),
            ("• 철회 비용 O(1)", 12, LIGHT_GRAY),
            ("• 진단/설계에 최적", 12, LIGHT_GRAY),
            ("", 6),
            ("단점", 14, ACCENT_RED, True),
            ("• 환경 수 폭발 가능 (2^n)", 12, LIGHT_GRAY),
            ("• 메모리 집약적", 12, LIGHT_GRAY),
            ("• 구현 복잡도 높음", 12, LIGHT_GRAY),
        ],
    ])

    # 적용 분야
    add_text(slide, Inches(0.5), Inches(4.1), Inches(12), Inches(0.4),
             "적합한 적용 분야", 18, WHITE, True)

    use_cases = [
        ("JTMS", JTMS_COLOR, "• 단순 규칙 시스템\n• 순방향 추론 엔진\n• 기호적 적분 (JSAINT)\n• N-Queens 퍼즐"),
        ("LTMS", LTMS_COLOR, "• 제약 만족 문제 (CSP)\n• SAT/SMT 문제\n• 양방향 추론이 필요한 경우\n• 마르크스 형제 퍼즐"),
        ("ATMS", ATMS_COLOR, "• 고장 진단 시스템\n• 설계 공간 탐색\n• 계획 수립 (블록 월드)\n• 다중 가설 추론"),
    ]
    for i, (name, color, desc) in enumerate(use_cases):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(4.6), Inches(3.9), Inches(2.3), CARD_BG, color)
        add_text(slide, x + Inches(0.1), Inches(4.7), Inches(3.7), Inches(0.35),
                 name, 16, color, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), Inches(5.1), Inches(3.6), Inches(1.7),
                 desc, 13, LIGHT_GRAY)


def slide_11_code_mapping(prs):
    """코드 구조 대응표"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "09", "코드 구조 대응표 — 파일 레벨 매핑")

    # 테이블 헤더
    cols = [
        ("기능", Inches(0.5), Inches(2.5), DIM_GRAY),
        ("JTMS (jtms/)", Inches(3.1), Inches(3.0), JTMS_COLOR),
        ("LTMS (ltms/)", Inches(6.2), Inches(3.0), LTMS_COLOR),
        ("ATMS (atms/)", Inches(9.3), Inches(3.5), ATMS_COLOR),
    ]
    for label, x, w, color in cols:
        add_rect(slide, x, Inches(1.3), w, Inches(0.45), CARD_BG, color)
        add_text(slide, x + Inches(0.05), Inches(1.33), w - Inches(0.1), Inches(0.4),
                 label, 13, color, True, PP_ALIGN.CENTER)

    file_rows = [
        ("핵심 엔진", "jtms.lisp (637줄)", "ltms.lisp (746줄)", "atms.lisp (1079줄)"),
        ("DB 인터페이스", "jdata.lisp (260줄)", "ldata.lisp (305줄)", "adata.lisp (233줄)"),
        ("규칙 엔진", "jrules.lisp (250줄)", "lrules.lisp (286줄)", "arules.lisp (372줄)"),
        ("통합 계층", "jinter.lisp (78줄)", "linter.lisp (101줄)", "ainter.lisp (173줄)"),
        ("패턴 매칭", "funify.lisp (130줄)", "funify.lisp (131줄)", "funify.lisp (129줄)"),
        ("로딩", "jtre.lisp (54줄)", "ltre.lisp (51줄)", "atre.lisp (51줄)"),
        ("고유 모듈", "JSAINT (적분 솔버)\njqueens (N-Queens)", "cltms (완전 LTMS)\ncwa, dds, indirect", "plan-a, plan-e (계획)\nblocks (블록 월드)"),
        ("총 규모", "~2,588줄 / 15파일", "~3,552줄 / 15파일", "~3,008줄 / 14파일"),
    ]

    for i, (func, j_file, l_file, a_file) in enumerate(file_rows):
        y = Inches(1.85 + i * 0.65)
        h = Inches(0.58)
        bg = CARD_BG if i % 2 == 0 else RGBColor(0x1E, 0x24, 0x40)

        add_rect(slide, Inches(0.5), y, Inches(2.5), h, bg, RGBColor(0x33, 0x44, 0x66))
        add_text(slide, Inches(0.6), y + Inches(0.03), Inches(2.3), h,
                 func, 11, WHITE, True, PP_ALIGN.LEFT)

        for j, val in enumerate([j_file, l_file, a_file]):
            x = [Inches(3.1), Inches(6.2), Inches(9.3)][j]
            w = [Inches(3.0), Inches(3.0), Inches(3.5)][j]
            add_rect(slide, x, y, w, h, bg, RGBColor(0x33, 0x44, 0x66))
            add_text(slide, x + Inches(0.05), y + Inches(0.03), w - Inches(0.1), h,
                     val, 10, LIGHT_GRAY, False, PP_ALIGN.CENTER, "Consolas")


def slide_12_summary(prs):
    """요약 및 핵심 정리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_header(slide, "10", "요약 — 한눈에 보는 비교")

    # 핵심 비교 3열
    summaries = [
        ("JTMS", JTMS_COLOR,
         "Justification-based",
         "\"지금 이 가정 하에서\n무엇이 참인가?\"",
         "IN / OUT\n(2값, 단일 시나리오)",
         "순방향 전파\n(propagate-inness)",
         "단순 규칙 시스템\nJSAINT, N-Queens"),
        ("LTMS", LTMS_COLOR,
         "Logic-based",
         "\"이 제약들 하에서 무엇이\n반드시 참/거짓인가?\"",
         "TRUE / FALSE / UNKNOWN\n(3값, CNF 절 기반)",
         "BCP + 단위 전파\n+ Resolution",
         "제약 만족 (CSP)\nSAT 솔버의 원형"),
        ("ATMS", ATMS_COLOR,
         "Assumption-based",
         "\"모든 가능한 가정 조합에서\n무엇이 참인가?\"",
         "환경 집합 [{A},{B,C},...]\n(다중 시나리오 동시)",
         "weave (환경 교차곱)\n+ nogood 전파",
         "고장 진단\n설계 탐색, 계획 수립"),
    ]

    labels = ["정식 명칭", "핵심 질문", "진리값 표현", "핵심 알고리즘", "적합한 분야"]

    for i, (name, color, fullname, question, truth, algo, domain) in enumerate(summaries):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(1.3), Inches(3.9), Inches(5.5), CARD_BG, color)
        add_text(slide, x + Inches(0.1), Inches(1.35), Inches(3.7), Inches(0.5),
                 name, 24, color, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), Inches(1.8), Inches(3.7), Inches(0.3),
                 fullname, 12, DIM_GRAY, False, PP_ALIGN.CENTER)

        items = [question, truth, algo, domain]
        for j, (label, item) in enumerate(zip(labels[1:], items)):
            y = Inches(2.25 + j * 1.2)
            add_text(slide, x + Inches(0.15), y, Inches(3.5), Inches(0.3),
                     label, 11, color, True)
            add_text(slide, x + Inches(0.15), y + Inches(0.3), Inches(3.5), Inches(0.8),
                     item, 12, LIGHT_GRAY)

    # 하단 한줄 요약
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), CARD_BG, ACCENT_PURPLE)
    add_text(slide, Inches(0.7), Inches(7.02), Inches(11.8), Inches(0.3),
             "JTMS는 가장 단순  |  LTMS는 제약 전파가 가장 강력  |  ATMS는 다중 시나리오 관리에 가장 효율적",
             13, WHITE, True, PP_ALIGN.CENTER)


def slide_13_closing(prs):
    """마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), ACCENT_BLUE)

    add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             "Thank You", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.5), Inches(3.0), Inches(4.3), Pt(2), DIM_GRAY)

    add_text(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.8),
             "문제의 성격에 맞는 TMS를 선택하는 것이 핵심입니다.", 20, WHITE, True, PP_ALIGN.CENTER)

    # 3개 시스템 한줄 정리
    final_msgs = [
        ("JTMS", "단순한 규칙 추론에", JTMS_COLOR),
        ("LTMS", "제약 전파와 SAT에", LTMS_COLOR),
        ("ATMS", "다중 가설 탐색에", ATMS_COLOR),
    ]
    for i, (name, use, color) in enumerate(final_msgs):
        x = Inches(1.5 + i * 3.8)
        add_rect(slide, x, Inches(4.5), Inches(3.0), Inches(1.0), CARD_BG, color)
        add_text(slide, x + Inches(0.1), Inches(4.55), Inches(2.8), Inches(0.4),
                 name, 20, color, True, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), Inches(5.0), Inches(2.8), Inches(0.4),
                 use, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "Building Problem Solvers  |  Forbus & de Kleer", 16, DIM_GRAY, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  메인 실행
# ══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_toc(prs)
    slide_03_what_is_tms(prs)
    slide_04_truth_representation(prs)
    slide_05_data_structures(prs)
    slide_06_algorithms(prs)
    slide_07_same_problem(prs)
    slide_08_assumption_mgmt(prs)
    slide_09_inference_spectrum(prs)
    slide_10_pros_cons(prs)
    slide_11_code_mapping(prs)
    slide_12_summary(prs)
    slide_13_closing(prs)

    output_path = "/home/user/bps/ltms/TMS_비교_JTMS_LTMS_ATMS.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")
    print(f"총 {len(prs.slides)}장의 슬라이드")


if __name__ == "__main__":
    main()
