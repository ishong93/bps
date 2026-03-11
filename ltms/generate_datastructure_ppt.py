#!/usr/bin/env python3
"""JTMS vs LTMS vs ATMS 핵심 자료구조 상세 비교 교육용 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK     = RGBColor(0x11, 0x14, 0x27)
BG_CARD     = RGBColor(0x1A, 0x1F, 0x3A)
BG_CARD_ALT = RGBColor(0x15, 0x19, 0x30)
BG_CODE     = RGBColor(0x0C, 0x0F, 0x1C)

BLUE        = RGBColor(0x00, 0xD2, 0xFF)
GREEN       = RGBColor(0x00, 0xE6, 0x96)
ORANGE      = RGBColor(0xFF, 0x9F, 0x43)
RED         = RGBColor(0xFF, 0x6B, 0x6B)
PURPLE      = RGBColor(0xA5, 0x5E, 0xFF)
YELLOW      = RGBColor(0xFF, 0xE0, 0x66)
CYAN        = RGBColor(0x56, 0xE8, 0xE0)
PINK        = RGBColor(0xFF, 0x79, 0xC6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY       = RGBColor(0xBB, 0xBB, 0xCC)
DGRAY       = RGBColor(0x77, 0x77, 0x99)
BORDER      = RGBColor(0x2A, 0x30, 0x55)

JTMS_C = BLUE
LTMS_C = GREEN
ATMS_C = ORANGE

SW = Inches(13.333)
SH = Inches(7.5)


# ═══════════════════════════════════════════════════════════
#  헬퍼 함수
# ═══════════════════════════════════════════════════════════

def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=BG_CARD, border=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def multi(slide, l, t, w, h, lines, dsz=14, dclr=LGRAY):
    """lines: str 또는 (text, size, color, bold, align) 튜플 리스트"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, sz, clr, bold, align = item, dsz, dclr, False, PP_ALIGN.LEFT
        else:
            text  = item[0]
            sz    = item[1] if len(item) > 1 else dsz
            clr   = item[2] if len(item) > 2 else dclr
            bold  = item[3] if len(item) > 3 else False
            align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = "맑은 고딕"
        p.alignment = align
        p.space_after = Pt(2)
    return box


def code(slide, l, t, w, h, text, sz=11, clr=GREEN):
    s = rect(slide, l, t, w, h, BG_CODE, BORDER)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_right = Pt(10)
    tf.margin_bottom = Pt(6)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(1)
    return s


def code_multi(slide, l, t, w, h, lines, sz=11):
    """lines: (text, color) 튜플 리스트"""
    s = rect(slide, l, t, w, h, BG_CODE, BORDER)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_right = Pt(10)
    tf.margin_bottom = Pt(6)
    for i, (text, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(1)
    return s


def header(slide, num, title, subtitle=None):
    txt(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.6),
        f"{num}  {title}", 30, BLUE, True)
    rect(slide, Inches(0.7), Inches(0.88), Inches(2.5), Pt(3), BLUE)
    if subtitle:
        txt(slide, Inches(0.7), Inches(0.95), Inches(11), Inches(0.4),
            subtitle, 14, DGRAY)


def badge(slide, l, t, w, h, label, color):
    rect(slide, l, t, w, h, BG_CARD, color)
    txt(slide, l, t, w, h, label, 16, color, True, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def col3_headers(slide, y, h=Inches(0.42)):
    for i, (name, clr) in enumerate([("JTMS", JTMS_C), ("LTMS", LTMS_C), ("ATMS", ATMS_C)]):
        x = Inches(3.3 + i * 3.4)
        rect(slide, x, y, Inches(3.2), h, BG_CARD, clr)
        txt(slide, x, y, Inches(3.2), h, name, 15, clr, True, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def table_row(slide, y, label, cells, h=Inches(0.55), label_w=Inches(2.8)):
    """label + 3 cells (JTMS/LTMS/ATMS)"""
    rect(slide, Inches(0.3), y, label_w, h, BG_CARD_ALT, BORDER)
    txt(slide, Inches(0.4), y + Inches(0.02), label_w - Inches(0.2), h,
        label, 12, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
    colors = [JTMS_C, LTMS_C, ATMS_C]
    for i, cell in enumerate(cells):
        x = Inches(3.3 + i * 3.4)
        rect(slide, x, y, Inches(3.2), h, BG_CARD_ALT, BORDER)
        txt(slide, x + Inches(0.1), y + Inches(0.02), Inches(3.0), h,
            cell, 11, LGRAY, False, PP_ALIGN.CENTER, "Consolas", MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════
#  슬라이드 생성
# ═══════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    rect(s, Inches(0), Inches(0), SW, Pt(4), BLUE)
    rect(s, Inches(0), SH - Pt(4), SW, Pt(4), BLUE)

    txt(s, Inches(1), Inches(1.2), Inches(11.3), Inches(0.9),
        "JTMS  vs  LTMS  vs  ATMS", 46, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(2.2), Inches(11.3), Inches(0.6),
        "핵심 자료구조 상세 비교", 28, LGRAY, False, PP_ALIGN.CENTER)

    rect(s, Inches(5), Inches(3.2), Inches(3.3), Pt(2), DGRAY)

    txt(s, Inches(1), Inches(3.6), Inches(11.3), Inches(0.5),
        "defstruct 레벨에서 살펴보는 세 시스템의 설계 철학", 18, DGRAY, False, PP_ALIGN.CENTER)

    # 3 badges
    items = [
        ("JTMS", "3 structs\njtms · tms-node · just", JTMS_C),
        ("LTMS", "3 structs\nltms · tms-node · clause", LTMS_C),
        ("ATMS", "4 structs\natms · tms-node · just · env", ATMS_C),
    ]
    for i, (name, desc, clr) in enumerate(items):
        x = Inches(0.8 + i * 4.2)
        rect(s, x, Inches(4.5), Inches(3.8), Inches(1.6), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(4.6), Inches(3.6), Inches(0.5),
            name, 24, clr, True, PP_ALIGN.CENTER)
        txt(s, x + Inches(0.1), Inches(5.2), Inches(3.6), Inches(0.7),
            desc, 14, LGRAY, False, PP_ALIGN.CENTER, "Consolas")

    txt(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
        "Building Problem Solvers  |  Forbus & de Kleer", 14, DGRAY, False, PP_ALIGN.CENTER)


def slide_02_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "", "목차")

    chapters = [
        ("01", "최상위 시스템 구조체 비교", "jtms / ltms / atms 필드 대조"),
        ("02", "노드 구조체 — 가장 핵심적인 차이", "tms-node의 label·support 필드"),
        ("03", "노드 필드 전체 대조표", "13개 필드 일대일 매핑"),
        ("04", "추론 단위: just vs clause vs env", "설계 철학이 갈리는 지점"),
        ("05", "JTMS just 구조 상세", "순방향 전제→결론 그래프"),
        ("06", "LTMS clause 구조 상세", "CNF 절 + pvs/sats 카운터"),
        ("07", "ATMS just + env 구조 상세", "다중 세계관 동시 관리"),
        ("08", "연결 관계 다이어그램", "구조체 간 참조 관계"),
        ("09", "pvs 카운터 작동 원리", "LTMS BCP의 핵심 메커니즘"),
        ("10", "env 라벨 작동 원리", "ATMS 환경 집합의 최소성"),
        ("11", "설계 철학 요약 비교", "한눈에 보는 종합 비교표"),
    ]
    for i, (num, title, desc) in enumerate(chapters):
        row, col = i // 2, i % 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.3 + row * 0.95)
        rect(s, x, y, Inches(5.8), Inches(0.82), BG_CARD, BORDER)
        txt(s, x + Inches(0.1), y + Inches(0.06), Inches(0.5), Inches(0.35),
            num, 13, BLUE, True)
        txt(s, x + Inches(0.55), y + Inches(0.03), Inches(5.0), Inches(0.38),
            title, 15, WHITE, True)
        txt(s, x + Inches(0.55), y + Inches(0.42), Inches(5.0), Inches(0.3),
            desc, 11, DGRAY)


def slide_03_system_struct(prs):
    """최상위 시스템 구조체"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "01", "최상위 시스템 구조체 비교", "jtms / ltms / atms 의 defstruct 필드 대조")

    col3_headers(s, Inches(1.3))

    rows = [
        ("title",            "title",             "title",              "title"),
        ("노드 카운터",       "node-counter",      "node-counter",       "node-counter"),
        ("추론단위 카운터",   "just-counter (1개)", "clause-counter (1개)", "just-counter\n+ env-counter (2개)"),
        ("노드 저장",         "nodes (리스트)",     "nodes (해시 테이블)",  "nodes (리스트)"),
        ("추론단위 저장",     "justs (리스트)",     "clauses (리스트/trie)", "justs (리스트)"),
        ("가정 관리",         "assumptions ✓",     "— (없음)",           "assumptions ✓"),
        ("모순 관리",         "contradictions",    "pending-\ncontradictions", "contradictions\n+ contra-node"),
        ("모순 핸들러",       "contradiction-\nhandler (단수)",  "contradiction-\nhandlers (복수)",  "— (없음)\nnogood로 처리"),
        ("외부 연동",         "enqueue-procedure", "enqueue-procedure",  "enqueue-procedure"),
    ]

    y0 = Inches(1.82)
    rh = Inches(0.56)
    for i, (label, j, l, a) in enumerate(rows):
        y = y0 + i * (rh + Inches(0.04))
        table_row(s, y, label, [j, l, a], rh)

    # 고유 필드 강조
    txt(s, Inches(0.3), Inches(7.35) - Inches(1.1), Inches(12.7), Inches(0.35),
        "★ 시스템별 고유 필드", 16, PURPLE, True)

    extras = [
        ("LTMS 고유", "complete, violated-clauses, queue,\nconses, cons-size, delay-sat", LTMS_C),
        ("ATMS 고유", "env-table, nogood-table,\nempty-env, contra-node", ATMS_C),
    ]
    for i, (label, desc, clr) in enumerate(extras):
        x = Inches(0.3 + i * 6.5)
        rect(s, x, Inches(7.35) - Inches(0.7), Inches(6.2), Inches(0.65), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(7.35) - Inches(0.68), Inches(1.5), Inches(0.6),
            label, 12, clr, True, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(1.6), Inches(7.35) - Inches(0.68), Inches(4.4), Inches(0.6),
            desc, 11, LGRAY, False, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)


def slide_04_node_core_diff(prs):
    """노드 구조체 — 핵심 차이"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "02", "노드(tms-node) — 가장 핵심적인 차이", "label과 support 필드가 모든 것을 결정합니다")

    # label 비교
    txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
        "★ label 필드 — 진리값 표현", 20, YELLOW, True)

    labels_data = [
        ("JTMS", JTMS_C,
         ":IN  |  :OUT",
         "2값 논리\n\n\"지금 이 순간\n 믿는다 / 안 믿는다\"\n\n한 번에 하나의 세계만"),
        ("LTMS", LTMS_C,
         ":TRUE | :FALSE | :UNKNOWN",
         "3값 논리\n\n\"참이다 / 거짓이다 /\n 아직 모른다\"\n\n:UNKNOWN이 핵심 추가"),
        ("ATMS", ATMS_C,
         "[{A1}, {A2,A3}, ...]",
         "환경 집합\n\n\"어떤 가정 조합에서\n 참인지 모두 기록\"\n\n모든 세계를 동시 관리"),
    ]
    for i, (name, clr, val, desc) in enumerate(labels_data):
        x = Inches(0.5 + i * 4.2)
        rect(s, x, Inches(1.8), Inches(3.9), Inches(3.2), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(1.85), Inches(3.7), Inches(0.35),
            name, 18, clr, True, PP_ALIGN.CENTER)
        code(s, x + Inches(0.15), Inches(2.25), Inches(3.6), Inches(0.45),
             f"label = {val}", 12, clr)
        txt(s, x + Inches(0.15), Inches(2.8), Inches(3.6), Inches(2.0),
            desc, 14, LGRAY, False, PP_ALIGN.CENTER)

    # support 비교
    txt(s, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
        "★ support 필드 — \"왜 이 값인가?\"", 20, YELLOW, True)

    support_data = [
        ("JTMS", JTMS_C,
         "support → just 구조체\n       또는 :ENABLED-ASSUMPTION",
         "하나의 정당화를 가리킨다"),
        ("LTMS", LTMS_C,
         "support → clause 구조체\n       또는 :ENABLED-ASSUMPTION",
         "하나의 절을 가리킨다"),
        ("ATMS", ATMS_C,
         "support 필드 없음!\nlabel 자체가 모든 근거를 담고 있음",
         "label = [{A1}, {B2}] 그 자체가 증거"),
    ]
    for i, (name, clr, val, desc) in enumerate(support_data):
        x = Inches(0.5 + i * 4.2)
        rect(s, x, Inches(5.7), Inches(3.9), Inches(1.6), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(5.75), Inches(3.7), Inches(0.3),
            name, 14, clr, True, PP_ALIGN.CENTER)
        code(s, x + Inches(0.15), Inches(6.05), Inches(3.6), Inches(0.55),
             val, 10, LGRAY)
        txt(s, x + Inches(0.15), Inches(6.65), Inches(3.6), Inches(0.5),
            desc, 11, DGRAY, False, PP_ALIGN.CENTER)


def slide_05_node_full_table(prs):
    """노드 필드 전체 대조표"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "03", "노드(tms-node) 필드 전체 대조표")

    col3_headers(s, Inches(1.2))

    rows = [
        ("index",          "index",              "index",              "index"),
        ("datum",          "datum",              "datum",              "datum"),
        ("label",          ":IN | :OUT",         ":TRUE|:FALSE|:UNKNOWN", "[env1, env2, ...]"),
        ("support",        "→ just",             "→ clause",           "— (없음)"),
        ("정당화 참조",     "justs []",           "— (없음)",           "justs []"),
        ("의존 관계",       "consequences []",    "— (없음)",           "consequences []"),
        ("절 참조",         "— (없음)",           "true-clauses []\nfalse-clauses []", "— (없음)"),
        ("리터럴 캐시",     "— (없음)",           "true-literal\nfalse-literal", "— (없음)"),
        ("규칙 트리거",     "in-rules []\nout-rules []", "true-rules []\nfalse-rules []", "rules [] (하나)"),
        ("assumption?",    "✓",                  "✓",                  "✓"),
        ("contradictory?", "✓",                  "— (없음)",           "✓"),
        ("mark",           "✓ (sweep용)",        "✓ (sweep용)",        "— (없음)"),
        ("역참조",          "→ jtms",             "→ ltms",             "→ atms"),
    ]

    y0 = Inches(1.72)
    rh = Inches(0.4)
    for i, (label, j, l, a) in enumerate(rows):
        y = y0 + i * (rh + Inches(0.04))
        h = Inches(0.55) if "\n" in j or "\n" in l or "\n" in a else rh
        table_row(s, y, label, [j, l, a], h)
        # Shift subsequent rows if this row was taller
        if h > rh:
            y0 += (h - rh)


def slide_06_inference_units(prs):
    """추론 단위: just vs clause vs env"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "04", "추론 단위 — 설계 철학이 갈리는 지점",
           "JTMS·ATMS는 justification, LTMS는 clause, ATMS에만 environment 추가")

    items = [
        ("JTMS", JTMS_C, "just\n(정당화)",
         "전제가 모두 :IN이면\n→ 결론도 :IN\n\n순방향 규칙 그래프\nA ∧ B → C",
         "just\n├─ index\n├─ informant\n├─ consequence → node\n└─ antecedents → [node...]"),
        ("LTMS", LTMS_C, "clause\n(CNF 절)",
         "리터럴 중 하나 이상\n만족되어야 함 (OR)\n\n양방향 제약 전파\n(A ∨ ¬B ∨ C)",
         "clause\n├─ index\n├─ informant\n├─ literals → [(n.:T)...]\n├─ pvs   ← 미결정 수\n├─ sats  ← 만족 수\n├─ length\n└─ status"),
        ("ATMS", ATMS_C, "just + env\n(정당화+환경)",
         "just = JTMS와 동일\nenv = 가정들의 조합\n= \"하나의 세계관\"\n\n모든 세계 동시 추적",
         "just (= JTMS 동일)\n├─ antecedents\n└─ consequence\n\nenv (고유!)\n├─ index, count\n├─ assumptions → [node...]\n├─ nodes → [이 세계의 참]\n├─ nogood?\n└─ rules"),
    ]

    for i, (name, clr, unit, desc, struct) in enumerate(items):
        x = Inches(0.4 + i * 4.25)
        w = Inches(4.0)

        rect(s, x, Inches(1.35), w, Inches(5.8), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(1.4), w - Inches(0.2), Inches(0.35),
            name, 20, clr, True, PP_ALIGN.CENTER)

        # 추론 단위 이름
        badge(s, x + Inches(0.6), Inches(1.85), w - Inches(1.2), Inches(0.55), unit, clr)

        # 설명
        txt(s, x + Inches(0.15), Inches(2.5), w - Inches(0.3), Inches(1.7),
            desc, 13, LGRAY, False, PP_ALIGN.CENTER)

        # 구조 코드
        code(s, x + Inches(0.15), Inches(4.3), w - Inches(0.3), Inches(2.7),
             struct, 11, clr)


def slide_07_jtms_just_detail(prs):
    """JTMS just 상세"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "05", "JTMS just(정당화) 구조 상세", "순방향 전제→결론 그래프")

    # 코드
    code(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.2),
         "(defstruct (just (:PRINT-FUNCTION print-just))\n"
         "  (index 0)        ; 고유 번호\n"
         "  informant         ; 규칙 이름 (예: \"modus-ponens\")\n"
         "  consequence       ; 결론 노드 (tms-node 1개)\n"
         "  antecedents)      ; 전제 노드들 (tms-node 리스트)", 13, JTMS_C)

    # 의미 설명
    rect(s, Inches(6.3), Inches(1.3), Inches(6.5), Inches(2.2), BG_CARD, JTMS_C)
    multi(s, Inches(6.5), Inches(1.4), Inches(6.1), Inches(2.0), [
        ("의미", 16, JTMS_C, True),
        ("", 4),
        ("\"전제(antecedents)가 모두 :IN이면\"", 14, WHITE, True),
        ("\"→ 결론(consequence)도 :IN으로 만들어라\"", 14, WHITE, True),
        ("", 6),
        ("• informant: 이 정당화를 만든 규칙 이름", 13, LGRAY),
        ("• 순방향만 가능: A → B 있어도 ¬B → ¬A 불가", 13, LGRAY),
    ])

    # 예제 다이어그램
    txt(s, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
        "예제: \"A이면 D이다\" + \"D이면 G이다\"", 18, WHITE, True)

    code(s, Inches(0.5), Inches(4.3), Inches(6.0), Inches(2.8),
         ";; 정당화 생성\n"
         "(justify-node 'R1 node-D (list node-A))\n"
         "  → just-1: antecedents=[A], consequence=D\n"
         "\n"
         "(justify-node 'R2 node-G (list node-D))\n"
         "  → just-2: antecedents=[D], consequence=G\n"
         "\n"
         ";; A를 가정으로 활성화\n"
         "(enable-assumption node-A)\n"
         "  → A=:IN → just-1 만족 → D=:IN\n"
         "           → just-2 만족 → G=:IN", 12, JTMS_C)

    # 그래프
    rect(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8), BG_CARD, JTMS_C)
    txt(s, Inches(7.0), Inches(4.4), Inches(5.4), Inches(0.3),
        "의존 그래프", 16, JTMS_C, True, PP_ALIGN.CENTER)

    # 노드 그림
    code(s, Inches(7.2), Inches(4.9), Inches(5.2), Inches(2.0),
         "  node-A (:IN)          \n"
         "    │ assumption?=T     \n"
         "    │ support=:ENABLED-ASSUMPTION\n"
         "    ▼                   \n"
         "  ┌─just-1─┐           \n"
         "  │ R1      │           \n"
         "  └────┬────┘           \n"
         "       ▼                \n"
         "  node-D (:IN)  ──────► just-2 ──► node-G (:IN)\n"
         "    support=just-1        R2        support=just-2", 10, CYAN)


def slide_08_ltms_clause_detail(prs):
    """LTMS clause 상세"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "06", "LTMS clause(절) 구조 상세", "CNF 절 + pvs/sats 카운터 — BCP의 핵심")

    # 코드
    code(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.5),
         "(defstruct (clause (:PRINT-FUNCTION print-clause))\n"
         "  (index 0)       ; 고유 번호\n"
         "  (informant nil)  ; 규칙 이름\n"
         "  (literals nil)   ; [(node . :TRUE), (node . :FALSE), ...]\n"
         "  (pvs 0)          ; ★ 미결정(potentially violating) 수\n"
         "  (length 0)       ; 전체 리터럴 수\n"
         "  (sats 0)         ; ★ 만족된 리터럴 수\n"
         "  (status nil))    ; :SUBSUMED | :QUEUED | :DIRTY | nil", 13, LTMS_C)

    # 의미 설명
    rect(s, Inches(6.6), Inches(1.3), Inches(6.2), Inches(2.5), BG_CARD, LTMS_C)
    multi(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.3), [
        ("핵심 의미", 16, LTMS_C, True),
        ("", 4),
        ("절 = 리터럴의 OR (CNF 형식)", 14, WHITE, True),
        ("\"리터럴 중 하나 이상 만족되어야 함\"", 14, WHITE),
        ("", 4),
        ("• literals: (노드 . 기대값) 쌍의 리스트", 13, LGRAY),
        ("  (A . :TRUE) = \"A가 참이면 이 절 만족\"", 12, DGRAY),
        ("  (B . :FALSE) = \"B가 거짓이면 이 절 만족\"", 12, DGRAY),
        ("", 4),
        ("JTMS의 just와 달리 양방향 추론 가능!", 13, YELLOW, True),
    ])

    # pvs / sats 설명
    txt(s, Inches(0.5), Inches(4.1), Inches(12), Inches(0.4),
        "★ pvs / sats 카운터 — LTMS 최적화의 핵심", 18, YELLOW, True)

    items = [
        ("pvs\n(potentially\n violating)", ORANGE,
         "\"아직 결정 안 된 리터럴 수\"\n\npvs = 0  →  모순 발생!\n    모든 리터럴이 위반됨\n\npvs = 1  →  단위 전파!\n    남은 하나가 반드시 참\n\npvs ≥ 2  →  아직 대기\n    결정할 수 없음"),
        ("sats\n(satisfied)", GREEN,
         "\"이미 만족된 리터럴 수\"\n\nsats > 0  →  절 만족됨\n    더 이상 전파 불필요\n\nsats = 0  →  미만족\n    pvs 카운터 감시 계속\n\n(satisfied-clause?)\n  = (> sats 0)"),
        ("상호 작용", PURPLE,
         "노드 값이 결정될 때마다:\n\n기대값과 일치 → sats++\n기대값과 불일치 → pvs--\n\n[pvs-- 후 검사]\n  pvs=0 → 모순!\n  pvs=1 → 단위전파!\n\n이 카운터 덕분에\nO(1)로 전파 결정 가능"),
    ]
    for i, (title, clr, desc) in enumerate(items):
        x = Inches(0.5 + i * 4.2)
        rect(s, x, Inches(4.6), Inches(3.9), Inches(2.7), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(4.65), Inches(3.7), Inches(0.6),
            title, 14, clr, True, PP_ALIGN.CENTER, "Consolas")
        txt(s, x + Inches(0.15), Inches(5.2), Inches(3.6), Inches(2.0),
            desc, 12, LGRAY)


def slide_09_atms_env_detail(prs):
    """ATMS just + env 상세"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "07", "ATMS just + env 구조 상세", "다중 세계관 동시 관리 — ATMS에만 존재하는 env")

    # just (왼쪽)
    txt(s, Inches(0.5), Inches(1.2), Inches(5), Inches(0.35),
        "just — JTMS와 동일한 구조", 16, ATMS_C, True)

    code(s, Inches(0.5), Inches(1.6), Inches(5.5), Inches(1.5),
         "(defstruct (just (:PRINT-FUNCTION print-just))\n"
         "  (index 0)        ; 고유 번호\n"
         "  (informant nil)   ; 규칙 이름\n"
         "  (consequence nil) ; 결론 노드 (1개)\n"
         "  (antecedents nil)); 전제 노드들 (리스트)", 13, ATMS_C)

    # env (오른쪽)
    txt(s, Inches(6.3), Inches(1.2), Inches(6.5), Inches(0.35),
        "★ env — ATMS 고유 구조: \"하나의 세계관\"", 16, ORANGE, True)

    code(s, Inches(6.3), Inches(1.6), Inches(6.5), Inches(1.5),
         "(defstruct (env (:PREDICATE env?))\n"
         "  (index 0)         ; E0, E1, E2...\n"
         "  (count 0)         ; 가정 수: {A1,A3} → 2\n"
         "  (assumptions nil) ; 가정 노드 리스트 (정렬됨)\n"
         "  (nodes nil)       ; 이 환경에서 참인 노드들\n"
         "  (nogood? nil)     ; nil=일관 | non-nil=모순\n"
         "  (rules nil))      ; nogood 시 실행할 규칙", 12, ORANGE)

    # env 핵심 설명
    txt(s, Inches(0.5), Inches(3.3), Inches(12), Inches(0.35),
        "★ env의 핵심 필드 해설", 18, YELLOW, True)

    fields = [
        ("assumptions", CYAN,
         "이 세계에서 \"참이라고 가정한\" 노드들\n\n{A1, A3} = \"A1과 A3을 참으로 가정\"\n\n• 인덱스 오름차순 정렬 (비교 최적화)\n• 구조적 공유: 같은 가정 조합 = 같은 env"),
        ("nodes", GREEN,
         "이 환경에서 참으로 도출된 모든 노드\n\n{A1,A3}에서 D가 참이면\n  → D ∈ env.nodes\n\n역방향 참조: 환경에서 노드 찾기 가능"),
        ("nogood?", RED,
         "nil = 이 세계는 일관적\nnon-nil = 이 세계는 모순!\n\n값이 just → 직접 모순 (모순의 근거)\n값이 env → 상위집합이 nogood\n\nnogood이 되면 모든 노드의\nlabel에서 이 환경이 제거됨"),
    ]
    for i, (name, clr, desc) in enumerate(fields):
        x = Inches(0.5 + i * 4.2)
        rect(s, x, Inches(3.75), Inches(3.9), Inches(3.3), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(3.8), Inches(3.7), Inches(0.35),
            name, 16, clr, True, PP_ALIGN.CENTER, "Consolas")
        txt(s, x + Inches(0.15), Inches(4.2), Inches(3.6), Inches(2.7),
            desc, 12, LGRAY)


def slide_10_connections(prs):
    """연결 관계 다이어그램"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "08", "연결 관계 다이어그램", "구조체 간 참조 관계 — 화살표가 \"가리킨다\"를 의미")

    # JTMS
    rect(s, Inches(0.3), Inches(1.2), Inches(4.0), Inches(0.4), BG_CARD, JTMS_C)
    txt(s, Inches(0.4), Inches(1.22), Inches(3.8), Inches(0.4),
        "JTMS", 16, JTMS_C, True, PP_ALIGN.CENTER)

    code(s, Inches(0.3), Inches(1.7), Inches(4.0), Inches(4.5),
         "  jtms\n"
         "  ├── nodes[] ──────────► tms-node\n"
         "  └── justs[] ──────────► just\n"
         "\n"
         "  tms-node\n"
         "  ├── label: :IN / :OUT\n"
         "  ├── support ──────────► just (1개)\n"
         "  ├── justs[] ──────────► [just...]\n"
         "  │   (이 노드가 결론인 정당화들)\n"
         "  ├── consequences[] ──► [just...]\n"
         "  │   (이 노드가 전제인 정당화들)\n"
         "  └── jtms ─────────────► jtms\n"
         "\n"
         "  just\n"
         "  ├── antecedents ──────► [tms-node...]\n"
         "  └── consequence ──────► tms-node\n"
         "\n"
         "  ★ 양방향 참조로 IN→OUT 전파 가능", 10, JTMS_C)

    # LTMS
    rect(s, Inches(4.6), Inches(1.2), Inches(4.0), Inches(0.4), BG_CARD, LTMS_C)
    txt(s, Inches(4.7), Inches(1.22), Inches(3.8), Inches(0.4),
        "LTMS", 16, LTMS_C, True, PP_ALIGN.CENTER)

    code(s, Inches(4.6), Inches(1.7), Inches(4.0), Inches(4.5),
         "  ltms\n"
         "  ├── nodes{} ─────────► tms-node (hash)\n"
         "  └── clauses[] ────────► clause\n"
         "\n"
         "  tms-node\n"
         "  ├── label: :T / :F / :U\n"
         "  ├── support ──────────► clause (1개)\n"
         "  ├── true-clauses[] ──► [clause...]\n"
         "  │   (이 노드의 TRUE 리터럴이 있는 절)\n"
         "  ├── false-clauses[] ─► [clause...]\n"
         "  │   (이 노드의 FALSE 리터럴이 있는 절)\n"
         "  └── ltms ─────────────► ltms\n"
         "\n"
         "  clause\n"
         "  └── literals ─────────► [(node . :T/F)...]\n"
         "\n"
         "  ★ 노드↔절 양방향 참조로 BCP 구현", 10, LTMS_C)

    # ATMS
    rect(s, Inches(8.9), Inches(1.2), Inches(4.0), Inches(0.4), BG_CARD, ATMS_C)
    txt(s, Inches(9.0), Inches(1.22), Inches(3.8), Inches(0.4),
        "ATMS", 16, ATMS_C, True, PP_ALIGN.CENTER)

    code(s, Inches(8.9), Inches(1.7), Inches(4.0), Inches(4.5),
         "  atms\n"
         "  ├── nodes[] ──────────► tms-node\n"
         "  ├── justs[] ──────────► just\n"
         "  ├── env-table[[]] ───► env (버킷)\n"
         "  └── nogood-table[[]] ► env (버킷)\n"
         "\n"
         "  tms-node\n"
         "  ├── label ────────────► [env, env, ...]\n"
         "  │   (최소 환경들의 집합!)\n"
         "  ├── justs[] ──────────► [just...]\n"
         "  ├── consequences[] ──► [just...]\n"
         "  └── atms ─────────────► atms\n"
         "\n"
         "  just: antecedents ──► [node...]\n"
         "        consequence ──► node\n"
         "\n"
         "  env\n"
         "  ├── assumptions ──────► [node...]\n"
         "  └── nodes ────────────► [node...]\n"
         "\n"
         "  ★ env가 node와 양방향 참조", 10, ATMS_C)

    # 하단 비교
    rect(s, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.7), BG_CARD, PURPLE)
    multi(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6), [
        ("핵심 연결 패턴 차이", 14, PURPLE, True),
        ("JTMS: node ↔ just (이분 그래프)    LTMS: node ↔ clause (이분 그래프)    ATMS: node ↔ just + node ↔ env (삼중 연결)", 12, LGRAY),
    ])


def slide_11_pvs_mechanism(prs):
    """pvs 카운터 작동 원리"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "09", "pvs 카운터 작동 원리 (LTMS 전용)", "Boolean Constraint Propagation의 핵심 메커니즘")

    # 예제 설정
    rect(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.7), BG_CARD, LTMS_C)
    multi(s, Inches(0.7), Inches(1.3), Inches(11.8), Inches(0.6), [
        ("예제:  절 C1 = (A ∨ ¬B ∨ D)   \"A가 참이거나, B가 거짓이거나, D가 참\"", 16, WHITE, True),
        ("literals = [(A . :TRUE), (B . :FALSE), (D . :TRUE)]     length = 3", 13, LGRAY, False, PP_ALIGN.LEFT, "Consolas"),
    ])

    # 4 단계
    steps = [
        ("초기 상태", LGRAY,
         "A = :UNKNOWN\nB = :UNKNOWN\nD = :UNKNOWN",
         "pvs = 3\nsats = 0",
         "3개 모두 미결정\n절 상태: 대기"),
        ("B = :TRUE 설정", ORANGE,
         "A = :UNKNOWN\nB = :TRUE  ← 결정!\nD = :UNKNOWN",
         "pvs = 2  (3→2)\nsats = 0",
         "B.:FALSE 기대인데\nB가 :TRUE → 위반!\npvs-- (불일치)"),
        ("A = :FALSE 설정", RED,
         "A = :FALSE ← 결정!\nB = :TRUE\nD = :UNKNOWN",
         "pvs = 1  (2→1)\nsats = 0",
         "A.:TRUE 기대인데\nA가 :FALSE → 위반!\npvs-- → pvs=1!\n★ 단위전파 발동!"),
        ("단위전파 결과", GREEN,
         "A = :FALSE\nB = :TRUE\nD = :TRUE  ← 강제!\n",
         "pvs = 0\nsats = 1  (0→1)",
         "D만 남았으므로\nD = :TRUE 강제!\n(D.:TRUE 만족)\n절 C1 만족됨 ✓"),
    ]

    for i, (title, clr, nodes, counters, desc) in enumerate(steps):
        x = Inches(0.3 + i * 3.25)
        w = Inches(3.1)

        rect(s, x, Inches(2.2), w, Inches(0.4), BG_CARD, clr)
        txt(s, x, Inches(2.22), w, Inches(0.38),
            f"Step {i+1}: {title}", 13, clr, True, PP_ALIGN.CENTER)

        code(s, x + Inches(0.05), Inches(2.7), w - Inches(0.1), Inches(1.0),
             nodes, 11, LGRAY)

        # pvs/sats 강조
        rect(s, x + Inches(0.05), Inches(3.8), w - Inches(0.1), Inches(0.55), BG_CARD,
             RED if "1" in counters.split("\n")[0][-1:] else BORDER)
        txt(s, x + Inches(0.1), Inches(3.82), w - Inches(0.2), Inches(0.5),
            counters, 12, YELLOW, True, PP_ALIGN.CENTER, "Consolas")

        txt(s, x + Inches(0.1), Inches(4.45), w - Inches(0.2), Inches(1.3),
            desc, 12, LGRAY)

    # 핵심 규칙
    txt(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.35),
        "★ pvs 카운터 규칙 정리", 18, YELLOW, True)

    rules = [
        ("pvs = 0", "모순 발생!", "모든 리터럴 위반됨\ncontradiction handler 호출", RED),
        ("pvs = 1", "단위 전파!", "남은 1개 리터럴을 강제로 만족\n→ 해당 노드 값 결정", GREEN),
        ("pvs ≥ 2", "대기", "아직 결정 불가\n다른 노드 결정 시 재검사", DGRAY),
    ]

    for i, (cond, action, desc, clr) in enumerate(rules):
        x = Inches(0.5 + i * 4.2)
        rect(s, x, Inches(6.2), Inches(3.9), Inches(1.1), BG_CARD, clr)
        txt(s, x + Inches(0.1), Inches(6.25), Inches(1.5), Inches(0.35),
            cond, 14, clr, True, PP_ALIGN.LEFT, "Consolas")
        txt(s, x + Inches(1.6), Inches(6.25), Inches(2.1), Inches(0.35),
            action, 14, WHITE, True)
        txt(s, x + Inches(0.1), Inches(6.65), Inches(3.7), Inches(0.6),
            desc, 11, LGRAY)


def slide_12_env_label(prs):
    """env 라벨 작동 원리"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "10", "env 라벨 작동 원리 (ATMS 전용)", "환경 집합의 최소성(minimality) 유지")

    # 예제 설정
    rect(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.6), BG_CARD, ATMS_C)
    multi(s, Inches(0.7), Inches(1.3), Inches(11.8), Inches(0.5), [
        ("예제: assume(A), assume(B), assume(C)   →   justify(J1, D, [A]),  justify(J2, D, [B])", 15, WHITE, True),
    ])

    # 4 단계
    steps = [
        ("가정 생성", LGRAY,
         "assume(A)\nassume(B)\nassume(C)\n\nA.label = [{A}]\nB.label = [{B}]\nC.label = [{C}]",
         "각 가정 노드는 자기 자신만\n포함하는 환경을 label로 가짐\n\n\"A는 A를 가정하면 참\""),
        ("J1: A → D", JTMS_C,
         "justify(J1, D, [A])\n\nweave 계산:\n  antecedents = [A]\n  A.label = [{A}]\n  → 교차곱 = [{A}]\n\nD.label = [{A}]",
         "전제들의 label 교차곱으로\n결론의 환경을 계산\n\n\"A를 가정하면 D가 참\""),
        ("J2: B → D", ATMS_C,
         "justify(J2, D, [B])\n\nweave 계산:\n  antecedents = [B]\n  B.label = [{B}]\n  → 교차곱 = [{B}]\n\nD.label = [{A}, {B}]",
         "기존 label에 새 환경 추가\n\n\"A를 가정해도 D가 참,\n B를 가정해도 D가 참\"\n\n두 독립 경로 동시 기록!"),
        ("최소성 유지", PURPLE,
         "만약 J3: A,B → D 추가하면?\n\nweave = [{A,B}]\n\n{A,B}는 {A}의 상위집합\n{A,B}는 {B}의 상위집합\n\nD.label = [{A}, {B}]\n  ({A,B}는 제거됨!)",
         "최소성(minimality) 규칙:\n\n이미 더 작은 환경이 있으면\n상위집합은 추가하지 않음\n\n\"A만으로 충분한데\n A+B를 기록할 필요 없다\""),
    ]

    for i, (title, clr, content, desc) in enumerate(steps):
        x = Inches(0.3 + i * 3.25)
        w = Inches(3.1)

        rect(s, x, Inches(2.1), w, Inches(0.4), BG_CARD, clr)
        txt(s, x, Inches(2.12), w, Inches(0.38),
            f"Step {i+1}: {title}", 13, clr, True, PP_ALIGN.CENTER)

        code(s, x + Inches(0.05), Inches(2.6), w - Inches(0.1), Inches(2.4),
             content, 10, clr if clr != LGRAY else CYAN)

        txt(s, x + Inches(0.1), Inches(5.1), w - Inches(0.2), Inches(1.5),
            desc, 11, LGRAY)

    # nogood 처리
    txt(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.35),
        "★ nogood 처리 — 환경이 모순될 때", 16, RED, True)

    rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), BG_CARD, RED)
    txt(s, Inches(0.7), Inches(7.02), Inches(11.8), Inches(0.3),
        "{A,C}가 nogood → {A,C}.nogood?=just   →   {A,C}를 포함하는 모든 상위집합({A,B,C} 등)도 자동 nogood   →   모든 노드 label에서 제거",
        11, LGRAY, False, PP_ALIGN.LEFT, "Consolas")


def slide_13_summary(prs):
    """설계 철학 요약"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    header(s, "11", "설계 철학 요약 — 한눈에 보는 종합 비교")

    col3_headers(s, Inches(1.2))

    rows = [
        ("구조체 수",          "3개",                    "3개",                      "4개 (env 추가)"),
        ("추론 단위",          "just (전제→결론)",       "clause (리터럴의 OR)",      "just + env"),
        ("노드 label",        ":IN | :OUT (2값)",       ":TRUE|:FALSE|:UNKNOWN (3값)", "[env...] (환경 집합)"),
        ("노드 support",      "→ just (1개)",           "→ clause (1개)",            "없음 (label이 근거)"),
        ("그래프 구조",        "node ↔ just\n(이분 그래프)", "node ↔ clause\n(이분 그래프)", "node ↔ just\n+ node ↔ env"),
        ("핵심 카운터",        "없음",                   "pvs / sats\n(BCP 최적화)",   "env.count\n(버킷팅 최적화)"),
        ("메모리 특성",        "가장 가벼움",            "중간\n(리터럴 캐시 추가)",    "가장 무거움\n(env 테이블)"),
        ("규칙 트리거",        "in-rules\nout-rules",    "true-rules\nfalse-rules",   "rules (하나)\n∅→비∅ 전환 시"),
    ]

    y0 = Inches(1.72)
    rh = Inches(0.55)
    for i, (label, j, l, a) in enumerate(rows):
        y = y0 + i * (rh + Inches(0.04))
        h = Inches(0.7) if "\n" in j or "\n" in l or "\n" in a else rh
        table_row(s, y, label, [j, l, a], h)
        if h > rh:
            y0 += (h - rh)

    # 핵심 한줄
    rect(s, Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.45), BG_CARD, PURPLE)
    txt(s, Inches(0.5), Inches(6.88), Inches(12.3), Inches(0.4),
        "JTMS는 가장 단순 (3 struct)  |  LTMS는 pvs 카운터가 핵심  |  ATMS는 env 구조체가 모든 것을 바꾼다",
        14, WHITE, True, PP_ALIGN.CENTER)


def slide_14_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    rect(s, Inches(0), Inches(0), SW, Pt(4), BLUE)

    txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
        "핵심 정리", 42, BLUE, True, PP_ALIGN.CENTER)

    rect(s, Inches(4.5), Inches(2.5), Inches(4.3), Pt(2), DGRAY)

    insights = [
        ("label 필드", "가 세 시스템의 정체성을 결정한다", WHITE),
        ("  JTMS:", "  :IN/:OUT → 단일 시나리오, 순방향", JTMS_C),
        ("  LTMS:", "  :T/:F/:U → 단일 시나리오, 양방향", LTMS_C),
        ("  ATMS:", "  [env...] → 다중 시나리오, 동시", ATMS_C),
        ("", "", DGRAY),
        ("LTMS의 pvs", " 카운터 = O(1) 전파 결정의 핵심", YELLOW),
        ("ATMS의 env", " 구조체 = 다중 세계 동시 관리의 핵심", ORANGE),
    ]

    for i, (key, val, clr) in enumerate(insights):
        y = Inches(3.0 + i * 0.45)
        txt(s, Inches(2.5), y, Inches(3.5), Inches(0.4),
            key, 16, clr, True, PP_ALIGN.RIGHT)
        txt(s, Inches(6.0), y, Inches(5.0), Inches(0.4),
            val, 16, LGRAY)

    txt(s, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
        "자료구조를 이해하면, 알고리즘은 자연스럽게 따라온다.", 18, WHITE, True, PP_ALIGN.CENTER)

    txt(s, Inches(1), Inches(7.0), Inches(11.3), Inches(0.3),
        "Building Problem Solvers  |  Forbus & de Kleer", 14, DGRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
#  메인
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_toc(prs)
    slide_03_system_struct(prs)
    slide_04_node_core_diff(prs)
    slide_05_node_full_table(prs)
    slide_06_inference_units(prs)
    slide_07_jtms_just_detail(prs)
    slide_08_ltms_clause_detail(prs)
    slide_09_atms_env_detail(prs)
    slide_10_connections(prs)
    slide_11_pvs_mechanism(prs)
    slide_12_env_label(prs)
    slide_13_summary(prs)
    slide_14_closing(prs)

    out = "/home/user/bps/ltms/TMS_자료구조_상세비교.pptx"
    prs.save(out)
    print(f"PPT 생성 완료: {out}")
    print(f"총 {len(prs.slides)}장의 슬라이드")


if __name__ == "__main__":
    main()
