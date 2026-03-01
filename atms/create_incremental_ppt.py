#!/usr/bin/env python3
"""Incremental Accumulation 프로그래밍 기법과 ATMS 적용 사례 교육자료 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK    = RGBColor(0x1E, 0x1E, 0x2E)
BG_SURFACE = RGBColor(0x2A, 0x2A, 0x3C)
ACCENT1    = RGBColor(0x89, 0xB4, 0xFA)   # 밝은 파랑
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # 밝은 초록
ACCENT3    = RGBColor(0xF9, 0xE2, 0xAF)   # 밝은 노랑
ACCENT4    = RGBColor(0xF3, 0x8B, 0xA8)   # 밝은 분홍
WHITE      = RGBColor(0xCD, 0xD6, 0xF4)
GRAY       = RGBColor(0x93, 0x99, 0xB2)
CODE_BG    = RGBColor(0x31, 0x32, 0x44)
ACCENT5    = RGBColor(0xCB, 0xA6, 0xF7)   # 밝은 보라
TEAL       = RGBColor(0x94, 0xE2, 0xD5)   # 청록
PEACH      = RGBColor(0xFA, 0xB3, 0x87)   # 피치

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(shape, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Malgun Gothic", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_code_block(slide, left, top, width, height, lines):
    shape = add_shape(slide, left, top, width, height, fill_color=CODE_BG, border_color=RGBColor(0x58, 0x5B, 0x70))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT3
        p.font.name = "Consolas"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
    return shape


def make_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    set_text(tb, title, size=40, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = add_text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(1.0))
    set_text(tb2, subtitle, size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT1
    line.line.fill.background()
    return slide


def make_section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
    set_text(tb, f"#{number}", size=60, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.0))
    set_text(tb2, title, size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        tb3 = add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
        set_text(tb3, subtitle, size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def make_content_slide(title, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_SURFACE
    bar.line.fill.background()
    tb = add_text_box(slide, Inches(0.6), Inches(0.1), Inches(10), Inches(0.7))
    set_text(tb, title, size=24, color=ACCENT1, bold=True)
    if page_num:
        tb_pg = add_text_box(slide, Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.6))
        set_text(tb_pg, page_num, size=14, color=GRAY, alignment=PP_ALIGN.RIGHT, font_name="Consolas")
    return slide


# ================================================================
# 슬라이드 1: 표지
# ================================================================
make_title_slide(
    "Incremental Accumulation",
    "변경분만 추적하는 프로그래밍 기법과 ATMS에서의 적용"
)

# ================================================================
# 슬라이드 2: 목차
# ================================================================
slide = make_content_slide("목차", "2")
items = [
    ("1", "Incremental Accumulation이란", "기법의 정의, 핵심 원리, 비유"),
    ("2", "전체 재계산 vs Incremental      ", "두 접근법의 구체적 비교"),
    ("3", "ATMS 전파 파이프라인 전체 흐름    ", "justify → propagate → weave → update-label"),
    ("4", "propagate: Delta 전달             ", "변경분만 전달하는 전파 엔진"),
    ("5", "weave: 단계별 교차곱 누적         ", "전건별 순차 처리와 중간 필터링"),
    ("6", "update-label: Delta 필터링        ", "Lazy Deletion과 최소성 유지"),
    ("7", "union-env / cons-env              ", "환경 합치기의 조기 종료와 구조적 공유"),
    ("8", "nogood 역전파                     ", "new-nogood의 역방향 incremental"),
    ("9", "5가지 구현 기법 종합              ", "Delta 추적, Lazy Deletion, 조기 종료, 구조적 공유, 단계별 필터"),
]
y = Inches(1.1)
for num, title, desc in items:
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.65))
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.03), Inches(0.5), Inches(0.3))
    set_text(tb, num, size=18, color=ACCENT2, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(5.5), Inches(0.3))
    set_text(tb2, title, size=16, color=WHITE, bold=True)
    tb3 = add_text_box(slide, Inches(1.5), y + Inches(0.33), Inches(10), Inches(0.28))
    set_text(tb3, desc, size=12, color=GRAY)
    y += Inches(0.69)


# ================================================================
# 슬라이드 3: Incremental Accumulation이란
# ================================================================
make_section_slide("1", "Incremental Accumulation이란",
                   "\"전체를 다시 만들지 말고, 바뀐 것만 반영하라\"")

slide = make_content_slide("Incremental Accumulation의 정의", "3")

# 정의 카드
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5), border_color=ACCENT1)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.2))
tf = set_text(tb, "정의", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "상태 전체를 매번 재계산하지 않고, 변경분(Delta, △)만 계산하여 기존 상태에 점진적으로 축적하는 기법",
              size=18, color=ACCENT3, bold=True)
add_paragraph(tf, "의존하는 다른 부분에도 변경분만 전달하여, 불필요한 재계산을 최소화합니다.", size=14, color=GRAY)

# 3가지 핵심 원리
card = add_shape(slide, Inches(0.5), Inches(3.0), Inches(3.8), Inches(4.0), border_color=ACCENT2)
tb = add_text_box(slide, Inches(0.8), Inches(3.15), Inches(3.3), Inches(3.7))
tf = set_text(tb, "원리 1: Delta 추적", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, '"무엇이 바뀌었는가?"만 식별', size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "비유: 회계 장부", size=16, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "매일 모든 거래를 처음부터", size=14, color=WHITE)
add_paragraph(tf, "다시 합산하지 않는다.", size=14, color=WHITE)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "어제 잔고: 1,000원", size=14, color=GRAY)
add_paragraph(tf, "오늘 입금: +500원  ← △", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "오늘 잔고: 1,500원", size=14, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ △(+500)만 반영하면 됨", size=14, color=ACCENT3)

card = add_shape(slide, Inches(4.6), Inches(3.0), Inches(4.0), Inches(4.0), border_color=ACCENT4)
tb = add_text_box(slide, Inches(4.9), Inches(3.15), Inches(3.5), Inches(3.7))
tf = set_text(tb, "원리 2: 선택적 전파", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, '"영향받는 곳에만 알려라"', size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "비유: 뉴스 알림", size=16, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "모든 사람에게 신문을", size=14, color=WHITE)
add_paragraph(tf, "다시 배달하지 않는다.", size=14, color=WHITE)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "변경: 속보 1건 추가 ← △", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "→ 구독자에게만 푸시 알림", size=14, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ △이 없으면 전파도 없음", size=14, color=ACCENT3)

card = add_shape(slide, Inches(8.9), Inches(3.0), Inches(3.9), Inches(4.0), border_color=ACCENT5)
tb = add_text_box(slide, Inches(9.2), Inches(3.15), Inches(3.4), Inches(3.7))
tf = set_text(tb, "원리 3: 단계별 필터링", size=18, color=ACCENT5, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, '"불필요한 것은 즉시 제거"', size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "비유: 공장 품질 검사", size=16, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "완제품을 다 만들고 나서", size=14, color=WHITE)
add_paragraph(tf, "불량 검사를 하지 않는다.", size=14, color=WHITE)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "각 공정마다 검사 ← 필터", size=14, color=ACCENT5, font_name="Consolas")
add_paragraph(tf, "→ 불량은 다음 공정 진입 X", size=14, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ 최종 단계의 부담 감소", size=14, color=ACCENT3)


# ================================================================
# 슬라이드 4: 전체 재계산 vs Incremental
# ================================================================
make_section_slide("2", "전체 재계산 vs Incremental", "두 접근법의 구체적 비교")

slide = make_content_slide("전체 재계산 vs Incremental Accumulation", "4")

# 좌: 전체 재계산
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(5.5))
tf = set_text(tb, "X  전체 재계산 (Batch Recomputation)", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "매번 label 전체를 처음부터 계산", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "의사코드:", size=15, color=GRAY)
add_paragraph(tf, "  on_change(node):", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    new_label = []", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    for just in node.justs:", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "      envs = cross_product(", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "        all antecedent labels)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "      new_label += filter(envs)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    node.label = minimize(new_label)", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "    for dep in dependents(node):", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "      on_change(dep)  # 전체 재귀!", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "문제점", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, "  1. 모든 정당화의 교차곱 재계산", size=14, color=WHITE)
add_paragraph(tf, "  2. 변화 없어도 후속 노드에 전파", size=14, color=WHITE)
add_paragraph(tf, "  3. O(m^n) 조합 폭발", size=14, color=WHITE)
add_paragraph(tf, "  4. 최종 단계에서만 필터링 → 낭비", size=14, color=WHITE)

# 우: Incremental
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "O  Incremental Accumulation", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "변경분(△)만 계산하여 기존에 병합", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "의사코드 (ATMS 실제 방식):", size=15, color=GRAY)
add_paragraph(tf, "  propagate(just, changed, △):", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    new = weave(changed,", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "       △, antecedents)  # △만!", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "    △' = update_label(", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "       node, new)  # 실제 추가분만", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "    if △' is empty:", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "      return  # 전파 중단!", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "    for dep in dependents:", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "      propagate(dep_just,", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "        node, △')  # △'만 전달", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "장점", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "  1. 변경분에 대한 교차곱만 계산", size=14, color=WHITE)
add_paragraph(tf, "  2. 변화 없으면 전파 즉시 중단", size=14, color=WHITE)
add_paragraph(tf, "  3. 각 단계에서 필터링 → 폭발 방지", size=14, color=WHITE)
add_paragraph(tf, "  4. 구조적 공유로 메모리 절감", size=14, color=WHITE)


# ================================================================
# 슬라이드 5: ATMS 전파 파이프라인 전체 흐름
# ================================================================
make_section_slide("3", "ATMS 전파 파이프라인 전체 흐름",
                   "justify → propagate → weave → update → update-label")

slide = make_content_slide("△(Delta)가 흘러가는 전체 경로", "5")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.8), Inches(5.5))
tf = set_text(tb, "파이프라인: 각 단계에서 △만 전달", size=19, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, '  justify-node("덧셈규칙", C, [A, B])               ▸ 사용자 호출',
              size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "     │  △₀ = [{}]                                      빈 환경 = 시작점",
              size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  propagate(just, nil, △₀=[{}])                      ▸ △₀를 weave에 전달",
              size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "     │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  weave(nil, △₀=[{}], [A, B])                        ▸ △₀ × A.label × B.label",
              size=14, color=TEAL, font_name="Consolas")
add_paragraph(tf, "     │  A.label=[{A1}], B.label=[{A2}]",
              size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  {} × {A1} × {A2} = {A1,A2}",
              size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  △₁ = [{A1,A2}]                                 결과(후보)",
              size=13, color=TEAL, font_name="Consolas")
add_paragraph(tf, "     ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  update(△₁=[{A1,A2}], C, just)                      ▸ 모순 체크 → label 갱신",
              size=14, color=PEACH, font_name="Consolas")
add_paragraph(tf, "     │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  update-label(C, △₁=[{A1,A2}])                      ▸ 기존 label에 병합",
              size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "     │  C.label: [] → [{A1,A2}]                        실제 변경 발생",
              size=13, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "     │  △₂ = [{A1,A2}]                                 새로 추가된 것만 반환!",
              size=13, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "     ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  propagate(other_just, C, △₂=[{A1,A2}])             ▸ C의 후속 노드에 △₂만 전달",
              size=14, color=ACCENT5, font_name="Consolas")
add_paragraph(tf, "     │  C의 전체 label이 아님! 새로 추가된 부분만!",
              size=13, color=ACCENT5, font_name="Consolas")
add_paragraph(tf, "     │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     └─ △₂가 비어있으면? → 전파 중단! (unless new-envs (return-from update nil))",
              size=13, color=ACCENT2, font_name="Consolas")


# ================================================================
# 슬라이드 6: propagate — Delta 전달
# ================================================================
make_section_slide("4", "propagate: Delta 전달",
                   "변경분만 전달하는 전파 엔진")

slide = make_content_slide("propagate: 변경분만 전달하는 전파 엔진", "6")

# 코드 + 분석
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5), border_color=ACCENT1)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(5.3), Inches(0.3))
set_text(tb, "propagate 코드 (atms.lisp:415)", size=16, color=ACCENT1, bold=True)
add_code_block(slide, Inches(0.7), Inches(1.75), Inches(5.4), Inches(1.8), [
    "(defun propagate (just antecedent envs",
    "                  &aux new-envs)",
    "  (if (setq new-envs",
    "            (weave antecedent envs",
    "                   (just-antecedents just)))",
    "      (update new-envs",
    "             (just-consequence just)",
    "             just)))",
])

# 오른쪽: 핵심 포인트
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.5), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(2.2))
tf = set_text(tb, "핵심: 매개변수가 곧 Delta", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "envs = 전체 label이 아닌, 변경분만!", size=15, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "  초기: propagate(just, nil, [{}])", size=14, color=TEAL, font_name="Consolas")
add_paragraph(tf, "        → 빈 환경 1개 = 시작 신호", size=13, color=GRAY)
add_paragraph(tf, "  재귀: propagate(just, C, new-envs)", size=14, color=PEACH, font_name="Consolas")
add_paragraph(tf, "        → update-label이 반환한 △만 전달", size=13, color=GRAY)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "weave가 nil 반환 → update 호출 안 함 = 전파 중단", size=14, color=ACCENT4)

# 하단: 호출 관계
card = add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0), border_color=TEAL)
tb = add_text_box(slide, Inches(0.8), Inches(4.15), Inches(11.8), Inches(2.7))
tf = set_text(tb, "antecedent 매개변수의 역할 — 중복 계산 방지", size=17, color=TEAL, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "상황: A와 B가 전건, C가 결론. A의 label이 변경됨", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  propagate(just, A, △_A)    ← A의 변경분", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    └─ weave(A, △_A, [A, B])", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "         │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "         ├─ node=A이면? (eq node antecedent) → 건너뜀!  △_A에 이미 A 반영됨",
              size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "         │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "         └─ node=B이면? → B.label 전체와 교차곱  (B는 안 바뀌었으니 전체 사용)",
              size=14, color=PEACH, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "  즉, △_A × B.label(전체) = 새 환경 후보들", size=15, color=ACCENT3, bold=True)
add_paragraph(tf, '  A.label(전체) × B.label(전체)가 아님! — 이것이 "incremental"의 핵심', size=14, color=ACCENT4)


# ================================================================
# 슬라이드 7: weave — 단계별 교차곱 누적
# ================================================================
make_section_slide("5", "weave: 단계별 교차곱 누적",
                   "전건별 순차 처리와 중간 필터링")

slide = make_content_slide("weave: 교차곱을 단계별로 누적하며 필터링", "7")

# 왼쪽: 코드
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=TEAL)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(5.3), Inches(0.3))
set_text(tb, "weave 코드 (atms.lisp:520)", size=16, color=TEAL, bold=True)
add_code_block(slide, Inches(0.7), Inches(1.7), Inches(5.4), Inches(5.1), [
    "(defun weave (antecedent envs antecedents",
    "              &aux new-envs new-env)",
    "  (setq envs (copy-list envs))",
    "  (dolist (node antecedents)  ; 전건 하나씩",
    "    (unless (eq node antecedent) ; 이미 반영된 것 건너뜀",
    "      (setq new-envs nil)",
    "      (dolist (env envs)      ; 현재 △",
    "        (if env",
    "          (dolist (node-env   ; × 이 전건의 label",
    "                   (tms-node-label node))",
    "            (setq new-env",
    "              (union-env env node-env))",
    "            (unless (env-nogood? new-env) ;필터1",
    "              (do ((nnew-envs new-envs",
    "                    (cdr nnew-envs)))",
    "                  ((null nnew-envs)",
    "                   (push new-env new-envs))",
    "                (when (car nnew-envs)",
    "                  (case (compare-env new-env",
    "                          (car nnew-envs))",
    "                    ((:EQ :S21) (return nil))",
    "                    (:S12  ; 필터2: 최소성",
    "                     (rplaca nnew-envs",
    "                             nil)))))))))",
    "      (setq envs  ; △ 갱신",
    "        (delete nil new-envs :TEST #'eq))",
    "      (unless envs  ; △=[] → 전체 실패",
    "        (return-from weave nil))))",
    "  envs)",
])

# 오른쪽: 단계별 예시
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=ACCENT3)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "3개 필터의 역할", size=18, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "필터 1  nogood 차단", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, "  (unless (env-nogood? new-env) ...)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 모순 환경은 즉시 탈락", size=13, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "필터 2  최소성 유지", size=16, color=ACCENT5, bold=True)
add_paragraph(tf, "  :EQ/:S21 → 새것 탈락 (기존이 더 작음)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  :S12    → 기존 것 제거 (새것이 더 작음)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 중간 결과도 항상 최소 유지!", size=13, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "필터 3  조기 종료", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "  (unless envs", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (return-from weave nil))", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 중간에 결과가 공집합이면 즉시 종료", size=13, color=GRAY)
add_paragraph(tf, "", size=12, color=WHITE)

add_paragraph(tf, "단계별 누적 예시", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "  △₀ = [{}]       (시작점)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  × A.label=[{A1},{A2}]", size=14, color=TEAL, font_name="Consolas")
add_paragraph(tf, "  △₁ = [{A1},{A2}] (필터 적용 후)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  × B.label=[{A3}]", size=14, color=TEAL, font_name="Consolas")
add_paragraph(tf, "  △₂ = [{A1,A3},{A2,A3}] (최종)", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  각 단계에서 필터링 → 조합 폭발 방지", size=14, color=ACCENT3, bold=True)


# ================================================================
# 슬라이드 8: update-label — Delta 필터링
# ================================================================
make_section_slide("6", "update-label: Delta 필터링",
                   "Lazy Deletion과 최소성 유지")

slide = make_content_slide("update-label: 실제로 추가된 것만 반환", "8")

# 왼쪽: 코드
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(5.3), Inches(0.3))
set_text(tb, "update-label 코드 (atms.lisp:474)", size=16, color=ACCENT4, bold=True)
add_code_block(slide, Inches(0.7), Inches(1.7), Inches(5.4), Inches(5.1), [
    "(defun update-label (node new-envs",
    "                     &aux envs)",
    "  (setq envs (tms-node-label node))",
    "",
    "  ;; 이중 루프: 새 환경 × 기존 label",
    "  (do ((new-envs new-envs (cdr new-envs)))",
    "      ((null new-envs))",
    "    (do ((nenvs envs (cdr nenvs)))",
    "        ((null nenvs)",
    "         (push (car new-envs) envs))",
    "      (cond",
    "        ((null (car nenvs)))  ; 이미 제거됨",
    "        ((null (car new-envs))) ; 이미 탈락",
    "        ((case (compare-env",
    "                 (car new-envs)",
    "                 (car nenvs))",
    "           ((:EQ :S21)  ; 새것 >= 기존",
    "            (rplaca new-envs nil)) ;탈락",
    "           (:S12  ; 새것 < 기존",
    "            (setf (env-nodes ...)",
    "              (delete node ...))",
    "            (rplaca nenvs nil)))))) ;기존제거",
    "",
    "  ;; 최종: nil 정리",
    "  (setq new-envs",
    "    (delete nil new-envs :TEST #'eq))",
    "  (dolist (new-env new-envs)",
    "    (push node (env-nodes new-env)))",
    "  (setf (tms-node-label node)",
    "    (delete nil envs :TEST #'eq))",
    "  new-envs)  ; ← △만 반환!",
])

# 오른쪽: Lazy Deletion + 예시
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.6), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(2.3))
tf = set_text(tb, "Lazy Deletion 패턴", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "단계 1: 제거할 것을 nil로 표시 (논리적 삭제)", size=15, color=WHITE)
add_paragraph(tf, "  (rplaca new-envs nil)   ; 새것 탈락", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (rplaca nenvs nil)      ; 기존 것 제거", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "단계 2: 마지막에 한 번만 물리적 삭제", size=15, color=WHITE)
add_paragraph(tf, "  (delete nil new-envs :TEST #'eq)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (delete nil envs :TEST #'eq)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ 순회 중 리스트 구조 변경 X → 안전!", size=14, color=TEAL, bold=True)

card = add_shape(slide, Inches(6.7), Inches(4.1), Inches(6.1), Inches(2.9), border_color=ACCENT3)
tb = add_text_box(slide, Inches(7.0), Inches(4.25), Inches(5.5), Inches(2.6))
tf = set_text(tb, "구체적 예시", size=18, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "기존 label = [{A1}, {A2,A3}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "새 후보    = [{A2}, {A1,A4}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "비교 1: {A2} vs {A1} → 독립 → 통과", size=13, color=GRAY)
add_paragraph(tf, "비교 2: {A2} vs {A2,A3} → :S12 → {A2,A3} 제거!", size=13, color=ACCENT4)
add_paragraph(tf, "비교 3: {A1,A4} vs {A1} → :S21 → {A1,A4} 탈락!", size=13, color=ACCENT4)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "결과 label = [{A1}, {A2}]     (최소성 유지)", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "반환 △    = [{A2}]            (실제 추가분만!)", size=14, color=ACCENT2, font_name="Consolas")


# ================================================================
# 슬라이드 9: union-env / cons-env
# ================================================================
make_section_slide("7", "union-env / cons-env",
                   "환경 합치기의 조기 종료와 구조적 공유")

slide = make_content_slide("union-env: 조기 종료 / cons-env: 구조적 공유", "9")

# 좌: union-env
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=PEACH)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(5.3), Inches(0.3))
set_text(tb, "union-env: 한 번에 하나씩, 조기 종료", size=17, color=PEACH, bold=True)

add_code_block(slide, Inches(0.7), Inches(1.75), Inches(5.4), Inches(2.2), [
    "(defun union-env (e1 e2)",
    "  (when (> (env-count e1)",
    "           (env-count e2))",
    "    (psetq e1 e2 e2 e1))  ; 작은 쪽을 e1으로",
    "  (dolist (assume (env-assumptions e1))",
    "    (setq e2 (cons-env assume e2)) ;하나씩 추가",
    "    (if (env-nogood? e2)  ; 모순 발견?",
    "        (return nil)))   ; → 즉시 종료!",
    "  e2)",
])

tb = add_text_box(slide, Inches(0.8), Inches(4.15), Inches(5.2), Inches(2.7))
tf = set_text(tb, "조기 종료 효과 (Incremental)", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "E1={A1,A2,A3}, E2={A4}", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  Step 1: A1 추가 → {A1,A4}  일관 O", size=13, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "  Step 2: A2 추가 → {A1,A2,A4} 모순 X", size=13, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "  → nil 반환! (A3은 검사하지 않음)", size=13, color=ACCENT4)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "전체를 합친 뒤 검사하면:", size=14, color=GRAY)
add_paragraph(tf, "  {A1,A2,A3,A4} 생성 → 검사 → nogood", size=13, color=GRAY)
add_paragraph(tf, "  → A3 추가까지 불필요하게 수행", size=13, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ 하나씩 추가하며 검사 = 낭비 최소화!", size=14, color=ACCENT3, bold=True)

# 우: cons-env
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=ACCENT5)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(0.3))
set_text(tb, "cons-env: 구조적 공유 (Structural Sharing)", size=17, color=ACCENT5, bold=True)

add_code_block(slide, Inches(6.9), Inches(1.75), Inches(5.9), Inches(2.0), [
    "(defun cons-env (assumption env &aux nassumes)",
    "  (setq nassumes",
    "    (ordered-insert assumption",
    "       (env-assumptions env)",
    "       #'assumption-order))",
    "  (or (lookup-env nassumes)   ; 이미 있으면 재사용!",
    "      (create-env             ; 없으면 새로 생성",
    "        (tms-node-atms assumption)",
    "        nassumes)))",
])

tb = add_text_box(slide, Inches(7.0), Inches(3.95), Inches(5.5), Inches(2.9))
tf = set_text(tb, "구조적 공유 효과", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(or (lookup-env ...) (create-env ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "패턴: Memoization / Flyweight", size=15, color=ACCENT5, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  같은 가정 조합 {A1,A3}이 여러 경로에서 생성:", size=14, color=WHITE)
add_paragraph(tf, "    경로1: {A1} × {A3} → {A1,A3}", size=13, color=TEAL, font_name="Consolas")
add_paragraph(tf, "    경로2: {A3} × {A1} → {A1,A3}", size=13, color=TEAL, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  lookup-env: 이미 있음 → 기존 객체 반환", size=14, color=ACCENT2)
add_paragraph(tf, "  → 중복 생성 X, 동일 객체 공유", size=14, color=ACCENT2)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "효과: 메모리 절감 + eq 비교 가능", size=15, color=ACCENT3, bold=True)
add_paragraph(tf, "(같은 환경 = 같은 포인터)", size=14, color=GRAY)


# ================================================================
# 슬라이드 10: nogood 역전파
# ================================================================
make_section_slide("8", "nogood 역전파",
                   "new-nogood의 역방향 Incremental 처리")

slide = make_content_slide("new-nogood: 역방향 Incremental 전파", "10")

# 왼쪽: 코드
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(5.3), Inches(0.3))
set_text(tb, "new-nogood 코드 (atms.lisp:760)", size=16, color=ACCENT4, bold=True)
add_code_block(slide, Inches(0.7), Inches(1.75), Inches(5.4), Inches(5.0), [
    "(defun new-nogood (atms cenv just",
    "                   &aux count)",
    "  (setf (env-nogood? cenv) just)",
    "  (remove-env-from-labels cenv atms)",
    "  ;; nogood-table에 등록",
    "  (setf (atms-nogood-table atms)",
    "    (insert-in-table",
    "      (atms-nogood-table atms) cenv))",
    "",
    "  (setq count (env-count cenv))",
    "",
    "  ;; 기존 nogood 중 상위집합 제거 (최소성)",
    "  (dolist (entry (atms-nogood-table atms))",
    "    (when (> (car entry) count)",
    "      (dolist (old (cdr entry))",
    "        (if (subset-env? cenv old)",
    "          (setf (cdr entry)",
    "            (delete old (cdr entry)",
    "                   :COUNT 1))))))",
    "",
    "  ;; env-table에서 상위집합 전부 nogood 표시",
    "  (dolist (entry (atms-env-table atms))",
    "    (when (> (car entry) count)",
    "      (dolist (old (cdr entry))",
    "        (when (and (not (env-nogood? old))",
    "                   (subset-env? cenv old))",
    "          (setf (env-nogood? old) cenv)",
    "          (remove-env-from-labels",
    "            old atms))))))",
])

# 오른쪽: 분석
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=TEAL)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "역방향 Incremental 패턴", size=18, color=TEAL, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "정방향: justify → label에 env 추가", size=15, color=ACCENT2, bold=True)
add_paragraph(tf, "역방향: nogood → label에서 env 제거", size=15, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, "△ = 새로 발견된 nogood 1개 (cenv)", size=16, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "Step 1: cenv를 모든 label에서 제거", size=15, color=ACCENT1, bold=True)
add_paragraph(tf, "  remove-env-from-labels(cenv)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → cenv가 있는 노드들의 label에서만 제거", size=13, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "Step 2: nogood-table 최소성 유지", size=15, color=ACCENT1, bold=True)
add_paragraph(tf, "  {A1} 발견 → {A1,A2}, {A1,A3} 제거", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 더 작은 nogood가 있으므로 큰 것은 중복", size=13, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "Step 3: 상위집합도 자동 nogood", size=15, color=ACCENT1, bold=True)
add_paragraph(tf, "  {A1} nogood → {A1,A4,A5}도 nogood!", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 부분이 모순이면 전체도 모순", size=13, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, "최적화: count 기반 조기 검색 범위 제한", size=15, color=ACCENT5, bold=True)
add_paragraph(tf, "  (when (> (car entry) count) ...)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → cenv보다 큰 환경만 검사 (작은 것은", size=13, color=GRAY)
add_paragraph(tf, "    cenv의 상위집합이 될 수 없으므로)", size=13, color=GRAY)


# ================================================================
# 슬라이드 11: 5가지 구현 기법 종합
# ================================================================
make_section_slide("9", "5가지 구현 기법 종합",
                   "ATMS에서 Incremental Accumulation을 구현하는 방법들")

slide = make_content_slide("ATMS의 Incremental Accumulation 5가지 구현 기법", "11")

techniques = [
    ("#1 Delta 추적 및 전달", ACCENT2,
     "propagate, update에서 △만 전달",
     "propagate(just, C, new-envs)  ← new-envs = △",
     "전체 label이 아닌 \"새로 추가된 것\"만 후속 노드에 전달. △이 []이면 전파 중단."),
    ("#2 Lazy Deletion", ACCENT4,
     "update-label, weave에서 rplaca nil 패턴",
     "(rplaca item nil) → 나중에 (delete nil list)",
     "순회 중 리스트 구조를 변경하지 않고 논리적 삭제 후 한 번에 물리적 정리."),
    ("#3 조기 종료 (Early Exit)", TEAL,
     "union-env, weave, update에서 불필요한 계산 차단",
     "(if (env-nogood? e2) (return nil))",
     "모순 발견, 결과 공집합, 변화 없음 등의 조건에서 즉시 종료하여 낭비 방지."),
    ("#4 구조적 공유 (Memoization)", ACCENT5,
     "cons-env의 lookup-env / create-env 패턴",
     "(or (lookup-env nassumes) (create-env ...))",
     "같은 가정 조합은 같은 객체를 재사용. 중복 생성 방지 + eq 비교 가능 + 메모리 절감."),
    ("#5 단계별 필터링", ACCENT1,
     "weave의 전건별 순차 처리 + 중간 최소성 검사",
     "nogood 필터 → 최소성 필터 → 공집합 검사",
     "각 전건 처리 시마다 3중 필터 적용. 최종 단계의 조합 폭발을 사전에 차단."),
]

y = Inches(1.1)
for title, color, where, code, desc in techniques:
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.22), border_color=color)
    # 제목
    tb = add_text_box(slide, Inches(0.8), y + Inches(0.06), Inches(3.6), Inches(0.3))
    set_text(tb, title, size=16, color=color, bold=True, font_name="Consolas")
    # 적용 위치
    tb = add_text_box(slide, Inches(4.5), y + Inches(0.04), Inches(8.1), Inches(0.3))
    set_text(tb, where, size=14, color=WHITE)
    # 코드 예시
    tb = add_text_box(slide, Inches(0.8), y + Inches(0.38), Inches(11.8), Inches(0.3))
    set_text(tb, code, size=12, color=ACCENT3, font_name="Consolas")
    # 설명
    tb = add_text_box(slide, Inches(0.8), y + Inches(0.62), Inches(11.8), Inches(0.55))
    set_text(tb, desc, size=12, color=GRAY)
    y += Inches(1.27)


# ================================================================
# 슬라이드 12: 성능 비교 정리
# ================================================================
slide = make_content_slide("전체 재계산 vs Incremental: 성능 비교", "12")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.0), border_color=ACCENT1)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(2.7))
tf = set_text(tb, "시나리오: 3개 전건, 각각 3개 환경", size=18, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)

add_paragraph(tf, "  전체 재계산                                     Incremental Accumulation", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  ─────────────────────────────                   ─────────────────────────────", size=12, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  모든 조합: 3 × 3 × 3 = 27개 생성               단계 1: 3개 (필터링 후 ≤3)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  → 전부 생성 후 필터링                           단계 2: ≤3 × 3 = ≤9 (필터링 후 ≤9)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  → 전체 label 비교 → 전체 전파                   단계 3: ≤9 × 3 = ≤27 (필터링 후 ≤27)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  복잡도: O(m^n)                                  복잡도: O(n × m²) — 필터링 효과로 실제 더 적음", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  변화 없어도 후속 전체 전파                       △=[] → 전파 즉시 중단", size=14, color=ACCENT3, font_name="Consolas")

# 하단: 마무리
card = add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.7), border_color=ACCENT2)
tb = add_text_box(slide, Inches(0.8), Inches(4.65), Inches(11.8), Inches(2.4))
tf = set_text(tb, "핵심 요약", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "Incremental Accumulation은 \"변경분 △만 계산하고 전달하는\" 프로그래밍 기법입니다.", size=17, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "ATMS에서는 이 기법이 전파 파이프라인의 모든 단계에 적용되어 있습니다:", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  propagate: △만 전달  →  weave: 단계별 필터링  →  update-label: 실제 추가분만 반환  →  재귀: △만 전파",
              size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "이 기법 덕분에 ATMS는 복잡한 추론 네트워크에서도 효율적으로 label을 유지할 수 있습니다.",
              size=16, color=TEAL, bold=True)


# ================================================================
# 저장
# ================================================================
output_path = "/home/user/bps/atms/incremental-accumulation.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
