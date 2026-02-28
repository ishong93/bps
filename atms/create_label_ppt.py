#!/usr/bin/env python3
"""node의 label이 갖는 성격과 특징 교육자료 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK    = RGBColor(0x1E, 0x1E, 0x2E)
BG_SURFACE = RGBColor(0x2A, 0x2A, 0x3C)
ACCENT1    = RGBColor(0x89, 0xB4, 0xFA)   # 밝은 파랑
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # 밝은 초록
ACCENT3    = RGBColor(0xF9, 0xE2, 0xAF)   # 밝은 노랑
ACCENT4    = RGBColor(0xF3, 0x8B, 0xA8)   # 밝은 분홍
WHITE      = RGBColor(0xCD, 0xD6, 0xF4)
GRAY       = RGBColor(0x93, 0x99, 0xB2)
CODE_BG    = RGBColor(0x31, 0x32, 0x44)
ACCENT5    = RGBColor(0xCB, 0xA6, 0xF7)   # 밝은 보라
TEAL       = RGBColor(0x94, 0xE2, 0xD5)   # 청록
PEACH      = RGBColor(0xFA, 0xB3, 0x87)   # 피치

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(shape, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Malgun Gothic", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_code_block(slide, left, top, width, height, lines):
    shape = add_shape(slide, left, top, width, height, fill_color=CODE_BG, border_color=RGBColor(0x58, 0x5B, 0x70))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT3
        p.font.name = "Consolas"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
    return shape


def make_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    set_text(tb, title, size=40, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = add_text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(1.0))
    set_text(tb2, subtitle, size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT1
    line.line.fill.background()
    return slide


def make_section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
    set_text(tb, f"#{number}", size=60, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.0))
    set_text(tb2, title, size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        tb3 = add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
        set_text(tb3, subtitle, size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def make_content_slide(title, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_SURFACE
    bar.line.fill.background()
    tb = add_text_box(slide, Inches(0.6), Inches(0.1), Inches(10), Inches(0.7))
    set_text(tb, title, size=24, color=ACCENT1, bold=True)
    if page_num:
        tb_pg = add_text_box(slide, Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.6))
        set_text(tb_pg, page_num, size=14, color=GRAY, alignment=PP_ALIGN.RIGHT, font_name="Consolas")
    return slide


# ================================================================
# 슬라이드 1: 표지
# ================================================================
make_title_slide(
    "node.label의 성격과 특징",
    "ATMS에서 노드의 label이 의미하는 것, 유지하는 불변식, 그리고 활용 방식"
)

# ================================================================
# 슬라이드 2: 목차
# ================================================================
slide = make_content_slide("목차", "2")
items = [
    ("1", "label이란 무엇인가               ", "정의와 직관적 의미"),
    ("2", "label의 구조                      ", "env 리스트, 가정 집합, 양방향 연결"),
    ("3", "3가지 핵심 불변식                  ", "최소성, 무모순성, 양방향 일관성"),
    ("4", "노드 유형별 초기 label              ", "일반 노드, 가정 노드, 모순 노드"),
    ("5", "label의 크기가 갖는 의미            ", "빈 리스트 / 빈 환경 / 다중 환경의 해석"),
    ("6", "label로 판별하는 것들               ", "true-node?, in-node?, node-consistent-with?"),
    ("7", "label이 갱신되는 경로               ", "justify → propagate → weave → update-label"),
    ("8", "label의 5가지 핵심 성격             ", "완전성, 최소성, 무모순성, 양방향, 전파 기반"),
]
y = Inches(1.2)
for num, title, desc in items:
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.72))
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(0.5), Inches(0.3))
    set_text(tb, num, size=19, color=ACCENT2, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1.5), y + Inches(0.03), Inches(5.5), Inches(0.35))
    set_text(tb2, title, size=17, color=WHITE, bold=True)
    tb3 = add_text_box(slide, Inches(1.5), y + Inches(0.37), Inches(10), Inches(0.3))
    set_text(tb3, desc, size=13, color=GRAY)
    y += Inches(0.76)


# ================================================================
# 슬라이드 3: label이란 무엇인가
# ================================================================
make_section_slide("1", "label이란 무엇인가", "\"이 노드를 참으로 만드는 최소 환경들의 리스트\"")

slide = make_content_slide("label의 정의와 직관적 의미", "3")

# 상단: 한 줄 정의
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.4), border_color=ACCENT1)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.1))
tf = set_text(tb, "정의", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "label = 이 노드를 참(true)으로 만드는 최소 환경(environment)들의 리스트", size=18, color=ACCENT3, bold=True)
add_paragraph(tf, "각 환경은 \"어떤 가정들을 참이라고 받아들이면 이 노드가 참이 되는가\"를 나타냅니다.", size=15, color=GRAY)

# 중단: 비유
card = add_shape(slide, Inches(0.5), Inches(2.9), Inches(5.8), Inches(4.2))
tb = add_text_box(slide, Inches(0.8), Inches(3.05), Inches(5.3), Inches(3.9))
tf = set_text(tb, "직관적 비유", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "node = 명제 (\"환자가 기침한다\")", size=16, color=WHITE)
add_paragraph(tf, "label = 그 명제가 참이 되는 이유 목록", size=16, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "비유: label은 법정의 \"증거 목록\"", size=17, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  피고인이 유죄인 근거:", size=15, color=WHITE)
add_paragraph(tf, "    증거1: {목격자 증언}", size=15, color=TEAL, font_name="Consolas")
add_paragraph(tf, "    증거2: {CCTV 영상, 지문}", size=15, color=TEAL, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  → 어느 하나만으로도 유죄를 입증 가능", size=15, color=GRAY)
add_paragraph(tf, "  → 각 증거는 독립적인 근거", size=15, color=GRAY)
add_paragraph(tf, '  → 이것이 label = [{목격자}, {CCTV,지문}]', size=15, color=ACCENT3, font_name="Consolas")

# 중단 오른쪽: 코드 예시
card = add_shape(slide, Inches(6.7), Inches(2.9), Inches(6.1), Inches(4.2))
tb = add_text_box(slide, Inches(7.0), Inches(3.05), Inches(5.5), Inches(3.9))
tf = set_text(tb, "ATMS에서의 구체적 예시", size=20, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "가정 노드:", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "  A1 = \"감기에 걸림\"", size=15, color=WHITE)
add_paragraph(tf, "  A2 = \"알레르기 있음\"", size=15, color=WHITE)
add_paragraph(tf, "  A3 = \"미세먼지 심함\"", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "노드 \"기침\".label = [{A1}, {A2}]", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "  의미:", size=15, color=GRAY)
add_paragraph(tf, '  "감기를 가정하면 기침이 참"  또는', size=15, color=WHITE)
add_paragraph(tf, '  "알레르기를 가정하면 기침이 참"', size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "노드 \"재채기\".label = [{A2,A3}]", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "  의미:", size=15, color=GRAY)
add_paragraph(tf, '  "알레르기 + 미세먼지를 동시에 가정해야', size=15, color=WHITE)
add_paragraph(tf, '   재채기가 참" (한 가지로는 부족)', size=15, color=WHITE)


# ================================================================
# 슬라이드 4: label의 구조
# ================================================================
make_section_slide("2", "label의 구조", "env 리스트, 가정 집합, 양방향 참조")

slide = make_content_slide("label의 데이터 구조", "4")

# 전체 구조도
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(5.5))
tf = set_text(tb, "계층 구조: node → label → env → assumptions", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, "  tms-node  \"X=6\"", size=16, color=ACCENT2, bold=True, font_name="Consolas")
add_paragraph(tf, "    │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ├─ label ─────── [env₁,  env₂,  env₃]     ← env 구조체들의 리스트", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "    │                  │       │       │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │         assumptions:  assumptions:  assumptions:", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │            (A1)      (A2 A3)     (A4)         ← 가정 노드의 정렬된 리스트", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    │                  │       │       │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │            nodes:     nodes:     nodes:", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │          (X=6 Y=3)  (X=6 Z=1) (X=6)          ← 이 환경에서 참인 노드들", size=14, color=PEACH, font_name="Consolas")
add_paragraph(tf, "    │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ├─ justs ─────── 이 노드를 결론으로 하는 정당화들", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ├─ consequences ─ 이 노드를 전건으로 사용하는 정당화들", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    └─ datum ─────── \"X=6\" (외부 데이터)", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "  양방향 참조:", size=17, color=TEAL, bold=True)
add_paragraph(tf, "    node.label → [env₁, env₂, ...]    \"이 노드가 참인 환경들\"   (정방향)", size=15, color=TEAL, font_name="Consolas")
add_paragraph(tf, "    env.nodes  → [node₁, node₂, ...]  \"이 환경에서 참인 노드들\" (역방향)", size=15, color=PEACH, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  이 양방향 연결은 DB의 인덱스와 동일한 원리입니다.", size=14, color=GRAY)


# ================================================================
# 슬라이드 5: 3가지 핵심 불변식
# ================================================================
make_section_slide("3", "3가지 핵심 불변식", "label이 항상 만족해야 하는 조건")

slide = make_content_slide("불변식 1: 최소성 (Minimality)", "5")

# 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.3), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(2.0))
tf = set_text(tb, "규칙: label 내의 어떤 두 환경도 부분집합 관계가 아니다", size=19, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "만약 {A1} ∈ label 이면, {A1, A2} ∉ label", size=17, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "원리: \"적은 가정으로 참이 되면, 더 많은 가정을 추가할 이유가 없다\"", size=16, color=WHITE)
add_paragraph(tf, "이것은 오컴의 면도날(Occam's Razor)의 형식화입니다.", size=15, color=GRAY)

# 3가지 케이스
cases = [
    ("O  올바른 label", ACCENT2,
     "label = [{A1}, {A2}]",
     "{A1}과 {A2}는 서로 부분집합이 아님\n(독립적인 두 근거)"),
    ("X  위반 사례", ACCENT4,
     "label = [{A1}, {A1,A2}]",
     "{A1} ⊂ {A1,A2} → 위반!\n{A1}만으로 충분한데\n{A1,A2}를 유지할 이유 없음"),
    ("→  수정 결과", ACCENT1,
     "label = [{A1}]",
     "update-label이 자동으로\n{A1,A2}를 제거하고\n{A1}만 유지"),
]
x = Inches(0.5)
for title, color, code, note in cases:
    card = add_shape(slide, x, Inches(3.9), Inches(3.9), Inches(3.2), border_color=color)
    tb = add_text_box(slide, x + Inches(0.25), Inches(4.05), Inches(3.4), Inches(0.4))
    set_text(tb, title, size=16, color=color, bold=True)
    add_code_block(slide, x + Inches(0.15), Inches(4.45), Inches(3.6), Inches(0.8), code.split('\n'))
    tb2 = add_text_box(slide, x + Inches(0.25), Inches(5.4), Inches(3.4), Inches(1.5))
    set_text(tb2, note, size=14, color=WHITE)
    x += Inches(4.15)


# 슬라이드 6: 불변식 2, 3
slide = make_content_slide("불변식 2, 3: 무모순성과 양방향 일관성", "6")

# 불변식 2
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=ACCENT2)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(5.5))
tf = set_text(tb, "불변식 2: 무모순성 (No Nogood)", size=19, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "규칙", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "label 내의 모든 환경은 nogood가 아니다.", size=16, color=ACCENT3)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "의미", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "nogood = 내부 모순이 있는 환경", size=15, color=WHITE)
add_paragraph(tf, '예: "감기 O"와 "감기 X"를 동시에 가정', size=15, color=WHITE)
add_paragraph(tf, "→ 모순이므로 의미 없는 조합", size=15, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "유지 방법", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "1) weave에서 nogood 필터로 사전 차단:", size=15, color=WHITE)
add_paragraph(tf, "   (unless (env-nogood? new-env) ...)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "2) 사후 발견 시 즉시 제거:", size=15, color=WHITE)
add_paragraph(tf, "   remove-env-from-labels", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "   → 해당 env를 모든 노드의 label에서 삭제", size=13, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "3) 모순 노드의 label에 env 추가 시:", size=15, color=WHITE)
add_paragraph(tf, "   → 그 env가 자동으로 nogood 등록됨", size=13, color=GRAY)

# 불변식 3
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=TEAL)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "불변식 3: 양방향 참조 일관성", size=19, color=TEAL, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "규칙", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "env ∈ node.label  ⟺  node ∈ env.nodes", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "의미", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "정방향과 역방향 참조가 항상 동기화되어야 합니다.", size=15, color=WHITE)
add_paragraph(tf, "한쪽만 수정하면 dangling reference 발생!", size=15, color=ACCENT4)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "유지 방법", size=17, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "추가 시 (update-label):", size=15, color=ACCENT2, bold=True)
add_paragraph(tf, "  (push env envs)               ; node→env", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (push node (env-nodes env))    ; env→node", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "제거 시 (update-label):", size=15, color=ACCENT4, bold=True)
add_paragraph(tf, "  (rplaca nenvs nil)             ; node→env", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (delete node (env-nodes ...))  ; env→node", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "노드 삭제 시 (remove-node):", size=15, color=GRAY, bold=True)
add_paragraph(tf, "  (dolist (env (tms-node-label node))", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (delete node (env-nodes env)))", size=13, color=ACCENT3, font_name="Consolas")


# ================================================================
# 슬라이드 7: 노드 유형별 초기 label
# ================================================================
make_section_slide("4", "노드 유형별 초기 label", "일반 노드, 가정 노드, 모순 노드의 시작 상태")

slide = make_content_slide("노드 유형별 초기 label", "7")

# 일반 노드
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.8), border_color=ACCENT1)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(3.3), Inches(5.5))
tf = set_text(tb, "일반 노드", size=22, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "초기 label", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "nil (빈 리스트)", size=18, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "의미", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "아직 참인 근거가 없음", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "예시", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "(tms-create-node", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, '  atms "X=6")', size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "X=6.label = []", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "이후 변화", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "justify-node으로 정당화가", size=14, color=WHITE)
add_paragraph(tf, "등록되면 전파를 통해", size=14, color=WHITE)
add_paragraph(tf, "환경이 추가됨", size=14, color=WHITE)

# 가정 노드
card = add_shape(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(5.8), border_color=ACCENT2)
tb = add_text_box(slide, Inches(4.85), Inches(1.35), Inches(3.5), Inches(5.5))
tf = set_text(tb, "가정 노드", size=22, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "초기 label", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, "[{자기자신}]", size=18, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "의미", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, '"나를 가정하면 나는 참"', size=15, color=WHITE)
add_paragraph(tf, "동어반복이지만 ATMS의 출발점", size=14, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "예시", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, "(tms-create-node", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, '  atms "M1정상"', size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  :assumptionp t)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "A1.label = [{A1}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "코드", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, "(push", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (create-env atms", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (list node))", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (tms-node-label node))", size=13, color=ACCENT3, font_name="Consolas")

# 모순 노드
card = add_shape(slide, Inches(8.9), Inches(1.2), Inches(3.9), Inches(5.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(9.15), Inches(1.35), Inches(3.4), Inches(5.5))
tf = set_text(tb, "모순 노드", size=22, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "초기 label", size=16, color=ACCENT5, bold=True)
add_paragraph(tf, "nil (빈 리스트)", size=18, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "의미", size=16, color=ACCENT5, bold=True)
add_paragraph(tf, "아직 모순이 발견되지 않음", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "특수 행동", size=16, color=ACCENT5, bold=True)
add_paragraph(tf, "label에 env가 추가되면", size=14, color=WHITE)
add_paragraph(tf, "그 env는 자동으로", size=14, color=WHITE)
add_paragraph(tf, "nogood로 등록됨!", size=14, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ 모순이 참인 환경은", size=14, color=GRAY)
add_paragraph(tf, "  불가능한 환경이므로", size=14, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "코드 (update 함수)", size=16, color=ACCENT5, bold=True)
add_paragraph(tf, "(when (contradictory?", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "       consequence)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (dolist (env new-envs)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (new-nogood ..)))", size=13, color=ACCENT3, font_name="Consolas")


# ================================================================
# 슬라이드 8: label의 크기가 갖는 의미
# ================================================================
make_section_slide("5", "label의 크기가 갖는 의미", "빈 리스트, 빈 환경, 단일/다중 환경의 해석")

slide = make_content_slide("label의 크기가 의미하는 것", "8")

states = [
    ("nil", "빈 리스트", ACCENT4,
     "label = []",
     "어떤 가정하에서도\n이 노드가 참이 아님",
     "\"근거 없음\"\n\n아직 정당화가\n등록되지 않았거나,\n모든 근거가 nogood로\n사라진 상태"),
    ("[{}]", "빈 환경 1개", ACCENT2,
     "label = [{}]",
     "아무 가정 없이도\n무조건 참 (true-node?)",
     "\"관측값, 공리\"\n\njustify-node에서\n전건 없이 등록된 노드\n예: 측정값 F=10"),
    ("[{A1}]", "단일 환경", ACCENT1,
     "label = [{A1}]",
     "하나의 가정 조합으로만\n참이 됨",
     "\"유일한 근거\"\n\n이 가정이 무너지면\n노드도 거짓이 됨"),
    ("[{A1},{A2}]", "다중 환경", ACCENT5,
     "label = [{A1}, {A2}]",
     "독립적 근거가\n여러 개 존재",
     "\"복수의 근거\"\n\n하나가 무너져도\n다른 근거가 유지됨\n→ 더 견고한 믿음"),
]

x = Inches(0.3)
for title, subtitle, color, code, meaning, note in states:
    card = add_shape(slide, x, Inches(1.2), Inches(3.05), Inches(5.8), border_color=color)
    tb = add_text_box(slide, x + Inches(0.15), Inches(1.35), Inches(2.75), Inches(0.45))
    set_text(tb, title, size=22, color=color, bold=True, font_name="Consolas")
    tb = add_text_box(slide, x + Inches(0.15), Inches(1.75), Inches(2.75), Inches(0.3))
    set_text(tb, subtitle, size=13, color=GRAY)

    add_code_block(slide, x + Inches(0.1), Inches(2.15), Inches(2.85), Inches(0.6), [code])

    tb = add_text_box(slide, x + Inches(0.2), Inches(2.85), Inches(2.65), Inches(1.2))
    set_text(tb, meaning, size=14, color=WHITE)

    tb = add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(2.65), Inches(2.6))
    set_text(tb, note, size=13, color=GRAY)

    x += Inches(3.2)


# ================================================================
# 슬라이드 9: label로 판별하는 것들
# ================================================================
make_section_slide("6", "label로 판별하는 것들", "true-node?, in-node?, node-consistent-with?")

slide = make_content_slide("label 기반 질의 함수들", "9")

# true-node?
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.8), border_color=ACCENT2)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(3.3), Inches(5.5))
tf = set_text(tb, "true-node?", size=20, color=ACCENT2, bold=True, font_name="Consolas")
add_paragraph(tf, "\"무조건 참인가?\"", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "검사 방법", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "label의 첫 원소가", size=15, color=WHITE)
add_paragraph(tf, "빈 환경 {} 인지 확인", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "(eq (car label)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    empty-env)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "예시", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "F=10.label = [{}]", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "→ true-node? = T", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "X=6.label = [{A1}]", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "→ true-node? = NIL", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "의미", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "가정이 필요 없는 확실한 사실", size=14, color=WHITE)
add_paragraph(tf, "(관측값, 공리 등)", size=14, color=GRAY)

# in-node?
card = add_shape(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(5.8), border_color=ACCENT3)
tb = add_text_box(slide, Inches(4.85), Inches(1.35), Inches(3.5), Inches(5.5))
tf = set_text(tb, "in-node?", size=20, color=ACCENT3, bold=True, font_name="Consolas")
add_paragraph(tf, "\"참인 상태인가?\"", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "2가지 모드", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "env 없이 호출:", size=15, color=ACCENT2, bold=True)
add_paragraph(tf, "  label이 비어있지 않은가?", size=14, color=WHITE)
add_paragraph(tf, "  (not (null label))", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 어떤 가정하에서든 참?", size=13, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "env와 함께 호출:", size=15, color=ACCENT4, bold=True)
add_paragraph(tf, "  label 중 env의", size=14, color=WHITE)
add_paragraph(tf, "  부분집합이 있는가?", size=14, color=WHITE)
add_paragraph(tf, "  (some (subset-env?", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "         le env) label)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 이 환경에서 참?", size=13, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "예: label=[{A1}], env={A1,A2}", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {A1} ⊂ {A1,A2} → T", size=13, color=ACCENT2, font_name="Consolas")

# node-consistent-with?
card = add_shape(slide, Inches(8.9), Inches(1.2), Inches(3.9), Inches(5.8), border_color=ACCENT5)
tb = add_text_box(slide, Inches(9.15), Inches(1.35), Inches(3.4), Inches(5.5))
tf = set_text(tb, "consistent-with?", size=18, color=ACCENT5, bold=True, font_name="Consolas")
add_paragraph(tf, "\"양립 가능한가?\"", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "검사 방법", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "label의 어떤 환경과", size=14, color=WHITE)
add_paragraph(tf, "env를 합쳤을 때", size=14, color=WHITE)
add_paragraph(tf, "nogood가 아닌 게 있는가?", size=14, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(some", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, " (not (env-nogood?", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "   (union-env le env)))", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, " label)", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "의미", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "이 노드가 참이면서", size=14, color=WHITE)
add_paragraph(tf, "동시에 env도 유지할 수", size=14, color=WHITE)
add_paragraph(tf, "있는 가능성이 있는가?", size=14, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "→ 모순 없이 공존 가능?", size=14, color=GRAY)


# ================================================================
# 슬라이드 10: label이 갱신되는 경로
# ================================================================
make_section_slide("7", "label이 갱신되는 경로", "justify → propagate → weave → update-label")

slide = make_content_slide("label 갱신 전파 체인", "10")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(5.5))
tf = set_text(tb, "정당화 등록에서 label 갱신까지", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, '  1. justify-node("덧셈규칙", C, [A, B])              사용자가 정당화 등록', size=15, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "     │", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  just 구조체 생성, 양방향 연결 등록", size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  2. propagate(just, nil, [{}])                       빈 환경에서 출발", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "     │", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  3. weave(nil, [{}], [A, B])                         전건들의 label 교차곱", size=15, color=TEAL, font_name="Consolas")
add_paragraph(tf, "     │", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  A.label=[{A1}], B.label=[{A2}]", size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  {} x {A1} x {A2} = {A1,A2}", size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  결과: [{A1,A2}]", size=13, color=TEAL, font_name="Consolas")
add_paragraph(tf, "     ▼", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  4. update([{A1,A2}], C, just)                       결론 노드에 전달", size=15, color=PEACH, font_name="Consolas")
add_paragraph(tf, "     │", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  5. update-label(C, [{A1,A2}])                       label에 병합 (최소성 유지)", size=15, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "     │", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     │  C.label = [] → [{A1,A2}]   (새로 추가됨)", size=13, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "     │  env-nodes({A1,A2})에 C 등록 (양방향)", size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     ▼", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  6. 반환: [{A1,A2}] (새로 추가된 것)                  후속 노드에 재귀 전파", size=15, color=ACCENT1, font_name="Consolas")
add_paragraph(tf, "     │", size=15, color=GRAY, font_name="Consolas")
add_paragraph(tf, "     └─ C를 전건으로 사용하는 정당화들에 propagate 재귀 호출", size=13, color=GRAY, font_name="Consolas")


# ================================================================
# 슬라이드 11: label의 5가지 핵심 성격
# ================================================================
make_section_slide("8", "label의 5가지 핵심 성격", "완전성, 최소성, 무모순성, 양방향, 전파 기반")

slide = make_content_slide("label의 5가지 핵심 성격 종합", "11")

properties = [
    ("#1 의미적 완전성", ACCENT2,
     "노드가 참이 되는 모든 최소 근거를 빠짐없이 담고 있음",
     "정당화가 추가될 때마다 전파 체인이 자동으로 모든 가능한 환경을 계산"),
    ("#2 최소성 (Minimality)", ACCENT4,
     "상위집합 환경을 허용하지 않아 공간/시간 효율 극대화",
     "update-label의 compare-env가 :S12/:S21/:EQ 비교로 강제"),
    ("#3 무모순성 (No Nogood)", ACCENT5,
     "내부 모순이 있는 환경은 label에 존재할 수 없음",
     "weave의 nogood 필터 + remove-env-from-labels로 사전/사후 차단"),
    ("#4 양방향 연결", TEAL,
     "node.label ↔ env.nodes가 항상 동기화되어 양방향 탐색 가능",
     "추가 시 push 쌍, 제거 시 delete+rplaca 쌍으로 일관성 유지"),
    ("#5 전파 기반 갱신", ACCENT1,
     "직접 설정이 아닌 정당화 체인을 통해 자동으로 계산됨",
     "justify-node → propagate → weave → update-label 전파 파이프라인"),
]

y = Inches(1.15)
for title, color, desc, detail in properties:
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.15), border_color=color)
    tb = add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(3.6), Inches(0.35))
    set_text(tb, title, size=17, color=color, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(4.5), y + Inches(0.05), Inches(8.1), Inches(0.4))
    set_text(tb2, desc, size=15, color=WHITE)
    tb3 = add_text_box(slide, Inches(4.5), y + Inches(0.42), Inches(8.1), Inches(0.7))
    set_text(tb3, detail, size=13, color=GRAY)
    y += Inches(1.22)

# 하단 마무리
tb = add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
tf = set_text(tb, "label은 ATMS의 핵심 데이터 — 이 5가지 성격이 전체 시스템의 정합성과 효율성을 보장합니다", size=15, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)


# ================================================================
# 저장
# ================================================================
output_path = "/home/user/bps/atms/node-label-characteristics.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
