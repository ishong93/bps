#!/usr/bin/env python3
"""update-label 함수의 프로그래밍 기법 교육자료 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK    = RGBColor(0x1E, 0x1E, 0x2E)   # 어두운 배경
BG_SURFACE = RGBColor(0x2A, 0x2A, 0x3C)   # 카드 배경
ACCENT1    = RGBColor(0x89, 0xB4, 0xFA)   # 밝은 파랑
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # 밝은 초록
ACCENT3    = RGBColor(0xF9, 0xE2, 0xAF)   # 밝은 노랑
ACCENT4    = RGBColor(0xF3, 0x8B, 0xA8)   # 밝은 분홍
WHITE      = RGBColor(0xCD, 0xD6, 0xF4)   # 밝은 흰색
GRAY       = RGBColor(0x93, 0x99, 0xB2)
CODE_BG    = RGBColor(0x31, 0x32, 0x44)   # 코드 배경

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # 모서리 둥글기
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(shape, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Malgun Gothic", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_code_block(slide, left, top, width, height, lines):
    """코드 블록을 추가"""
    shape = add_shape(slide, left, top, width, height, fill_color=CODE_BG, border_color=RGBColor(0x58, 0x5B, 0x70))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT3
        p.font.name = "Consolas"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
    return shape


def make_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # 제목
    tb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    set_text(tb, title, size=40, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)

    # 부제
    tb2 = add_text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(1.0))
    set_text(tb2, subtitle, size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 하단 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT1
    line.line.fill.background()

    return slide


def make_section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # 번호
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
    set_text(tb, f"#{number}", size=60, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")

    # 제목
    tb2 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.0))
    set_text(tb2, title, size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    if subtitle:
        tb3 = add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
        set_text(tb3, subtitle, size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    return slide


def make_content_slide(title, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # 상단 제목 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_SURFACE
    bar.line.fill.background()

    tb = add_text_box(slide, Inches(0.6), Inches(0.1), Inches(10), Inches(0.7))
    set_text(tb, title, size=24, color=ACCENT1, bold=True)

    if page_num:
        tb_pg = add_text_box(slide, Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.6))
        set_text(tb_pg, page_num, size=14, color=GRAY, alignment=PP_ALIGN.RIGHT, font_name="Consolas")

    return slide


# ================================================================
# 슬라이드 1: 표지
# ================================================================
make_title_slide(
    "update-label 함수의 프로그래밍 기법",
    "ATMS(Assumption-based Truth Maintenance System)  |  Common Lisp 구현 분석"
)

# ================================================================
# 슬라이드 2: 목차
# ================================================================
slide = make_content_slide("목차", "2")
items = [
    ("1", "Destructive Update (rplaca)    ", "리스트를 제자리에서 수정하는 파괴적 갱신"),
    ("2", "Lazy Deletion (지연 삭제)       ", "nil 마킹 후 일괄 정리 패턴"),
    ("3", "중첩 루프 + 조기 탈출           ", "do-do-cond 구조의 탐색 패턴"),
    ("4", "역참조 (Reverse Index)          ", "env → node 양방향 연결 관리"),
    ("5", "반환값을 통한 흐름 제어          ", "새로 추가된 환경만 반환하여 전파 결정"),
    ("6", "집합 최소성 유지 (Subsumption)  ", "부분집합 관계를 이용한 집합 정리"),
]
y = Inches(1.3)
for num, title, desc in items:
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85))
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(0.5), Inches(0.35))
    set_text(tb, num, size=20, color=ACCENT2, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1.5), y + Inches(0.05), Inches(5.5), Inches(0.4))
    set_text(tb2, title, size=18, color=WHITE, bold=True)
    tb3 = add_text_box(slide, Inches(1.5), y + Inches(0.45), Inches(10), Inches(0.35))
    set_text(tb3, desc, size=14, color=GRAY)
    y += Inches(0.95)

# ================================================================
# 슬라이드 3: 전체 코드 개관
# ================================================================
slide = make_content_slide("전체 코드 개관", "3")
code_lines = [
    "(defun update-label (node new-envs &aux envs)",
    "  (setq envs (tms-node-label node))",
    "  (do ((new-envs new-envs (cdr new-envs)))      ; 외부 루프",
    "      ((null new-envs))",
    "    (do ((nenvs envs (cdr nenvs)))               ; 내부 루프",
    "        ((null nenvs) (push (car new-envs) envs))",
    "      (cond",
    "        ((null (car nenvs)))           ; 삭제된 항목 건너뜀",
    "        ((null (car new-envs)))        ; 무시된 환경 건너뜀",
    '        ((case (compare-env (car new-envs) (car nenvs))',
    "           ((:EQ :S21) (rplaca new-envs nil))    ; 새 환경 버림",
    "           (:S12                                  ; 기존 환경 제거",
    "            (setf (env-nodes (car nenvs))",
    "              (delete node (env-nodes (car nenvs)) :COUNT 1))",
    "            (rplaca nenvs nil)))))))",
    "  ;; 정리 단계",
    '  (setq new-envs (delete nil new-envs :TEST #\'eq))',
    "  (dolist (new-env new-envs) (push node (env-nodes new-env)))",
    '  (setf (tms-node-label node) (delete nil envs :TEST #\'eq))',
    "  new-envs)   ; 실제로 추가된 환경만 반환",
]
add_code_block(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8), code_lines)

# 핵심 영역 주석 화살표
tb = add_text_box(slide, Inches(10.0), Inches(1.5), Inches(3), Inches(0.5))
set_text(tb, "← 초기화", size=13, color=ACCENT2, bold=True, font_name="Consolas")

tb = add_text_box(slide, Inches(10.0), Inches(4.5), Inches(3), Inches(0.5))
set_text(tb, "← 핵심 비교 로직", size=13, color=ACCENT4, bold=True, font_name="Consolas")

tb = add_text_box(slide, Inches(10.0), Inches(6.0), Inches(3), Inches(0.5))
set_text(tb, "← 정리 + 반환", size=13, color=ACCENT3, bold=True, font_name="Consolas")


# ================================================================
# 슬라이드 4: 기법 1 - Destructive Update
# ================================================================
make_section_slide("1", "Destructive Update (파괴적 갱신)", "rplaca로 리스트를 제자리에서 수정하기")

slide = make_content_slide("기법 1: Destructive Update (rplaca)", "4")

# 왼쪽: 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.3), Inches(5.0))
tf = set_text(tb, "rplaca란?", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, 'cons cell의 car 부분을 직접 교체하는 함수입니다.', size=16, color=WHITE)
add_paragraph(tf, '새 리스트를 만들지 않고 기존 리스트를 수정합니다.', size=16, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "update-label에서의 사용", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "• (rplaca new-envs nil)", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "   → 새 환경을 nil로 마킹 (\"삭제 표시\")", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "• (rplaca nenvs nil)", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "   → 기존 환경을 nil로 마킹", size=15, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "왜 사용하는가?", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "순회 중인 리스트를 안전하게 \"삭제 예약\"할 수 있음", size=15, color=WHITE)
add_paragraph(tf, "O(1) 시간으로 즉시 표시 가능", size=15, color=WHITE)

# 오른쪽: 코드 비교
add_code_block(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.2), [
    ";; Destructive (실제 코드)",
    "(rplaca new-envs nil)  ; 제자리 수정",
    "",
    ";; 결과:",
    ";; ({A1,A2} {A3}) → (nil {A3})",
    ";; 원래 리스트가 직접 변경됨!"
])

add_code_block(slide, Inches(6.7), Inches(3.7), Inches(6.1), Inches(2.0), [
    ";; Non-destructive (대안)",
    "(setq new-envs",
    "  (remove (car new-envs) new-envs))",
    "",
    ";; 새 리스트를 생성 → 메모리 할당 필요"
])

tb = add_text_box(slide, Inches(6.7), Inches(5.9), Inches(6.1), Inches(0.8))
tf = set_text(tb, "핵심: 순회 중 구조를 직접 변경하되, nil 마킹으로 안전성 확보", size=14, color=ACCENT4, bold=True)


# ================================================================
# 슬라이드 5: 기법 2 - Lazy Deletion
# ================================================================
make_section_slide("2", "Lazy Deletion (지연 삭제)", "\"지금 표시하고, 나중에 정리한다\"")

slide = make_content_slide("기법 2: Lazy Deletion (지연 삭제)", "5")

# 상단: 2단계 프로세스
card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.6), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(2.3))
tf = set_text(tb, "Phase 1: Mark (표시 단계)", size=20, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "순회 도중 삭제할 항목을 nil로 마킹", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(rplaca new-envs nil)  ; 새 환경 삭제 표시", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "(rplaca nenvs nil)     ; 기존 환경 삭제 표시", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "리스트 구조 자체는 유지 → 순회가 안전하게 계속됨", size=15, color=GRAY)

card2 = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.6), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(2.3))
tf = set_text(tb, "Phase 2: Sweep (정리 단계)", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "순회가 끝난 후 nil 항목 일괄 제거", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(delete nil new-envs :TEST #'eq)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "(delete nil envs :TEST #'eq)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "깨끗한 최종 상태 확보", size=15, color=GRAY)

# 하단: 시각화
card3 = add_shape(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8))
tb = add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.8), Inches(2.5))
tf = set_text(tb, "동작 시각화", size=18, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "초기 상태:    [{A1,A2},  {A3}]    ×    [{A1}]", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "                     기존 label              새 환경", size=13, color=GRAY, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "Mark 후:     [{A1,A2},  {A3}]    ×    [{A1}]       compare: {A1} ⊂ {A1,A2} → :S12", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "              → [nil,    {A3}]    ×    [{A1}]       기존 {A1,A2}를 nil로 마킹", size=16, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "Sweep 후:    [{A3}]               ×    [{A1}]       nil 제거 완료", size=16, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "최종 label:  [{A3}, {A1}]                           두 환경 공존", size=16, color=ACCENT1, font_name="Consolas")


# ================================================================
# 슬라이드 6: 기법 3 - 중첩 루프 + 조기 탈출
# ================================================================
make_section_slide("3", "중첩 루프 + 조기 탈출", "do-do-cond 구조의 우아한 탐색 패턴")

slide = make_content_slide("기법 3: 중첩 do 루프와 조기 탈출", "6")

# 왼쪽: 구조 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5))
tf = set_text(tb, "패턴 구조", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "외부 do 루프: new-envs 순회", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "  각 새 환경을 하나씩 처리", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  내부 do 루프: envs (기존 label) 순회", size=16, color=ACCENT3, bold=True)
add_paragraph(tf, "    새 환경과 기존 환경을 하나씩 비교", size=15, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "조기 탈출 메커니즘", size=20, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "내부 루프의 종료 조건:", size=16, color=WHITE)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "1. (null nenvs) → 모든 기존 환경과 비교 완료", size=15, color=WHITE)
add_paragraph(tf, "   → push로 새 환경을 label에 추가", size=14, color=ACCENT2)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "2. (null (car new-envs)) → 이미 nil로 마킹됨", size=15, color=WHITE)
add_paragraph(tf, "   → cond 절이 아무것도 안 함 → 자연 종료", size=14, color=ACCENT4)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, '3. rplaca 실행 → 마킹 후 cond의 case가 값 반환', size=15, color=WHITE)
add_paragraph(tf, "   → 내부 루프의 다음 반복으로 이동", size=14, color=ACCENT3)

# 오른쪽: 의사코드
add_code_block(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.8), [
    "# 의사코드 (pseudocode)",
    "",
    "for new_env in new_envs:",
    "  found_superset = False",
    "",
    "  for existing_env in envs:",
    "    if existing_env is nil:",
    "      continue    # 건너뜀",
    "    if new_env is nil:",
    "      break       # 이미 처리됨",
    "",
    "    match compare(new_env, existing_env):",
    "      case EQ, S21:    # 새 것 ⊇ 기존",
    "        new_env = nil  # 새 것 버림",
    "        break",
    "      case S12:        # 새 것 ⊂ 기존",
    "        existing_env = nil  # 기존 제거",
    "        # 계속 순회 (더 제거할 것 있을 수도)",
    "",
    "  if new_env is not nil:",
    "    envs.add(new_env)  # 살아남으면 추가",
])


# ================================================================
# 슬라이드 7: 기법 4 - 역참조 (Reverse Index)
# ================================================================
make_section_slide("4", "역참조 (Reverse Index)", "env ↔ node 양방향 연결 관리")

slide = make_content_slide("기법 4: 역참조 (Reverse Index)", "7")

# 상단: 데이터 구조 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.4))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(2.2))
tf = set_text(tb, "두 방향의 참조", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "node.label → [env1, env2, ...]         \"이 노드가 참인 환경들\"", size=16, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "env.nodes  → [node1, node2, ...]       \"이 환경에서 참인 노드들\"", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "양방향으로 연결하면 어느 쪽에서든 빠르게 탐색할 수 있습니다.", size=16, color=WHITE)
add_paragraph(tf, '이는 관계형 DB의 인덱스, 그래프의 adjacency list와 같은 원리입니다.', size=15, color=GRAY)

# 하단 왼쪽: 추가 시
card2 = add_shape(slide, Inches(0.5), Inches(4.0), Inches(5.8), Inches(3.0), border_color=ACCENT2)
tb = add_text_box(slide, Inches(0.8), Inches(4.15), Inches(5.3), Inches(2.7))
tf = set_text(tb, "환경 추가 시 (양방향 등록)", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, ";; node → env 방향", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "(push (car new-envs) envs)  ; label에 추가", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, ";; env → node 방향", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "(push node (env-nodes new-env))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "→ 둘 다 등록해야 일관성 유지!", size=15, color=WHITE)

# 하단 오른쪽: 제거 시
card3 = add_shape(slide, Inches(6.7), Inches(4.0), Inches(6.1), Inches(3.0), border_color=ACCENT4)
tb = add_text_box(slide, Inches(7.0), Inches(4.15), Inches(5.5), Inches(2.7))
tf = set_text(tb, "환경 제거 시 (양방향 해제)", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, ";; node → env 방향", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "(rplaca nenvs nil)  ; label에서 nil 마킹", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, ";; env → node 방향", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "(delete node (env-nodes (car nenvs)))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "→ 양쪽 모두 정리해야 dangling reference 방지!", size=15, color=WHITE)


# ================================================================
# 슬라이드 8: 기법 5 - 반환값을 통한 흐름 제어
# ================================================================
make_section_slide("5", "반환값을 통한 흐름 제어", "\"실제로 변한 것만 반환하여 전파 범위를 최소화\"")

slide = make_content_slide("기법 5: 반환값을 통한 흐름 제어", "8")

# 왼쪽: 흐름도
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.3), Inches(5.5))
tf = set_text(tb, "호출 흐름", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "propagate", size=18, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  └─ weave     → 교차곱 환경들", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  └─ update    → update-label 호출", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "       └─ update-label", size=16, color=ACCENT2, bold=True, font_name="Consolas")
add_paragraph(tf, "            반환: 실제 추가된 환경만!", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=14, color=WHITE)
add_paragraph(tf, "반환값의 의미", size=20, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "• 빈 리스트 () 반환", size=16, color=WHITE)
add_paragraph(tf, "  → 새로 추가된 것이 없음", size=15, color=GRAY)
add_paragraph(tf, "  → update가 전파를 멈춤 (불필요한 재귀 방지)", size=15, color=ACCENT4)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "• 환경 리스트 반환", size=16, color=WHITE)
add_paragraph(tf, "  → 실제로 label이 변경됨", size=15, color=GRAY)
add_paragraph(tf, "  → 후속 정당화들에 재귀 전파 필요", size=15, color=ACCENT2)

# 오른쪽: 코드
add_code_block(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.5), [
    ";; update-label의 반환",
    "(setq new-envs",
    "  (delete nil new-envs :TEST #'eq))",
    "...",
    "new-envs)  ; 살아남은 새 환경만 반환",
])

add_code_block(slide, Inches(6.7), Inches(4.0), Inches(6.1), Inches(3.0), [
    ";; update 함수에서의 활용",
    "(defun update (new-envs consequence just)",
    "  ...",
    "  ;; update-label이 빈 리스트 반환하면",
    "  ;; new-envs가 nil → 아래 루프 실행 안 됨",
    "  (setq new-envs",
    "    (update-label consequence new-envs))",
    "  (dolist (just (tms-node-consequences",
    "                  consequence))",
    "    ;; new-envs가 nil이면 전파 안 함!",
    "    (propagate just consequence new-envs)))",
])


# ================================================================
# 슬라이드 9: 기법 6 - 집합 최소성 유지 (Subsumption)
# ================================================================
make_section_slide("6", "집합 최소성 유지 (Subsumption)", "부분집합 관계를 이용한 효율적 집합 정리")

slide = make_content_slide("기법 6: Subsumption (포섭 관계 활용)", "9")

# 상단: 원리 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.0))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.8))
tf = set_text(tb, "최소성(Minimality) 원칙", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, 'label에는 항상 최소 환경만 유지합니다. {A1}이 있으면 {A1,A2}는 불필요합니다.', size=16, color=WHITE)
add_paragraph(tf, '"적은 가정으로 참이 되면, 더 많은 가정을 추가할 필요가 없다"는 원리입니다.', size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "이는 compare-env 함수로 두 환경의 부분집합 관계를 판별하여 구현합니다.", size=15, color=GRAY)

# 하단: 3가지 케이스
cases_data = [
    ("Case 1: :EQ / :S21", ACCENT4,
     "새 환경 ⊇ 기존 환경",
     "label = [{A1}]\nnew  = {A1,A2}\n\n→ {A1,A2} 버림\n  (이미 {A1}으로 충분)",
     "새 환경이 기존보다\n크거나 같으면 불필요"),
    ("Case 2: :S12", ACCENT2,
     "새 환경 ⊂ 기존 환경",
     "label = [{A1,A2}]\nnew  = {A1}\n\n→ {A1,A2} 제거\n→ {A1} 추가",
     "더 작은 환경이\n발견되면 교체"),
    ("Case 3: nil", ACCENT1,
     "비교 불가 (독립적)",
     "label = [{A1}]\nnew  = {A2}\n\n→ 둘 다 유지\n  label = [{A1},{A2}]",
     "독립적인 환경은\n모두 유지"),
]

x = Inches(0.5)
for title, color, subtitle, example, note in cases_data:
    card = add_shape(slide, x, Inches(3.6), Inches(3.9), Inches(3.5), border_color=color)
    tb = add_text_box(slide, x + Inches(0.25), Inches(3.75), Inches(3.4), Inches(0.4))
    set_text(tb, title, size=16, color=color, bold=True)
    tb2 = add_text_box(slide, x + Inches(0.25), Inches(4.15), Inches(3.4), Inches(0.35))
    set_text(tb2, subtitle, size=13, color=GRAY)

    # 예시 코드
    example_lines = example.split('\n')
    add_code_block(slide, x + Inches(0.15), Inches(4.55), Inches(3.6), Inches(1.7), example_lines)

    tb3 = add_text_box(slide, x + Inches(0.25), Inches(6.35), Inches(3.4), Inches(0.6))
    set_text(tb3, note, size=12, color=WHITE)

    x += Inches(4.15)


# ================================================================
# 슬라이드 10: 종합 - 기법 간 관계
# ================================================================
slide = make_content_slide("종합: 기법들의 상호 관계", "10")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.5))
tf = set_text(tb, "update-label에 적용된 6가지 기법의 관계도", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=14, color=WHITE)

add_paragraph(tf, "    ┌────────────────────────────────────────────────────────────────┐", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #6 Subsumption: 최소성 원칙이 전체 알고리즘의 목표를 결정     │", size=14, color=ACCENT1, font_name="Consolas")
add_paragraph(tf, "    └─────────────────────────┬──────────────────────────────────────┘", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "                              │ 구현", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "                              ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ┌──────────────────────────────────────────────────────────┐", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #3 중첩 루프: 모든 (새 환경, 기존 환경) 쌍을 비교      │", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    │       │                          │                       │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │       ▼                          ▼                       │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #1 rplaca: 즉시 마킹     #4 역참조: env↔node 동기화    │", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "    │       │                                                  │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │       ▼                                                  │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #2 Lazy Delete: 순회 후 nil 일괄 정리                   │", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "    └──────────────────────────────────────────────────────────┘", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "                              │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "                              ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ┌──────────────────────────────────────────────────────────┐", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #5 반환값 제어: 실제 변경분만 반환 → 전파 범위 최소화   │", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    └──────────────────────────────────────────────────────────┘", size=14, color=GRAY, font_name="Consolas")


# ================================================================
# 슬라이드 11: 핵심 요약
# ================================================================
slide = make_content_slide("핵심 요약", "11")

summaries = [
    ("#1 Destructive Update", "rplaca로 제자리 수정 → O(1) 마킹, 메모리 절약", ACCENT4),
    ("#2 Lazy Deletion", "Mark-Sweep 패턴: 순회 중 안전하게 삭제 처리", ACCENT2),
    ("#3 중첩 루프 + 조기 탈출", "do-do-cond 구조로 전수 비교, nil 체크로 불필요한 반복 스킵", ACCENT3),
    ("#4 역참조 (Reverse Index)", "env ↔ node 양방향 연결로 빠른 탐색 지원", ACCENT1),
    ("#5 반환값 흐름 제어", "변경된 것만 반환하여 불필요한 재귀 전파 차단", ACCENT3),
    ("#6 Subsumption", "부분집합 관계로 최소 환경만 유지 → 공간/시간 효율 극대화", ACCENT1),
]

y = Inches(1.2)
for title, desc, color in summaries:
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.85), border_color=color)
    tb = add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(4.0), Inches(0.35))
    set_text(tb, title, size=17, color=color, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(4.8), y + Inches(0.12), Inches(7.8), Inches(0.65))
    set_text(tb2, desc, size=16, color=WHITE)
    y += Inches(0.95)

# 하단 메시지
tb = add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8))
tf = set_text(tb, "이 기법들은 1980-90년대 AI 추론 시스템에서 발전한 실전 패턴으로,", size=15, color=GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "현대 프로그래밍의 GC, 인덱싱, 이벤트 전파 등의 기초가 되었습니다.", size=15, color=GRAY, alignment=PP_ALIGN.CENTER)


# ================================================================
# 저장
# ================================================================
output_path = "/home/user/bps/atms/update-label-techniques.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
