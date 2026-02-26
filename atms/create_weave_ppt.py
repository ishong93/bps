#!/usr/bin/env python3
"""weave 함수의 프로그래밍 기법 교육자료 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK    = RGBColor(0x1E, 0x1E, 0x2E)   # 어두운 배경
BG_SURFACE = RGBColor(0x2A, 0x2A, 0x3C)   # 카드 배경
ACCENT1    = RGBColor(0x89, 0xB4, 0xFA)   # 밝은 파랑
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # 밝은 초록
ACCENT3    = RGBColor(0xF9, 0xE2, 0xAF)   # 밝은 노랑
ACCENT4    = RGBColor(0xF3, 0x8B, 0xA8)   # 밝은 분홍
WHITE      = RGBColor(0xCD, 0xD6, 0xF4)   # 밝은 흰색
GRAY       = RGBColor(0x93, 0x99, 0xB2)
CODE_BG    = RGBColor(0x31, 0x32, 0x44)   # 코드 배경
ACCENT5    = RGBColor(0xCB, 0xA6, 0xF7)   # 밝은 보라

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(shape, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Malgun Gothic", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_code_block(slide, left, top, width, height, lines):
    """코드 블록을 추가"""
    shape = add_shape(slide, left, top, width, height, fill_color=CODE_BG, border_color=RGBColor(0x58, 0x5B, 0x70))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT3
        p.font.name = "Consolas"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
    return shape


def make_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    set_text(tb, title, size=40, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = add_text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(1.0))
    set_text(tb2, subtitle, size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT1
    line.line.fill.background()
    return slide


def make_section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
    set_text(tb, f"#{number}", size=60, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.0))
    set_text(tb2, title, size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        tb3 = add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
        set_text(tb3, subtitle, size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def make_content_slide(title, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_SURFACE
    bar.line.fill.background()
    tb = add_text_box(slide, Inches(0.6), Inches(0.1), Inches(10), Inches(0.7))
    set_text(tb, title, size=24, color=ACCENT1, bold=True)
    if page_num:
        tb_pg = add_text_box(slide, Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.6))
        set_text(tb_pg, page_num, size=14, color=GRAY, alignment=PP_ALIGN.RIGHT, font_name="Consolas")
    return slide


# ================================================================
# 슬라이드 1: 표지
# ================================================================
make_title_slide(
    "weave 함수의 프로그래밍 기법",
    "ATMS 환경 교차곱(Cross-Product) 알고리즘의 구현 기법 분석"
)

# ================================================================
# 슬라이드 2: 목차
# ================================================================
slide = make_content_slide("목차", "2")
items = [
    ("1", "Defensive Copy (방어적 복사)       ", "copy-list로 원본 데이터를 보호하는 패턴"),
    ("2", "3중 루프 교차곱 (Cross-Product)    ", "dolist-dolist-dolist 구조의 조합 생성"),
    ("3", "점진적 축적 (Incremental Fold)     ", "각 전건 처리 후 envs를 점진적으로 교체"),
    ("4", "다단계 필터링 파이프라인            ", "nogood 필터 → subsumption 필터의 2단 검증"),
    ("5", "Early Return (조기 반환)           ", "return-from으로 불필요한 계산 즉시 차단"),
    ("6", "Lazy Deletion + rplaca             ", "순회 중 마킹, 순회 후 정리하는 패턴"),
]
y = Inches(1.3)
for num, title, desc in items:
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85))
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(0.5), Inches(0.35))
    set_text(tb, num, size=20, color=ACCENT2, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1.5), y + Inches(0.05), Inches(5.5), Inches(0.4))
    set_text(tb2, title, size=18, color=WHITE, bold=True)
    tb3 = add_text_box(slide, Inches(1.5), y + Inches(0.45), Inches(10), Inches(0.35))
    set_text(tb3, desc, size=14, color=GRAY)
    y += Inches(0.95)

# ================================================================
# 슬라이드 3: 전체 코드 개관
# ================================================================
slide = make_content_slide("전체 코드 개관", "3")
code_lines = [
    "(defun weave (antecedent envs antecedents &aux new-envs new-env)",
    "  (setq envs (copy-list envs))                     ; #1 방어적 복사",
    "  (dolist (node antecedents)                        ; #2 외부: 전건 순회",
    "    (unless (eq node antecedent)                    ;    이미 반영된 것 건너뜀",
    "      (setq new-envs nil)",
    "      (dolist (env envs)                            ; #2 중간: 현재 envs 순회",
    '        (if env',
    "          (dolist (node-env (tms-node-label node))  ; #2 내부: 노드의 label 순회",
    "            (setq new-env (union-env env node-env)) ;    합집합 계산",
    "            (unless (env-nogood? new-env)           ; #4 필터1: nogood 제거",
    "              (do ((nnew-envs new-envs              ; #4 필터2: 최소성 검사",
    "                    (cdr nnew-envs)))",
    "                  ((null nnew-envs)",
    "                   (push new-env new-envs))         ;    살아남으면 추가",
    "                (when (car nnew-envs)",
    "                  (case (compare-env new-env (car nnew-envs))",
    "                    ((:EQ :S21) (return nil))       ; #5 새것 탈락 → 즉시 탈출",
    "                    (:S12                           ; #6 기존 것 제거 → rplaca",
    "                     (rplaca nnew-envs nil)))))))))",
    "      (setq envs                                   ; #3 점진적 축적",
    "        (delete nil new-envs :TEST #'eq))           ; #6 nil 정리",
    "      (unless envs (return-from weave nil))))       ; #5 빈 결과 → 조기 반환",
    "  envs)",
]
add_code_block(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0), code_lines)

# 우측 기법 번호 주석
annotations = [
    (Inches(1.4), "#1 방어적 복사", ACCENT2),
    (Inches(2.5), "#2 3중 루프", ACCENT5),
    (Inches(4.2), "#4 다단계 필터", ACCENT4),
    (Inches(5.6), "#3 점진적 축적", ACCENT3),
    (Inches(6.2), "#5 조기 반환", ACCENT1),
]


# ================================================================
# 슬라이드 4: 기법 1 - Defensive Copy
# ================================================================
make_section_slide("1", "Defensive Copy (방어적 복사)", "copy-list로 원본 데이터를 보호하기")

slide = make_content_slide("기법 1: Defensive Copy (방어적 복사)", "4")

# 왼쪽: 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.3), Inches(5.5))
tf = set_text(tb, "무엇을 복사하는가?", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(setq envs (copy-list envs))", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "함수에 전달된 envs 리스트를 복사합니다.", size=16, color=WHITE)
add_paragraph(tf, "weave 내부에서 envs를 변형(delete, setq)하기 때문에", size=16, color=WHITE)
add_paragraph(tf, "원본 리스트가 손상되지 않도록 보호합니다.", size=16, color=WHITE)
add_paragraph(tf, "", size=12, color=WHITE)
add_paragraph(tf, "왜 필요한가?", size=20, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "weave는 호출자(propagate)로부터 envs를 받습니다.", size=16, color=WHITE)
add_paragraph(tf, "복사 없이 직접 수정하면:", size=16, color=WHITE)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "  1. 호출자의 데이터가 예기치 않게 변경됨", size=15, color=ACCENT4)
add_paragraph(tf, "  2. 같은 envs를 다른 곳에서도 사용 중이면 버그 발생", size=15, color=ACCENT4)
add_paragraph(tf, "  3. 디버깅이 매우 어려워짐 (원격 부작용)", size=15, color=ACCENT4)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "원칙", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, '"입력을 변형할 것이라면, 먼저 복사하라"', size=16, color=ACCENT3)

# 오른쪽: 비교
add_code_block(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.5), [
    ";; weave: 복사 후 수정 (안전)",
    "(defun weave (antecedent envs antecedents)",
    "  (setq envs (copy-list envs))  ; 복사!",
    "  ;; 이후 envs를 마음껏 변형 가능",
    "  (setq envs (delete nil ...))  ; 안전",
    "  ...)",
])

add_code_block(slide, Inches(6.7), Inches(4.1), Inches(6.1), Inches(2.5), [
    ";; update-label: 복사 없이 직접 수정",
    "(defun update-label (node new-envs &aux envs)",
    "  (setq envs (tms-node-label node))",
    "  ;; 내부 데이터를 직접 수정하고",
    "  ;; 최종 결과를 다시 setf로 저장",
    "  (setf (tms-node-label node) ...))",
])

tb = add_text_box(slide, Inches(6.7), Inches(6.8), Inches(6.1), Inches(0.6))
tf = set_text(tb, "비교: update-label은 구조체 필드를 직접 갱신하므로 복사 불필요", size=13, color=GRAY)
add_paragraph(tf, "weave는 호출자의 리스트를 받으므로 복사 필수", size=13, color=ACCENT2)


# ================================================================
# 슬라이드 5: 기법 2 - 3중 루프 교차곱
# ================================================================
make_section_slide("2", "3중 루프 교차곱 (Cross-Product)", "Cartesian Product를 생성하는 중첩 dolist 패턴")

slide = make_content_slide("기법 2: 3중 루프 교차곱", "5")

# 왼쪽: 루프 구조 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5))
tf = set_text(tb, "3중 dolist 구조", size=20, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "외부 (dolist node antecedents)", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, '  "각 전건 노드를 하나씩 처리"', size=14, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  중간 (dolist env envs)", size=16, color=ACCENT5, bold=True)
add_paragraph(tf, '    "현재까지 축적된 환경들을 하나씩 꺼냄"', size=14, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "    내부 (dolist node-env (label node))", size=16, color=ACCENT3, bold=True)
add_paragraph(tf, '      "이 노드의 가능한 환경들과 결합"', size=14, color=GRAY)
add_paragraph(tf, "", size=12, color=WHITE)
add_paragraph(tf, "수학적 의미", size=20, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "envs = envs  x  label(node)", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "     = { e1 U e2 | e1 in envs, e2 in label(node) }", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "이것은 집합의 Cartesian Product(데카르트 곱)에", size=15, color=WHITE)
add_paragraph(tf, "합집합(union) 연산을 결합한 형태입니다.", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "곱셈처럼 결합하되, 원소는 환경(가정 집합)이므로", size=15, color=GRAY)
add_paragraph(tf, "개별 원소를 곱하는 대신 합집합으로 병합합니다.", size=15, color=GRAY)

# 오른쪽: 시각화
card2 = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.8))
tb = add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5.3), Inches(5.5))
tf = set_text(tb, "동작 예시", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "A.label = [{A1}, {A2}]", size=15, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "B.label = [{A3}]", size=15, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "C.label = [{A4}, {A5}]", size=15, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "초기: envs = [{}]", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "A 처리:", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "  {} U {A1} = {A1}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {} U {A2} = {A2}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  envs = [{A1}, {A2}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "B 처리:", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "  {A1} U {A3} = {A1,A3}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {A2} U {A3} = {A2,A3}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  envs = [{A1,A3}, {A2,A3}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "C 처리:", size=16, color=ACCENT1, bold=True)
add_paragraph(tf, "  {A1,A3} U {A4} = {A1,A3,A4}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {A1,A3} U {A5} = {A1,A3,A5}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {A2,A3} U {A4} = {A2,A3,A4}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {A2,A3} U {A5} = {A2,A3,A5}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  envs = [{A1,A3,A4}, {A1,A3,A5},", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "          {A2,A3,A4}, {A2,A3,A5}]", size=14, color=WHITE, font_name="Consolas")


# ================================================================
# 슬라이드 6: 기법 3 - 점진적 축적
# ================================================================
make_section_slide("3", "점진적 축적 (Incremental Fold)", "각 전건 처리 후 envs를 교체하는 누적 패턴")

slide = make_content_slide("기법 3: 점진적 축적 (Incremental Fold)", "6")

# 상단: 패턴 설명
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(2.3))
tf = set_text(tb, "Fold/Reduce 패턴", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "각 전건(antecedent)을 처리할 때마다 envs가 새로운 값으로 교체됩니다.", size=16, color=WHITE)
add_paragraph(tf, "이것은 함수형 프로그래밍의 fold(reduce)와 동일한 패턴입니다.", size=16, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "envs₀ = 초기값      →  (A와 교차곱)  →  envs₁  →  (B와 교차곱)  →  envs₂  →  ...  →  최종 envs", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "마치 숫자의 곱셈 누적과 같습니다: result = 1 * a * b * c  →  result = {} x A.label x B.label x C.label", size=14, color=GRAY)

# 하단 왼쪽: 실제 코드
add_code_block(slide, Inches(0.5), Inches(4.1), Inches(5.8), Inches(3.2), [
    ";; weave에서의 점진적 축적",
    "",
    "(dolist (node antecedents)     ; 각 전건에 대해",
    "  (setq new-envs nil)          ; 결과 초기화",
    "  (dolist (env envs)           ; 현재 envs와",
    "    (dolist (node-env ...)     ; 노드의 label을",
    "      ...                      ;   결합하여",
    "      (push new-env new-envs)));  new-envs에 축적",
    "",
    "  ;; ★ 핵심: envs를 new-envs로 교체 ★",
    "  (setq envs",
    "    (delete nil new-envs ...)) ; envs ← new-envs",
])

# 하단 오른쪽: 함수형 대비
add_code_block(slide, Inches(6.7), Inches(4.1), Inches(6.1), Inches(1.5), [
    "# 함수형 프로그래밍에서의 동일 패턴 (Python)",
    "from functools import reduce",
    "result = reduce(cross_product, antecedents, [{}])",
])

card2 = add_shape(slide, Inches(6.7), Inches(5.9), Inches(6.1), Inches(1.4), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(6.05), Inches(5.5), Inches(1.2))
tf = set_text(tb, "이 패턴의 장점", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "각 단계의 중간 결과를 바로 다음 단계의 입력으로 사용", size=14, color=WHITE)
add_paragraph(tf, "전건이 N개여도 동일한 2중 루프 구조로 처리 가능", size=14, color=WHITE)
add_paragraph(tf, "중간 단계에서 환경이 비면 즉시 종료 가능 (Early Return)", size=14, color=WHITE)


# ================================================================
# 슬라이드 7: 기법 4 - 다단계 필터링 파이프라인
# ================================================================
make_section_slide("4", "다단계 필터링 파이프라인", "nogood 필터 → subsumption 필터의 2단 검증")

slide = make_content_slide("기법 4: 다단계 필터링 파이프라인", "7")

# 상단: 전체 흐름
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.6))
tf = set_text(tb, "환경 생성 → 검증 파이프라인", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "union-env(env, node-env)  ──→  nogood 검사  ──→  subsumption 검사  ──→  new-envs에 추가", size=16, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "      생성                    1차 필터            2차 필터               최종 축적", size=13, color=GRAY, font_name="Consolas")

# 하단 왼쪽: 1차 필터
card2 = add_shape(slide, Inches(0.5), Inches(3.4), Inches(5.8), Inches(3.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.8), Inches(3.55), Inches(5.3), Inches(3.5))
tf = set_text(tb, "1차 필터: Nogood 검사", size=19, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(unless (env-nogood? new-env) ...)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "환경이 모순(nogood)인지 확인합니다.", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "예: {A1, A2}가 서로 모순이면", size=15, color=WHITE)
add_paragraph(tf, "  이 환경은 절대 참이 될 수 없으므로 즉시 버립니다.", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "특징:", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "  O(1) 검사 (환경에 nogood 플래그가 있음)", size=15, color=WHITE)
add_paragraph(tf, "  가장 먼저 실행 → 비용이 큰 subsumption 검사 전에", size=15, color=WHITE)
add_paragraph(tf, "  불가능한 환경을 빠르게 제거하여 성능 최적화", size=15, color=WHITE)

# 하단 오른쪽: 2차 필터
card3 = add_shape(slide, Inches(6.7), Inches(3.4), Inches(6.1), Inches(3.8), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(3.55), Inches(5.5), Inches(3.5))
tf = set_text(tb, "2차 필터: Subsumption 검사", size=19, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(do ((nnew-envs new-envs ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (case (compare-env new-env (car nnew-envs))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "이미 축적된 new-envs와 크기 비교:", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, ":EQ / :S21 → 새것 >= 기존것", size=15, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "  새 환경이 불필요 → 탈락 (return nil)", size=14, color=WHITE)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, ":S12       → 새것 < 기존것", size=15, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "  기존 환경이 불필요 → rplaca nil로 마킹", size=14, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "update-label과 동일한 최소성 유지 패턴!", size=15, color=ACCENT1, bold=True)


# ================================================================
# 슬라이드 8: 기법 5 - Early Return
# ================================================================
make_section_slide("5", "Early Return (조기 반환)", "\"곱셈에서 0이 나오면 나머지는 계산할 필요 없다\"")

slide = make_content_slide("기법 5: Early Return (조기 반환)", "8")

# 왼쪽: 두 가지 조기 반환
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5))
tf = set_text(tb, "weave에는 2종류의 조기 반환이 있다", size=18, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=12, color=WHITE)

add_paragraph(tf, "Return Type A: return-from weave", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(unless envs", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (return-from weave nil))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "어떤 전건의 교차곱 결과가 빈 리스트이면", size=15, color=WHITE)
add_paragraph(tf, "전체 weave를 즉시 종료하고 nil 반환합니다.", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "이유: 빈 집합과의 교차곱은 항상 빈 집합", size=14, color=GRAY)
add_paragraph(tf, "  {} x Anything = {} (곱셈의 0과 동일)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=14, color=WHITE)

add_paragraph(tf, "Return Type B: return nil (do 루프)", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "((:EQ :S21) (return nil))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "subsumption 검사에서 새 환경이 기존보다", size=15, color=WHITE)
add_paragraph(tf, "크거나 같으면 내부 do 루프를 즉시 탈출합니다.", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "이유: 이미 더 작은 환경이 있으므로 불필요", size=14, color=GRAY)

# 오른쪽: 시각화
card2 = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(2.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(7.3), Inches(1.35), Inches(5.3), Inches(2.5))
tf = set_text(tb, "Type A: 빈 결과 → 전체 함수 탈출", size=17, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "A.label = [{A1}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "B.label = []          ← label이 비어있음!", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "C.label = [{A4}]", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "A 처리: envs = [{A1}]", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "B 처리: {A1} x [] = []  → envs = nil!", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "  → return-from weave nil", size=14, color=ACCENT4, font_name="Consolas", bold=True)
add_paragraph(tf, "C는 처리하지 않음 (불필요)", size=14, color=GRAY)

card3 = add_shape(slide, Inches(7.0), Inches(4.3), Inches(5.8), Inches(2.7), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.3), Inches(4.45), Inches(5.3), Inches(2.4))
tf = set_text(tb, "Type B: 상위집합 → 내부 루프 탈출", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "new-envs = [{A1}]    (이미 축적됨)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "new-env  = {A1,A2}   (새로 생성됨)", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "compare: {A1,A2} vs {A1} → :S21", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  {A1,A2} ⊇ {A1} → 새것이 더 큼 → 탈락!", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "  → (return nil) → do 루프만 탈출", size=14, color=ACCENT2, font_name="Consolas", bold=True)
add_paragraph(tf, "  나머지 node-env와의 결합은 계속됨", size=14, color=GRAY)


# ================================================================
# 슬라이드 9: 기법 6 - Lazy Deletion + rplaca
# ================================================================
make_section_slide("6", "Lazy Deletion + rplaca", "순회 중 마킹, 순회 후 정리하는 패턴 (update-label과 공유)")

slide = make_content_slide("기법 6: Lazy Deletion + rplaca", "9")

# 상단: update-label과의 공통점
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.2))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.9))
tf = set_text(tb, "update-label과 동일한 Mark-Sweep 패턴", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "weave와 update-label은 동일한 Lazy Deletion 기법을 공유합니다.", size=16, color=WHITE)
add_paragraph(tf, "순회 중인 리스트를 직접 수정하면 순회가 깨질 수 있기 때문에,", size=16, color=WHITE)
add_paragraph(tf, "먼저 nil로 마킹(mark)하고 순회 후 한꺼번에 정리(sweep)합니다.", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "이 패턴은 BPS 코드베이스 전체에서 반복적으로 사용되는 관용구(idiom)입니다.", size=15, color=GRAY)

# 하단 왼쪽: weave에서의 사용
card2 = add_shape(slide, Inches(0.5), Inches(3.8), Inches(5.8), Inches(3.4), border_color=ACCENT5)
tb = add_text_box(slide, Inches(0.8), Inches(3.95), Inches(5.3), Inches(3.1))
tf = set_text(tb, "weave에서의 rplaca + delete", size=18, color=ACCENT5, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "Mark: 기존 환경 제거 표시", size=16, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "(:S12 (rplaca nnew-envs nil))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → new-envs 리스트 내의 항목을 nil로 마킹", size=13, color=WHITE)
add_paragraph(tf, "  → 더 작은 환경이 발견되면 기존 큰 것 표시", size=13, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "Sweep: nil 항목 일괄 제거", size=16, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "(setq envs", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (delete nil new-envs :TEST #'eq))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 각 전건 처리 후 nil을 한꺼번에 제거", size=13, color=WHITE)
add_paragraph(tf, "  → 깨끗한 envs를 다음 전건 처리에 전달", size=13, color=WHITE)

# 하단 오른쪽: Sentinel Skip
card3 = add_shape(slide, Inches(6.7), Inches(3.8), Inches(6.1), Inches(3.4), border_color=ACCENT3)
tb = add_text_box(slide, Inches(7.0), Inches(3.95), Inches(5.5), Inches(3.1))
tf = set_text(tb, "Sentinel Skip 패턴", size=18, color=ACCENT3, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "nil로 마킹된 항목을 건너뛰는 가드 코드:", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(dolist (env envs)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, '  (if env          ; nil이면 건너뜀', size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (dolist ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(when (car nnew-envs)  ; nil이면 건너뜀", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (case ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "delete로 정리하기 전에도 nil 항목이", size=15, color=WHITE)
add_paragraph(tf, "순회에 영향을 주지 않도록 보호합니다.", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "이것은 lazy deletion의 필수 동반 패턴입니다.", size=14, color=GRAY)


# ================================================================
# 슬라이드 10: update-label vs weave 비교
# ================================================================
slide = make_content_slide("비교: update-label vs weave 기법 공유", "10")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.5))
tf = set_text(tb, "두 함수의 기법 비교표", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=12, color=WHITE)

# 표 헤더
add_paragraph(tf, "  프로그래밍 기법                update-label       weave", size=16, color=ACCENT2, bold=True, font_name="Consolas")
add_paragraph(tf, "  ──────────────────────────────────────────────────────────────", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  Defensive Copy (방어적 복사)       X               O (copy-list)", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  rplaca (파괴적 갱신)               O               O", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  Lazy Deletion (지연 삭제)          O               O", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  Subsumption (최소성 검사)          O               O", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  중첩 루프 탐색                     2중 (do-do)     3중 (dolist x3) + do", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  역참조 관리 (env-nodes)            O               X", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  점진적 축적 (Fold)                 X               O (envs 교체)", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  다단계 필터링                      X               O (nogood + subsumption)", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  반환값 흐름 제어                   O               O (nil → 전파 중단)", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  return-from (함수 조기 반환)       X               O", size=15, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  ──────────────────────────────────────────────────────────────", size=14, color=GRAY, font_name="Consolas")

add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "  공유 기법: rplaca, Lazy Deletion, Subsumption → ATMS 코드베이스의 핵심 관용구", size=16, color=ACCENT2)
add_paragraph(tf, "  고유 기법: Defensive Copy, 3중 루프 교차곱, 점진적 축적, 다단계 필터링 → weave의 특수성", size=16, color=ACCENT5)


# ================================================================
# 슬라이드 11: 종합 관계도
# ================================================================
slide = make_content_slide("종합: 기법들의 상호 관계", "11")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.5))
tf = set_text(tb, "weave에 적용된 6가지 기법의 관계도", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=12, color=WHITE)

add_paragraph(tf, "    ┌────────────────────────────────────────────────────────────────┐", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #1 Defensive Copy: 원본 envs를 보호 (copy-list)               │", size=14, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "    └─────────────────────────┬──────────────────────────────────────┘", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "                              ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ┌────────────────────────────────────────────────────────────────┐", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #2 3중 루프: 전건 x envs x label → 모든 조합 생성            │", size=14, color=ACCENT5, font_name="Consolas")
add_paragraph(tf, "    │       │                                                        │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │       ▼                                                        │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #4 다단계 필터: nogood 필터 → subsumption 필터                │", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "    │       │                  │                                     │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │       │            #6 rplaca + Lazy Delete: 마킹 후 일괄 정리  │", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    │       │                  │                                     │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │       ▼                  ▼                                     │", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #3 점진적 축적: envs ← new-envs (각 전건 처리 후 교체)        │", size=14, color=ACCENT1, font_name="Consolas")
add_paragraph(tf, "    └─────────────────────────┬──────────────────────────────────────┘", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "                              ▼", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ┌────────────────────────────────────────────────────────────────┐", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │  #5 Early Return: envs가 비면 즉시 nil 반환 (곱셈의 0 원리)   │", size=14, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "    └────────────────────────────────────────────────────────────────┘", size=14, color=GRAY, font_name="Consolas")


# ================================================================
# 슬라이드 12: 핵심 요약
# ================================================================
slide = make_content_slide("핵심 요약", "12")

summaries = [
    ("#1 Defensive Copy", "copy-list로 호출자의 데이터 보호 → 예기치 않은 부작용 방지", ACCENT2),
    ("#2 3중 루프 교차곱", "dolist x3 구조로 Cartesian Product + union 조합 생성", ACCENT5),
    ("#3 점진적 축적", "Fold/Reduce 패턴: 각 전건 처리 후 envs를 교체하며 누적", ACCENT1),
    ("#4 다단계 필터링", "nogood(O(1)) → subsumption(비교) 순서로 비용 최적화", ACCENT4),
    ("#5 Early Return", "빈 결과 즉시 탈출 (곱셈의 0 원리) + 상위집합 탈락 즉시 탈출", ACCENT4),
    ("#6 Lazy Deletion", "rplaca nil 마킹 + delete nil 정리 → 순회 안전한 삭제 패턴", ACCENT3),
]

y = Inches(1.2)
for title, desc, color in summaries:
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.85), border_color=color)
    tb = add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(3.8), Inches(0.35))
    set_text(tb, title, size=17, color=color, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(4.6), y + Inches(0.12), Inches(8.0), Inches(0.65))
    set_text(tb2, desc, size=16, color=WHITE)
    y += Inches(0.95)

# 하단 메시지
tb = add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8))
tf = set_text(tb, "weave는 update-label의 기법(rplaca, lazy deletion, subsumption)을 재사용하면서,", size=15, color=GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "교차곱(cross-product)이라는 새로운 차원의 알고리즘을 그 위에 구축합니다.", size=15, color=GRAY, alignment=PP_ALIGN.CENTER)


# ================================================================
# 저장
# ================================================================
output_path = "/home/user/bps/atms/weave-techniques.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
