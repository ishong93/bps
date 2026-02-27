#!/usr/bin/env python3
"""weave vs update-label 기능 비교 교육자료 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK    = RGBColor(0x1E, 0x1E, 0x2E)
BG_SURFACE = RGBColor(0x2A, 0x2A, 0x3C)
ACCENT1    = RGBColor(0x89, 0xB4, 0xFA)   # 밝은 파랑
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # 밝은 초록
ACCENT3    = RGBColor(0xF9, 0xE2, 0xAF)   # 밝은 노랑
ACCENT4    = RGBColor(0xF3, 0x8B, 0xA8)   # 밝은 분홍
WHITE      = RGBColor(0xCD, 0xD6, 0xF4)
GRAY       = RGBColor(0x93, 0x99, 0xB2)
CODE_BG    = RGBColor(0x31, 0x32, 0x44)
ACCENT5    = RGBColor(0xCB, 0xA6, 0xF7)   # 밝은 보라
WEAVE_CLR  = RGBColor(0x74, 0xC7, 0xEC)   # 시안 (weave 대표색)
UL_CLR     = RGBColor(0xFA, 0xB3, 0x87)   # 오렌지 (update-label 대표색)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_SURFACE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(shape, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Malgun Gothic", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_code_block(slide, left, top, width, height, lines):
    shape = add_shape(slide, left, top, width, height, fill_color=CODE_BG, border_color=RGBColor(0x58, 0x5B, 0x70))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT3
        p.font.name = "Consolas"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
    return shape


def make_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    set_text(tb, title, size=40, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = add_text_box(slide, Inches(1), Inches(3.7), Inches(11.3), Inches(1.0))
    set_text(tb2, subtitle, size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT1
    line.line.fill.background()
    return slide


def make_section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
    set_text(tb, f"#{number}", size=60, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1.0))
    set_text(tb2, title, size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        tb3 = add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
        set_text(tb3, subtitle, size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def make_content_slide(title, page_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_SURFACE
    bar.line.fill.background()
    tb = add_text_box(slide, Inches(0.6), Inches(0.1), Inches(10), Inches(0.7))
    set_text(tb, title, size=24, color=ACCENT1, bold=True)
    if page_num:
        tb_pg = add_text_box(slide, Inches(11.5), Inches(0.15), Inches(1.5), Inches(0.6))
        set_text(tb_pg, page_num, size=14, color=GRAY, alignment=PP_ALIGN.RIGHT, font_name="Consolas")
    return slide


# ================================================================
# 슬라이드 1: 표지
# ================================================================
make_title_slide(
    "weave vs update-label",
    "ATMS 핵심 함수의 기능적 차이와 공통점 비교 분석"
)

# ================================================================
# 슬라이드 2: 목차
# ================================================================
slide = make_content_slide("목차", "2")
items = [
    ("1", "한눈에 보는 비교                 ", "역할, 입출력, 수학적 의미의 대비"),
    ("2", "호출 흐름에서의 위치             ", "propagate → weave → update → update-label"),
    ("3", "기능적 차이 1: 역할              ", "교차곱 생성 vs label 병합"),
    ("4", "기능적 차이 2: 비교 대상          ", "같은 라운드 결과 vs 기존 label"),
    ("5", "기능적 차이 3: 부수 효과          ", "순수 계산 vs 상태 변형"),
    ("6", "기능적 차이 4: 루프 구조          ", "3중 dolist + do vs 2중 do"),
    ("7", "기능적 차이 5: 조기 종료          ", "return-from weave vs 없음"),
    ("8", "공통점: 공유하는 핵심 관용구      ", "rplaca, lazy deletion, subsumption"),
    ("9", "종합 비교표                       ", "전체 기법 비교 매트릭스"),
]
y = Inches(1.15)
for num, title, desc in items:
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.63))
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.03), Inches(0.5), Inches(0.3))
    set_text(tb, num, size=18, color=ACCENT2, bold=True, font_name="Consolas")
    tb2 = add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(5.5), Inches(0.35))
    set_text(tb2, title, size=16, color=WHITE, bold=True)
    tb3 = add_text_box(slide, Inches(1.5), y + Inches(0.33), Inches(10), Inches(0.3))
    set_text(tb3, desc, size=13, color=GRAY)
    y += Inches(0.68)


# ================================================================
# 슬라이드 3: 한눈에 보는 비교
# ================================================================
slide = make_content_slide("한눈에 보는 비교", "3")

# weave 카드
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=WEAVE_CLR)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(5.5))
tf = set_text(tb, "weave", size=28, color=WEAVE_CLR, bold=True, font_name="Consolas")
add_paragraph(tf, "환경 교차곱 (Cross-Product)", size=16, color=GRAY)
add_paragraph(tf, "", size=12, color=WHITE)
add_paragraph(tf, "한 줄 요약", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "여러 전건 노드의 label을 교차곱하여", size=16, color=WHITE)
add_paragraph(tf, '"모든 전건이 동시에 참인 환경 조합"을 계산', size=16, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "비유", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "여러 서랍의 물건을 조합하기", size=16, color=ACCENT3)
add_paragraph(tf, "옷장에서 상의 x 하의 x 신발 조합 만들기", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "수학적 의미", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "envs = envs x label(A) x label(B) x ...", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "(합집합 연산이 결합된 데카르트 곱)", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "성격", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "순수 계산 (Pure Computation)", size=16, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "ATMS 상태를 변경하지 않음", size=14, color=GRAY)

# update-label 카드
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=UL_CLR)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "update-label", size=28, color=UL_CLR, bold=True, font_name="Consolas")
add_paragraph(tf, "Label 병합 (Merge)", size=16, color=GRAY)
add_paragraph(tf, "", size=12, color=WHITE)
add_paragraph(tf, "한 줄 요약", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "새 환경들을 노드의 label에 병합하되", size=16, color=WHITE)
add_paragraph(tf, "최소성(minimality)을 유지하며 추가/제거", size=16, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "비유", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "서랍에 새 물건 넣기 (중복 정리)", size=16, color=ACCENT3)
add_paragraph(tf, "더 작은 것이 들어오면 큰 것을 빼기", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "수학적 의미", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "label = minimal(label U new-envs)", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "(합집합 후 상위집합 제거)", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "성격", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "상태 변형 (State Mutation)", size=16, color=UL_CLR, bold=True)
add_paragraph(tf, "node.label, env-nodes를 직접 수정", size=14, color=GRAY)


# ================================================================
# 슬라이드 4: 호출 흐름에서의 위치
# ================================================================
make_section_slide("2", "호출 흐름에서의 위치", "propagate → weave → update → update-label")

slide = make_content_slide("호출 흐름: Producer → Consumer 관계", "4")

# 전체 흐름도
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.2))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(3.0))
tf = set_text(tb, "ATMS 전파 파이프라인", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "  propagate(just, antecedent, envs)", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    ├─ weave(antecedent, envs, antecedents)        ← 생산자 (Producer)", size=16, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, '    │    "전건들의 환경을 조합하여 가능한 환경 리스트 생성"', size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    │    반환: 교차곱 결과 환경 리스트", size=14, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, "    │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "    └─ update(new-envs, consequence, just)", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "         │", size=16, color=GRAY, font_name="Consolas")
add_paragraph(tf, "         └─ update-label(consequence, new-envs)    ← 소비자 (Consumer)", size=16, color=UL_CLR, font_name="Consolas")
add_paragraph(tf, '              "받은 환경을 노드의 label에 실제로 반영"', size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "              반환: 실제로 추가된 환경만", size=14, color=UL_CLR, font_name="Consolas")

# 하단: 핵심 포인트
card2 = add_shape(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.3), border_color=WEAVE_CLR)
tb = add_text_box(slide, Inches(0.8), Inches(4.95), Inches(5.3), Inches(2.0))
tf = set_text(tb, "weave가 nil을 반환하면?", size=18, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(if (setq new-envs (weave ...))  ", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (update new-envs ...))        ", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "weave가 nil → update 자체가 호출되지 않음", size=15, color=WHITE)
add_paragraph(tf, "→ update-label도 실행되지 않음 (전파 차단)", size=15, color=ACCENT4)

card3 = add_shape(slide, Inches(6.7), Inches(4.8), Inches(6.1), Inches(2.3), border_color=UL_CLR)
tb = add_text_box(slide, Inches(7.0), Inches(4.95), Inches(5.5), Inches(2.0))
tf = set_text(tb, "update-label이 nil을 반환하면?", size=18, color=UL_CLR, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "(setq new-envs (update-label ...)) ", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "(unless new-envs                    ", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (return-from update nil))         ", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "update-label이 nil → 후속 전파(propagate) 중단", size=15, color=WHITE)
add_paragraph(tf, "→ 이미 알려진 환경이면 재전파 불필요", size=15, color=ACCENT4)


# ================================================================
# 슬라이드 5: 차이 1 - 역할
# ================================================================
make_section_slide("3", "기능적 차이 1: 역할", "교차곱 생성 vs Label 병합")

slide = make_content_slide("차이 1: 역할 — 무엇을 하는가", "5")

# 왼쪽: weave
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5), border_color=WEAVE_CLR)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(2.2))
tf = set_text(tb, "weave: 조합 생성", size=20, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "입력: 여러 전건 노드 + 출발 환경들", size=16, color=WHITE)
add_paragraph(tf, "처리: 각 전건의 label을 교차곱", size=16, color=WHITE)
add_paragraph(tf, "출력: 모든 전건이 동시에 참인 환경 조합들", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "핵심 질문: \"어떤 가정 조합이면 모두 참인가?\"", size=15, color=ACCENT3)

# 오른쪽: update-label
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.5), border_color=UL_CLR)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(2.2))
tf = set_text(tb, "update-label: 병합 갱신", size=20, color=UL_CLR, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "입력: 하나의 노드 + 후보 환경들", size=16, color=WHITE)
add_paragraph(tf, "처리: 기존 label과 비교하며 병합", size=16, color=WHITE)
add_paragraph(tf, "출력: 실제로 새로 추가된 환경들", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "핵심 질문: \"이 환경이 정말 새로운 정보인가?\"", size=15, color=ACCENT3)

# 하단: 구체적 예시
card = add_shape(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.2))
tb = add_text_box(slide, Inches(0.8), Inches(4.25), Inches(11.8), Inches(3.0))
tf = set_text(tb, "구체적 예시: justification A,B → C", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "A.label = [{A1}, {A2}],  B.label = [{A3}]", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "1단계 weave:", size=17, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "   {} x {A1} = {A1},  {} x {A2} = {A2}     (A의 label과 교차곱)", size=15, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, "   {A1} x {A3} = {A1,A3},  {A2} x {A3} = {A2,A3}   (B의 label과 교차곱)", size=15, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, "   결과: [{A1,A3}, {A2,A3}]                 → \"이 조합이면 A,B 모두 참\"", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "2단계 update-label(C, [{A1,A3}, {A2,A3}]):", size=17, color=UL_CLR, bold=True)
add_paragraph(tf, "   기존 C.label = [{A1,A3}] (이미 있다면)", size=15, color=UL_CLR, font_name="Consolas")
add_paragraph(tf, "   {A1,A3} vs {A1,A3} → :EQ → 새것 탈락", size=15, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "   {A2,A3} vs {A1,A3} → 독립 → 추가!", size=15, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "   결과: [{A2,A3}]                           → \"이것만 진짜 새 정보\"", size=15, color=ACCENT3, font_name="Consolas")


# ================================================================
# 슬라이드 6: 차이 2 - 비교 대상
# ================================================================
make_section_slide("4", "기능적 차이 2: 비교 대상", "같은 라운드 결과끼리 vs 기존 label과")

slide = make_content_slide("차이 2: 비교 대상이 다르다", "6")

# 왼쪽: weave
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=WEAVE_CLR)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(5.5))
tf = set_text(tb, "weave: 새것 vs 같은 라운드 결과", size=18, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "비교 대상", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "방금 생성한 환경(new-env)을", size=15, color=WHITE)
add_paragraph(tf, "이번 라운드에서 축적된 new-envs와 비교", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "(do ((nnew-envs new-envs ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (compare-env new-env", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (car nnew-envs)))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "           new-env        (방금 생성)", size=14, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, "              vs", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "         (car nnew-envs)  (이번 라운드 축적)", size=14, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "이유", size=17, color=ACCENT4, bold=True)
add_paragraph(tf, "교차곱 과정에서 중복/상위집합이 생길 수 있음", size=15, color=WHITE)
add_paragraph(tf, "예: {A1}x{A3}={A1,A3}와 {A1,A2}x{A3}={A1,A2,A3}", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "→ {A1,A2,A3}은 {A1,A3}의 상위집합 → 탈락", size=14, color=ACCENT4, font_name="Consolas")

# 오른쪽: update-label
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=UL_CLR)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "update-label: 새것 vs 기존 label", size=18, color=UL_CLR, bold=True)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "비교 대상", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "후보 환경(car new-envs)을", size=15, color=WHITE)
add_paragraph(tf, "노드의 기존 label(envs)과 비교", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "(do ((nenvs envs ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (compare-env (car new-envs)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    (car nenvs)))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "        (car new-envs)   (후보 환경)", size=14, color=UL_CLR, font_name="Consolas")
add_paragraph(tf, "              vs", size=14, color=WHITE, font_name="Consolas")
add_paragraph(tf, "         (car nenvs)     (기존 label 항목)", size=14, color=UL_CLR, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "이유", size=17, color=ACCENT4, bold=True)
add_paragraph(tf, "이미 label에 있는 환경과 중복/포함 관계를", size=15, color=WHITE)
add_paragraph(tf, "확인해야 최소성을 유지할 수 있음", size=15, color=WHITE)
add_paragraph(tf, "예: label=[{A1}], new={A1,A2}", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "→ {A1,A2} ⊇ {A1} → 새것 불필요 → 탈락", size=14, color=ACCENT4, font_name="Consolas")


# ================================================================
# 슬라이드 7: 차이 3 - 부수 효과
# ================================================================
make_section_slide("5", "기능적 차이 3: 부수 효과", "순수 계산 vs 상태 변형")

slide = make_content_slide("차이 3: 부수 효과 (Side Effects)", "7")

# 왼쪽: weave
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=WEAVE_CLR)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(5.5))
tf = set_text(tb, "weave: 순수 계산 (거의 부수 효과 없음)", size=17, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "ATMS 상태 변경", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "  node.label 수정      X", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "  env-nodes 수정       X", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "입력 보호", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "  (copy-list envs)     O  방어적 복사", size=16, color=ACCENT2, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "내부 변형 (로컬)", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "  rplaca nnew-envs nil   (new-envs 내부만)", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  delete nil new-envs    (new-envs 정리)", size=15, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "  → 모두 함수 내부의 지역 데이터에만 적용", size=14, color=GRAY)
add_paragraph(tf, "  → 외부에서 관찰 가능한 변형 없음", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "함수형 프로그래밍 원칙에 가까움", size=16, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "\"같은 입력 → 같은 출력\"", size=15, color=ACCENT3)

# 오른쪽: update-label
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=UL_CLR)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "update-label: 상태 변형 (부수 효과 있음)", size=17, color=UL_CLR, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "ATMS 상태 변경", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "  node.label 수정      O  setf로 직접 갱신", size=16, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "  env-nodes 수정       O  양방향 참조 관리", size=16, color=ACCENT4, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "입력 보호", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "  copy-list            X  복사 불필요", size=16, color=WHITE, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "외부 변형 (글로벌)", size=17, color=ACCENT2, bold=True)
add_paragraph(tf, "  ;; node → env 방향", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  (setf (tms-node-label node) ...)", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  ;; env → node 방향 (추가)", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  (push node (env-nodes new-env))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  ;; env → node 방향 (제거)", size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "  (delete node (env-nodes ...))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "명령형 프로그래밍 스타일", size=16, color=UL_CLR, bold=True)
add_paragraph(tf, "\"상태를 직접 변경하여 시스템 갱신\"", size=15, color=ACCENT3)


# ================================================================
# 슬라이드 8: 차이 4 - 루프 구조
# ================================================================
make_section_slide("6", "기능적 차이 4: 루프 구조", "3중 dolist + do vs 2중 do")

slide = make_content_slide("차이 4: 루프 구조 비교", "8")

# 왼쪽: weave
add_code_block(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), [
    ";; weave: 3중 dolist + 1중 do = 4단계 중첩",
    "",
    "(dolist (node antecedents)        ; L1: 전건 노드들",
    "  ...",
    "  (dolist (env envs)              ; L2: 현재 축적 환경들",
    "    (if env",
    "      (dolist (node-env           ; L3: 노드의 label",
    "               (tms-node-label node))",
    "        (setq new-env",
    "          (union-env env node-env)); 합집합 생성",
    "        (unless (env-nogood? new-env)",
    "          (do ((nnew-envs ...))   ; L4: 최소성 검사",
    "            ...)))))",
    "",
    ";; L1: 무엇과 교차곱할 것인가",
    ";; L2: 현재까지의 결과에서 하나씩",
    ";; L3: 이 노드의 가능한 환경 하나씩",
    ";; L4: 중복/상위집합 제거 검사",
    "",
    ";; 역할: 조합 생성 + 필터링을",
    ";;       하나의 루프 구조에서 처리",
])

# 오른쪽: update-label
add_code_block(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), [
    ";; update-label: 2중 do",
    "",
    "(do ((new-envs new-envs           ; L1: 후보 환경들",
    "      (cdr new-envs)))",
    "    ((null new-envs))",
    "  (do ((nenvs envs               ; L2: 기존 label",
    "        (cdr nenvs)))",
    "      ((null nenvs)",
    "       (push (car new-envs) envs)); 살아남으면 추가",
    "    (cond",
    "      ((null (car nenvs)))",
    "      ((null (car new-envs)))",
    "      ((case (compare-env ...)",
    "         ((:EQ :S21)",
    "          (rplaca new-envs nil))",
    "         (:S12 ...))))))",
    "",
    ";; L1: 추가할 환경 하나씩",
    ";; L2: 기존 label 전체와 비교",
    "",
    ";; 역할: 각 후보를 기존 label 전체와",
    ";;       1:N 비교하여 추가/탈락 결정",
])


# ================================================================
# 슬라이드 9: 차이 5 - 조기 종료
# ================================================================
make_section_slide("7", "기능적 차이 5: 조기 종료", "weave만의 return-from 패턴")

slide = make_content_slide("차이 5: 조기 종료 메커니즘", "9")

# weave
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.8), border_color=WEAVE_CLR)
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5.3), Inches(5.5))
tf = set_text(tb, "weave: 2가지 조기 종료", size=20, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, "Type A: 함수 전체 탈출", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "(unless envs", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (return-from weave nil))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "어떤 전건의 교차곱 결과가 비면", size=15, color=WHITE)
add_paragraph(tf, "전체 함수를 즉시 종료합니다.", size=15, color=WHITE)
add_paragraph(tf, '이유: {} x Anything = {}', size=14, color=GRAY, font_name="Consolas")
add_paragraph(tf, "(곱셈에서 0이 나오면 나머지 불필요)", size=14, color=GRAY)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, "Type B: 내부 do 루프 탈출", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "((:EQ :S21) (return nil))", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "이미 더 작은 환경이 new-envs에 있으면", size=15, color=WHITE)
add_paragraph(tf, "현재 환경의 subsumption 검사를 즉시 중단.", size=15, color=WHITE)
add_paragraph(tf, "다른 env x node-env 조합은 계속 처리.", size=14, color=GRAY)

# update-label
card = add_shape(slide, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.8), border_color=UL_CLR)
tb = add_text_box(slide, Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.5))
tf = set_text(tb, "update-label: 조기 종료 없음", size=20, color=UL_CLR, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)

add_paragraph(tf, "모든 쌍을 반드시 비교", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "update-label에는 return-from이 없습니다.", size=16, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "각 새 환경을 기존 label의 모든 환경과", size=16, color=WHITE)
add_paragraph(tf, "빠짐없이 비교해야 합니다.", size=16, color=WHITE)
add_paragraph(tf, "", size=12, color=WHITE)

add_paragraph(tf, "왜 조기 종료가 불가능한가?", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "새 환경이 기존 label의 부분집합(:S12)이면", size=15, color=WHITE)
add_paragraph(tf, "기존 환경을 제거해야 합니다.", size=15, color=WHITE)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "하나의 새 환경이 여러 기존 환경의", size=15, color=WHITE)
add_paragraph(tf, "부분집합일 수 있으므로, 모든 기존 환경을", size=15, color=WHITE)
add_paragraph(tf, "확인하지 않으면 상위집합이 남을 수 있음.", size=15, color=WHITE)
add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "예: label = [{A1,A2}, {A1,A3}]", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    new  = {A1}", size=14, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "    → 둘 다 제거해야 함 → 끝까지 순회 필수", size=14, color=ACCENT4, font_name="Consolas")


# ================================================================
# 슬라이드 10: 공통점 - 공유 관용구
# ================================================================
make_section_slide("8", "공통점: 공유하는 핵심 관용구", "rplaca, Lazy Deletion, Subsumption — ATMS의 DNA")

slide = make_content_slide("공통점: 세 가지 핵심 관용구", "10")

# 관용구 1: rplaca
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.8), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.75), Inches(1.35), Inches(3.3), Inches(5.5))
tf = set_text(tb, "rplaca nil 마킹", size=18, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "두 함수 모두 삭제할 항목을", size=14, color=WHITE)
add_paragraph(tf, "nil로 마킹합니다.", size=14, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "weave:", size=14, color=WEAVE_CLR, bold=True)
add_paragraph(tf, "(rplaca nnew-envs nil)", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "→ 더 큰 환경 마킹", size=12, color=GRAY)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "update-label:", size=14, color=UL_CLR, bold=True)
add_paragraph(tf, "(rplaca new-envs nil)", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "→ 새 환경 탈락 마킹", size=12, color=GRAY)
add_paragraph(tf, "(rplaca nenvs nil)", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "→ 기존 환경 제거 마킹", size=12, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "목적", size=15, color=ACCENT2, bold=True)
add_paragraph(tf, "순회 중인 리스트의", size=13, color=WHITE)
add_paragraph(tf, "구조를 깨뜨리지 않고", size=13, color=WHITE)
add_paragraph(tf, "O(1)로 삭제 예약", size=13, color=WHITE)

# 관용구 2: Lazy Deletion
card = add_shape(slide, Inches(4.6), Inches(1.2), Inches(4.0), Inches(5.8), border_color=ACCENT2)
tb = add_text_box(slide, Inches(4.85), Inches(1.35), Inches(3.5), Inches(5.5))
tf = set_text(tb, "Lazy Deletion", size=18, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "Mark → Skip → Sweep", size=14, color=WHITE)
add_paragraph(tf, "3단계 패턴이 동일합니다.", size=14, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "Mark:", size=14, color=ACCENT4, bold=True)
add_paragraph(tf, "rplaca ... nil", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "Skip:", size=14, color=ACCENT1, bold=True)
add_paragraph(tf, "weave:", size=12, color=WEAVE_CLR, font_name="Consolas")
add_paragraph(tf, "  (if env ...)", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (when (car nnew-envs)..)", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "update-label:", size=12, color=UL_CLR, font_name="Consolas")
add_paragraph(tf, "  (null (car nenvs))", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  (null (car new-envs))", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "Sweep:", size=14, color=ACCENT2, bold=True)
add_paragraph(tf, "(delete nil ... :TEST #'eq)", size=12, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "목적", size=15, color=ACCENT2, bold=True)
add_paragraph(tf, "순회 안전성 보장 +", size=13, color=WHITE)
add_paragraph(tf, "일괄 정리 효율성", size=13, color=WHITE)

# 관용구 3: Subsumption
card = add_shape(slide, Inches(8.9), Inches(1.2), Inches(3.9), Inches(5.8), border_color=ACCENT1)
tb = add_text_box(slide, Inches(9.15), Inches(1.35), Inches(3.4), Inches(5.5))
tf = set_text(tb, "Subsumption", size=18, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "compare-env로 부분집합", size=14, color=WHITE)
add_paragraph(tf, "관계를 판별합니다.", size=14, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "동일한 판별 로직:", size=14, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, ":EQ  동일", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 새것 탈락", size=12, color=GRAY)
add_paragraph(tf, ":S21 새것 ⊇ 기존", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 새것 탈락", size=12, color=GRAY)
add_paragraph(tf, ":S12 새것 ⊂ 기존", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 기존 제거", size=12, color=GRAY)
add_paragraph(tf, "nil  독립적", size=13, color=ACCENT3, font_name="Consolas")
add_paragraph(tf, "  → 둘 다 유지", size=12, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "목적", size=15, color=ACCENT2, bold=True)
add_paragraph(tf, "환경 집합의 최소성", size=13, color=WHITE)
add_paragraph(tf, "(minimality) 유지:", size=13, color=WHITE)
add_paragraph(tf, '"적은 가정이면 충분한데', size=12, color=GRAY)
add_paragraph(tf, ' 많은 가정을 유지할', size=12, color=GRAY)
add_paragraph(tf, ' 이유가 없다"', size=12, color=GRAY)


# ================================================================
# 슬라이드 11: 종합 비교표
# ================================================================
slide = make_content_slide("종합 비교표", "11")

card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
tb = add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(5.5))
tf = set_text(tb, "전체 비교 매트릭스", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=10, color=WHITE)

header = "  항목                        weave                     update-label"
add_paragraph(tf, header, size=14, color=ACCENT2, bold=True, font_name="Consolas")
add_paragraph(tf, "  ───────────────────────────────────────────────────────────────────────", size=12, color=GRAY, font_name="Consolas")

rows = [
    ("  역할           ", "교차곱으로 환경 조합 생성   ", "기존 label에 환경 병합      "),
    ("  수학적 의미     ", "envs x label(A) x label(B)  ", "minimal(label U new-envs)   "),
    ("  비교 대상       ", "새것 vs 같은 라운드 축적분  ", "새것 vs 기존 label          "),
    ("  부수 효과       ", "없음 (순수 계산)            ", "node.label + env-nodes 변형 "),
    ("  입력 보호       ", "O (copy-list)               ", "X (불필요)                  "),
    ("  루프 구조       ", "3중 dolist + 1중 do         ", "2중 do                      "),
    ("  return-from     ", "O (빈 결과 → 전체 탈출)     ", "X                           "),
    ("  rplaca 마킹     ", "O (new-envs 내부)           ", "O (new-envs + envs 모두)    "),
    ("  lazy deletion   ", "O (delete nil)              ", "O (delete nil)              "),
    ("  subsumption     ", "O (compare-env)             ", "O (compare-env)             "),
    ("  역참조 관리     ", "X                           ", "O (env-nodes push/delete)   "),
    ("  점진적 축적     ", "O (envs <- new-envs)        ", "X                           "),
    ("  nogood 필터     ", "O (env-nogood?)             ", "X (사전 필터링됨)           "),
]

for col1, col2, col3 in rows:
    add_paragraph(tf, f"{col1}{col2}{col3}", size=13, color=WHITE, font_name="Consolas", space_before=Pt(3))

add_paragraph(tf, "  ───────────────────────────────────────────────────────────────────────", size=12, color=GRAY, font_name="Consolas")


# ================================================================
# 슬라이드 12: 핵심 요약
# ================================================================
slide = make_content_slide("핵심 요약", "12")

# 상단: 한 줄 요약
card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.6))
tb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(1.3))
tf = set_text(tb, "한 문장 요약", size=20, color=ACCENT1, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "weave는 \"어떤 가정 조합이 가능한가\"를 계산하는 생산자(Producer)이고,", size=17, color=WEAVE_CLR)
add_paragraph(tf, "update-label은 \"그 결과를 실제로 반영\"하는 소비자(Consumer)입니다.", size=17, color=UL_CLR)

# 차이점 요약
card2 = add_shape(slide, Inches(0.5), Inches(3.1), Inches(5.8), Inches(3.5), border_color=ACCENT4)
tb = add_text_box(slide, Inches(0.8), Inches(3.25), Inches(5.3), Inches(3.2))
tf = set_text(tb, "핵심 차이점 5가지", size=19, color=ACCENT4, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "1. 역할: 교차곱 생성 vs label 병합", size=15, color=WHITE)
add_paragraph(tf, "2. 비교 대상: 라운드 내부 vs 기존 label", size=15, color=WHITE)
add_paragraph(tf, "3. 부수 효과: 순수 계산 vs 상태 변형", size=15, color=WHITE)
add_paragraph(tf, "4. 루프 구조: 4단 중첩 vs 2단 중첩", size=15, color=WHITE)
add_paragraph(tf, "5. 조기 종료: return-from 있음 vs 없음", size=15, color=WHITE)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "→ 서로 다른 문제를 풀지만,", size=14, color=GRAY)
add_paragraph(tf, "  동일한 코딩 관용구를 공유합니다.", size=14, color=GRAY)

# 공통점 요약
card3 = add_shape(slide, Inches(6.7), Inches(3.1), Inches(6.1), Inches(3.5), border_color=ACCENT2)
tb = add_text_box(slide, Inches(7.0), Inches(3.25), Inches(5.5), Inches(3.2))
tf = set_text(tb, "핵심 공통점 3가지", size=19, color=ACCENT2, bold=True)
add_paragraph(tf, "", size=6, color=WHITE)
add_paragraph(tf, "1. rplaca nil 마킹", size=15, color=WHITE)
add_paragraph(tf, "   순회 중 O(1) 삭제 예약", size=14, color=GRAY)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "2. Lazy Deletion (Mark-Skip-Sweep)", size=15, color=WHITE)
add_paragraph(tf, "   순회 안전한 삭제 + 일괄 정리", size=14, color=GRAY)
add_paragraph(tf, "", size=4, color=WHITE)
add_paragraph(tf, "3. Subsumption (compare-env)", size=15, color=WHITE)
add_paragraph(tf, "   부분집합 관계로 최소성 유지", size=14, color=GRAY)
add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "→ 이 세 관용구는 BPS/ATMS 코드베이스", size=14, color=GRAY)
add_paragraph(tf, "  전체를 관통하는 핵심 DNA입니다.", size=14, color=GRAY)

# 하단 마무리
tb = add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5))
tf = set_text(tb, "Producer-Consumer 패턴 + 공유 관용구 = ATMS 전파 엔진의 설계 원리", size=16, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)


# ================================================================
# 저장
# ================================================================
output_path = "/home/user/bps/atms/weave-vs-update-label.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
